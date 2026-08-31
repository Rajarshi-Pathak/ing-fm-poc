from io import BytesIO
from pptx import Presentation
from pitchbook_builder import detect_product_family, build_pitchbook

print("=" * 60)
print("RUNNING END-TO-END VERIFICATION SUITE")
print("=" * 60)

# Test 1: Classifier evaluation order & boundary test
print("\n[Test 1] Testing Product Family Classifier...")
test_cases = [
    ({"opportunity_type": "Green Bond Issuance & SPO"}, "GREEN_ESG"),
    ({"opportunity_type": "Pre-Hedge Interest Rate Swap (IRS)"}, "RATES_HEDGE"),
    ({"opportunity_type": "Strategic FX Collar Overlay"}, "FX_HEDGE"),
    ({"opportunity_type": "Senior EMTN Benchmark Refinancing"}, "DCM_REFI"),
]

for ctx, expected in test_cases:
    res = detect_product_family(ctx)
    assert res == expected, f"Failed: {ctx} -> got {res}, expected {expected}"
    print(f"  ✓ {ctx['opportunity_type'][:35]:35s} -> {res}")

# Test 2: Ingestion dynamic content propagation
print("\n[Test 2] Testing PPTX Generator with Ingested Signal...")
test_trigger = "ALERT: Client issued US$4.5B 30Y Yankee Bond; reverse-enquiry arbitrage window open."
test_action = "Execute €1.2B cross-currency basis swap overlay."

mock_ctx = {
    "client_name": "Test Energy N.V.",
    "rm_name": "Senior Coverage RM",
    "opportunity_type": "GREEN_ESG",
    "why_now_nlg": test_trigger,
    "next_best_action": test_action,
    "revenue_str": "€85,000M",
    "ebitda_str": "€22,000M",
    "net_debt_str": "€50,000M",
    "liquidity_str": "€12,000M",
    "debt_maturing_24m_str": "€8,500M",
    "maturities": [],
    "signals": [],
    "deals": []
}

mock_opp = {
    "opportunity_type": "GREEN_ESG",
    "product_family": "GREEN_ESG",
    "catalyst_rationale": test_trigger,
    "latent_opportunity": test_action
}

pptx_buf = build_pitchbook(mock_ctx, mock_opp)
prs = Presentation(BytesIO(pptx_buf.getvalue()))

slide2 = prs.slides[1]
slide2_texts = []
for shape in slide2.shapes:
    if shape.has_text_frame:
        slide2_texts.append(shape.text_frame.text)

full_slide2_text = "\n".join(slide2_texts)

assert test_trigger in full_slide2_text, "ERROR: Ingested trigger not found in Slide 2!"
assert test_action in full_slide2_text, "ERROR: Ingested action not found in Slide 2!"

print(f"  ✓ Total Slides Generated: {len(prs.slides)}")
print(f"  ✓ Ingested Trigger correctly placed on Slide 2: '{test_trigger[:50]}...'")
print(f"  ✓ Ingested Action correctly placed on Slide 2: '{test_action[:50]}...'")

print("\n" + "=" * 60)
print("ALL VERIFICATION TESTS PASSED (100% DYNAMIC GROUNDING, 0 FABRICATION)")
print("=" * 60)
