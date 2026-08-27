import io
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("pitchbook_builder")

# ING Corporate Palette
ING_ORANGE = RGBColor(255, 98, 0)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_LIGHT_ORANGE = RGBColor(255, 235, 220)
ING_WHITE = RGBColor(255, 255, 255)
BG_LIGHT = RGBColor(248, 249, 250)
LINE_GRAY = RGBColor(220, 224, 230)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)


def add_header(slide, title, is_white=False):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = ING_WHITE if is_white else ING_DARK_SLATE


def add_logo(slide, is_white=False):
    logo_filename = "assets/ing_logo_white.png" if is_white else "assets/ing_logo_orange.png"
    if os.path.exists(logo_filename):
        try:
            slide.shapes.add_picture(logo_filename, Inches(10.8), Inches(0.4), width=Inches(1.8))
            return
        except Exception as exc:
            logger.warning(f"Could not insert logo: {exc}")
    # Fallback text
    tb = slide.shapes.add_textbox(Inches(10.8), Inches(0.4), Inches(1.8), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING"
    p.alignment = PP_ALIGN.RIGHT
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = ING_WHITE if is_white else ING_ORANGE


def add_footer(slide, is_white=False):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING Financial Markets • Strictly Confidential"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(148, 163, 184) if is_white else TEXT_MUTED


def fetch_pitchbook_bundle(canonical_id, client_id_raw, get_db_connection):
    cid = str(canonical_id or client_id_raw or "CLI102").strip()
    
    ctx = {
        "client_id": cid,
        "client_name": "ASML Holding N.V." if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else 
                       "BASF SE" if ("103" in cid or "BASF" in cid.upper()) else
                       "Enel S.p.A." if ("101" in cid or "ENEL" in cid.upper()) else
                       "Orsted A/S" if ("001" in cid or "ORSTED" in cid.upper()) else "Corporate Client",
        "tier": "Tier 1",
        "hq_country": "Netherlands" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "Europe",
        "revenue_str": "—",
        "ebitda_str": "—",
        "net_debt_str": "—",
        "liquidity_str": "—",
        "debt_maturing_24m_str": "—",
        "rm_name": "Daan Visser" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "Named Sector Originator",
        "maturities": [],
        "coverage_team": [],
        "product_family": "FX_HEDGE" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "DCM_REFI",
        "opportunity_type": "FX Hedging Programme" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "Capital Structuring",
        "trigger_source": "FX exposure shift (hedge gap ~$8bn)" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "Balance sheet review",
        "next_best_action": "Propose staged FX hedging programme to close ~$8bn gap" if ("002" in cid or "102" in cid or "ASML" in cid.upper()) else "Propose financing structure"
    }

    try:
        conn, connector = get_db_connection()
        if conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT client_name, tier, hq_country, revenue_eur_m, rm_name 
                    FROM ca.client_master 
                    WHERE client_id = %s OR client_name ILIKE %s LIMIT 1;
                """, (cid, f"%{cid}%"))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["client_name"] = row[0]
                    if row[1]: ctx["tier"] = row[1]
                    if row[2]: ctx["hq_country"] = row[2]
                    if row[3]: ctx["revenue_str"] = f"€{float(row[3]):,.0f}M"
                    if row[4]: ctx["rm_name"] = row[4]

                cur.execute("""
                    SELECT net_debt_eur_m, liquidity_eur_m, ebitda_eur_m, reported_revenue_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings 
                    WHERE client_id = %s 
                    ORDER BY reporting_period DESC LIMIT 1;
                """, (cid,))
                row = cur.fetchone()
                if row:
                    if row[0] is not None: ctx["net_debt_str"] = f"€{float(row[0]):,.0f}M"
                    if row[1] is not None: ctx["liquidity_str"] = f"€{float(row[1]):,.0f}M"
                    if row[2] is not None: ctx["ebitda_str"] = f"€{float(row[2]):,.0f}M"
                    if row[3] is not None and ctx["revenue_str"] == "—": ctx["revenue_str"] = f"€{float(row[3]):,.0f}M"
                    if row[4] is not None: ctx["debt_maturing_24m_str"] = f"€{float(row[4]):,.0f}M"

                cur.execute("""
                    SELECT opportunity_type, trigger_source, next_best_action, est_revenue_eur_000
                    FROM ca.ca_opportunity_scoring 
                    WHERE client_id = %s ORDER BY rank ASC LIMIT 1;
                """, (cid,))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["opportunity_type"] = row[0]
                    if row[1]: ctx["trigger_source"] = row[1]
                    if row[2]: ctx["next_best_action"] = row[2]
                    
                    comb = f"{ctx['opportunity_type']} {ctx['trigger_source']}".lower()
                    if "fx" in comb or "currency" in comb or "bunker" in comb:
                        ctx["product_family"] = "FX_HEDGE"
                    elif "green" in comb or "sustainable" in comb or "sll" in comb or "hybrid" in comb:
                        ctx["product_family"] = "GREEN_ESG"
                    elif "rate" in comb or "irs" in comb or "swap" in comb or "fixed coverage" in comb:
                        ctx["product_family"] = "RATES_HEDGE"
                    else:
                        ctx["product_family"] = "DCM_REFI"

                cur.execute("""
                    SELECT isin, instrument_type, amount_eur_m, maturity_year, coupon_rate_pct, currency
                    FROM ca.debt_maturity_schedule 
                    WHERE client_id = %s 
                    ORDER BY maturity_year ASC;
                """, (cid,))
                rows = cur.fetchall()
                if rows:
                    ctx["maturities"] = []
                    for r in rows:
                        ctx["maturities"].append({
                            "isin": r[0],
                            "instrument_type": r[1],
                            "amount_eur_m": float(r[2]),
                            "maturity_year": str(r[3]),
                            "coupon_rate_pct": float(r[4]) if r[4] is not None else 0.0,
                            "currency": r[5]
                        })

            try: cur.close()
            except Exception: pass
            try: conn.close()
            except Exception: pass
            try: connector.close()
            except Exception: pass

    except Exception as exc:
        logger.warning(f"Database query warning: {exc}")

    if not ctx.get("coverage_team"):
        if "002" in cid or "102" in cid or "ASML" in cid.upper():
            ctx["coverage_team"] = [
                {"name": "Daan Visser", "title": "Global Relationship Manager, Semiconductors (Amsterdam)"},
                {"name": "Marta Nowak", "title": "Director, FX & Financial Markets Derivatives (London)"},
                {"name": "Alex van der Meer", "title": "Managing Director, Technology Capital Markets (Amsterdam)"}
            ]
        elif "103" in cid or "BASF" in cid.upper():
            ctx["coverage_team"] = [
                {"name": "Lena Hoffmann", "title": "Global Relationship Manager, Chemicals (Frankfurt)"},
                {"name": "Roman Weiss", "title": "Managing Director, Rates & Derivatives Structuring (London)"}
            ]
        else:
            ctx["coverage_team"] = [
                {"name": "Giulia Romano", "title": "Global Relationship Manager, Utilities (Milan)"},
                {"name": "Luca Moretti", "title": "Managing Director, Head of DCM Origination (Milan)"},
                {"name": "Marta Nowak", "title": "Director, Financial Markets Derivatives (London)"}
            ]

    return ctx


def build_pitchbook(ctx, opp, compliance_bullets=None, overrides=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    client_name = ctx.get("client_name", "Corporate Client")
    ov = overrides or {}
    p_fam = ctx.get("product_family", "DCM_REFI")

    bond_notional = ov.get("notional_bond", "EUR 600,000,000")
    swap_notional = ov.get("notional_swap", "EUR 400,000,000")
    tenor_str = ov.get("tenor", "7 Years (T + 7Y)")
    spread_str = ov.get("spread", "Mid-Swap + 82 bps")
    market_date = ov.get("market_date", f"Market Snapshot as of {datetime.now().strftime('%d %B %Y, %H:%M CET')}")
    spread_disclaimer = ov.get("spread_disclaimer", "*Indicative pricing subject to market conditions, bookbuilding depth, and credit approval.*")
    
    swap_5y = ov.get("swap_5y", "2.62%")
    iboxx_bbb = ov.get("iboxx_bbb", "115 bps")
    bund_10y = ov.get("bund_10y", "2.61%")
    itraxx_main = ov.get("itraxx_main", "58 bps")
    
    scen_up = ov.get("rate_scenario_up", "4.55%")
    scen_lock = ov.get("rate_scenario_lock", "3.60% (locked)")
    scen_down = ov.get("rate_scenario_down", "3.15%")

    # ----------------------------------------------------
    # SLIDE 1: Cover Slide
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank)
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ING_DARK_SLATE
    bg.line.fill.background()

    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ING_ORANGE
    accent.line.fill.background()

    add_logo(s1, is_white=True)
    add_footer(s1, is_white=True)

    subtitle_banner = "Strategic FX Exposure & Layered Hedging Programme" if p_fam == "FX_HEDGE" else \
                      "Sustainable Capital Structuring & Green Hybrid EMTN" if p_fam == "GREEN_ESG" else \
                      "Rates Risk Advisory & Pre-Hedge Swap Overlay" if p_fam == "RATES_HEDGE" else "Refinancing & Capital Markets Execution Framework"

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf1 = tb1.text_frame
    p = tf1.paragraphs[0]
    p.text = "ING Wholesale Banking"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ING_ORANGE

    p = tf1.add_paragraph()
    p.text = subtitle_banner
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(8)

    p = tf1.add_paragraph()
    p.text = f"Discussion materials prepared for {client_name} Treasury"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(220, 226, 235)
    p.space_before = Pt(6)

    p = tf1.add_paragraph()
    p.text = f"Financial Markets Origination  ·  {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(24)

    # ----------------------------------------------------
    # SLIDE 2: Agenda
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank)
    add_header(s2, "Agenda")
    add_logo(s2)
    add_footer(s2)

    if p_fam == "FX_HEDGE":
        agenda_items = [
            ("01", "Executive Summary", "Strategic USD revenue exposure & hedge policy gap"),
            ("02", "Currency Exposure Profile", "Revenue mix, USD inflow expansion & ~$8bn gap"),
            ("03", "FX Volatility & Rates Backdrop", "EUR/USD spot, forward points & central bank policies"),
            ("04", "Hedging Rationale & Scenarios", "Layered forward collar vs unhedged revenue sensitivity"),
            ("05", "Structured FX Solutions", "Layered forwards, participating collars & corridors"),
            ("06", "Indicative Term Sheet", "Execution architecture & zero-net premium corridor"),
            ("07", "Why ING Financial Markets", "Global FX market making & derivatives franchise"),
            ("08", "Your Coverage Team & Next Steps", f"Lead contacts ({ctx.get('rm_name')}) & roadmap")
        ]
    else:
        agenda_items = [
            ("01", "Executive Summary", "Why we are here and what this pitchbook covers"),
            ("02", "Your Debt Maturity Profile", "Understanding the refinancing need we are here to manage"),
            ("03", "Market Backdrop", "Current rates, credit spreads, and issuance context"),
            ("04", "Refinancing Rationale", "A structured approach to liability management"),
            ("05", "Financing Solutions Suite", "Bonds, loans, liability management & hybrid capital"),
            ("06", "Illustrative Term Sheet", "Worked refinancing example with indicative pricing"),
            ("07", "Why ING", "Platform, balance sheet strength, and execution capability"),
            ("08", "Your Team & Next Steps", "Coverage contacts and proposed timeline")
        ]

    for idx, (num, title_a, sub_a) in enumerate(agenda_items):
        col = idx // 4
        row = idx % 4
        card_x = Inches(0.8 + (col * 5.9))
        card_y = Inches(1.6 + (row * 1.35))
        card_w = Inches(5.6)
        card_h = Inches(1.2)

        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_num = s2.shapes.add_textbox(card_x + Inches(0.15), card_y + Inches(0.18), Inches(0.9), Inches(0.8))
        tf_num = tb_num.text_frame
        tf_num.word_wrap = False
        p_num = tf_num.paragraphs[0]
        p_num.text = num
        p_num.font.bold = True
        p_num.font.size = Pt(26)
        p_num.font.color.rgb = ING_LIGHT_ORANGE

        tb_txt = s2.shapes.add_textbox(card_x + Inches(1.05), card_y + Inches(0.16), card_w - Inches(1.15), Inches(0.9))
        tf_txt = tb_txt.text_frame
        tf_txt.word_wrap = True

        p_title = tf_txt.paragraphs[0]
        p_title.text = title_a
        p_title.font.bold = True
        p_title.font.size = Pt(14)
        p_title.font.color.rgb = ING_ORANGE
        p_title.space_after = Pt(2)

        p_sub = tf_txt.add_paragraph()
        p_sub.text = sub_a
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: Executive Summary
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank)
    add_logo(s3)
    add_footer(s3)

    hero_panel = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
    hero_panel.fill.solid()
    hero_panel.fill.fore_color.rgb = ING_ORANGE
    hero_panel.line.fill.background()

    tb_lh = s3.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(3.8), Inches(4.5))
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    p = tf_lh.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ING_WHITE

    p = tf_lh.add_paragraph()
    p.text = "Strategic Risk Management with Intent"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(12)

    p = tf_lh.add_paragraph()
    p.text = f"Ground-truth analysis for {client_name} based on live treasury exposures, corporate filings, and global market dynamics."
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 235, 220)
    p.space_before = Pt(10)

    if p_fam == "FX_HEDGE":
        pillars = [
            ("1", "Exposure-led, not assumption-led", "Addressing the ~$8bn USD hedge gap arising from North American revenue expansion to 38% of total."),
            ("2", "Multi-tenor layered structure", "Tailored 12M–24M rolling FX forwards and participating collars to protect margins without locking in punitive downside."),
            ("3", "Competitive electronic execution", "Direct pricing via ING Financial Markets global FX liquidity pool across Amsterdam, London, New York and Singapore."),
            ("4", "Dedicated sector partner", f"Named coverage led by {ctx.get('rm_name')} with deep semiconductor capital markets and FX expertise.")
        ]
    else:
        mat_wall = ctx.get("debt_maturing_24m_str", "€4,800M")
        pillars = [
            ("1", "Maturity-led, not mandate-led", f"We start from your actual {mat_wall} maturity wall and cash flow profile to size and time the refinancing that fits."),
            ("2", "Right-sized structure", f"Tailored combination of {bond_notional} benchmark bond and {swap_notional} forward pre-hedge swap."),
            ("3", "Competitive, transparent execution", "Indicative levels shown are for illustration; live executable pricing is provided at bookbuild through our DCM syndicate desk."),
            ("4", "A long-term partner", "Balance sheet strength, global distribution network, and follow-the-sun syndicate coverage across major investor bases.")
        ]

    for idx, (p_num, p_head, p_body) in enumerate(pillars):
        y_pos = Inches(1.2 + (idx * 1.35))
        c_shp = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), y_pos + Inches(0.05), Inches(0.35), Inches(0.35))
        c_shp.fill.solid()
        c_shp.fill.fore_color.rgb = ING_ORANGE
        c_shp.line.fill.background()
        tf_c = c_shp.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = p_num
        p_c.font.bold = True
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = ING_WHITE

        tb_p = s3.shapes.add_textbox(Inches(5.7), y_pos, Inches(6.8), Inches(1.2))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        p_h = tf_p.paragraphs[0]
        p_h.text = p_head
        p_h.font.bold = True
        p_h.font.size = Pt(14)
        p_h.font.color.rgb = ING_ORANGE

        p_b = tf_p.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 4: Balance Sheet & Exposure Snapshot
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank)
    add_header(s4, "Capital Structure & Exposure Profile" if p_fam != "FX_HEDGE" else "Currency Breakdown & Treasury Foundation")
    add_logo(s4)
    add_footer(s4)

    metrics = [
        ("Annual Revenue", ctx.get("revenue_str", "€27,600M")),
        ("EBITDA", ctx.get("ebitda_str", "€6,226M")),
        ("Net Debt / (Cash)", ctx.get("net_debt_str", "€3,192M")),
        ("Available Liquidity", ctx.get("liquidity_str", "€1,008M"))
    ]

    for idx, (lbl, val) in enumerate(metrics):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.6))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s4.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.4))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ING_DARK_SLATE
        p.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 5: Specialised Product Slide
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank)
    add_logo(s5)
    add_footer(s5)

    if p_fam == "FX_HEDGE":
        add_header(s5, "FX Exposure Breakdown & Policy Gap Analysis")
        shp = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb = s5.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Commercial Inflow Mismatch & Unhedged Horizon"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ING_ORANGE

        bullets = [
            "• North American revenue expansion increased USD inflows to >$12.0B (38% of total group revenues).",
            "• Operational manufacturing base remains EUR-denominated (€7.5B, 45% of total expenses), creating structural margin volatility.",
            "• Active FX hedge ratio stands at 50%, below corporate treasury policy target of 75%–80%.",
            "• Unhedged gap of ~$6.0B to $8.0B is exposed to EUR/USD appreciation and volatility spikes.",
            "• Recommended Solution: Implement a staged 12M–24M rolling layered forward programme with zero-cost participating collar overlays."
        ]
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(10)
    else:
        add_header(s5, "Debt Maturity Profile & Refinancing Horizon")
        shp = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb = s5.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Maturity Wall Breakdown for {client_name}"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ING_ORANGE

        for m in ctx.get("maturities", [])[:4]:
            p = tf.add_paragraph()
            p.text = f"• {m['maturity_year']} Maturing Tranche: €{m['amount_eur_m']:,.0f}M ({m.get('instrument_type', 'Bond')}) @ {m.get('coupon_rate_pct', 0.0)}%"
            p.font.size = Pt(12)
            p.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 6: Sensitivity & Scenario Matrix
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank)
    add_header(s6, "Hedging Sensitivity & Payoff Scenarios" if p_fam == "FX_HEDGE" else "Refinancing Rationale & Rate Scenarios")
    add_logo(s6)
    add_footer(s6)

    table_shape = s6.shapes.add_table(4, 3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.2))
    tbl = table_shape.table

    headers = ["Market Scenario", "Structured Programme Payoff", "Unhedged Exposure Impact"] if p_fam == "FX_HEDGE" else ["Rate Scenario", "Refinance Today (Locked)", "Wait 6 Months (Unhedged)"]
    for c_idx, h_txt in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ING_DARK_SLATE
        p = cell.text_frame.paragraphs[0]
        p.text = h_txt
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_WHITE

    if p_fam == "FX_HEDGE":
        rows = [
            ("EUR/USD +5% (USD Weakens)", "Guaranteed Floor Protection (1.0850)", "-$450M Direct Revenue Reduction"),
            ("Spot Unchanged (1.0650)", "1.0650 Executable Forward Rate", "Neutral vs Budget Baseline"),
            ("EUR/USD -5% (USD Strengthens)", "Full Upside Participation up to 1.0450", "+$380M Direct FX Gain")
        ]
    else:
        rows = [
            ("Rates +100 bps", scen_lock, scen_up),
            ("Rates Unchanged", scen_lock, scen_lock),
            ("Rates -50 bps", scen_lock, scen_down)
        ]

    for r_idx, (c1, c2, c3) in enumerate(rows, start=1):
        tbl.cell(r_idx, 0).text_frame.paragraphs[0].text = c1
        tbl.cell(r_idx, 1).text_frame.paragraphs[0].text = c2
        tbl.cell(r_idx, 2).text_frame.paragraphs[0].text = c3
        for c in range(3):
            tbl.cell(r_idx, c).text_frame.paragraphs[0].font.size = Pt(11)

    # ----------------------------------------------------
    # SLIDE 7: Market Intelligence Backdrop
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank)
    add_header(s7, "Market Backdrop & Yield Environment")
    add_logo(s7)
    add_footer(s7)

    mkt_cards = [
        ("5Y EUR Swap Rate", swap_5y),
        ("10Y German Bund", bund_10y),
        ("iBoxx EUR Corp BBB", iboxx_bbb),
        ("iTraxx Europe Main", itraxx_main)
    ]
    for idx, (lbl, val) in enumerate(mkt_cards):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.6))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s7.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.4))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ING_ORANGE
        p.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 8: Indicative Term Sheet
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank)
    add_header(s8, "Indicative Transaction Term Sheet")
    add_logo(s8)
    add_footer(s8)

    table_shape = s8.shapes.add_table(6, 2, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
    tbl = table_shape.table

    if p_fam == "FX_HEDGE":
        ts_rows = [
            ("Client / Counterparty", client_name),
            ("Instrument Structure", "Layered FX Forwards & Zero-Cost Participating Collars"),
            ("Target Notional Programme", "$1,500,000,000 to $3,000,000,000 staged"),
            ("Hedge Tenor Horizon", "12 to 24 Months rolling quarterly"),
            ("Pricing Mechanism", "Zero Net Upfront Premium (Collared Corridor)"),
            ("Structuring & Execution Desk", "ING Wholesale Banking Financial Markets")
        ]
    else:
        ts_rows = [
            ("Issuer", client_name),
            ("Instrument Format", "Senior Unsecured Euro Medium Term Note (EMTN)"),
            ("Notional Sizing", bond_notional),
            ("Tenor / Maturity", tenor_str),
            ("Indicative Margin / Spread", spread_str),
            ("Sole Structurer & Bookrunner", "ING Bank N.V.")
        ]

    for r_idx, (k, v) in enumerate(ts_rows):
        c0 = tbl.cell(r_idx, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = BG_LIGHT
        p0 = c0.text_frame.paragraphs[0]
        p0.text = k
        p0.font.bold = True
        p0.font.size = Pt(11)
        p0.font.color.rgb = ING_DARK_SLATE

        c1 = tbl.cell(r_idx, 1)
        p1 = c1.text_frame.paragraphs[0]
        p1.text = v
        p1.font.size = Pt(11)
        p1.font.color.rgb = TEXT_DARK

    # ----------------------------------------------------
    # SLIDE 9: Coverage Team & Execution Roadmap
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank)
    add_header(s9, "Your ING Deal Team & Execution Roadmap")
    add_logo(s9)
    add_footer(s9)

    tb_t = s9.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p = tf_t.paragraphs[0]
    p.text = "Dedicated Coverage Team"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_ORANGE

    for member in ctx.get("coverage_team", []):
        p = tf_t.add_paragraph()
        p.text = f"• {member.get('name')}"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = ING_DARK_SLATE
        p.space_before = Pt(6)

        p_sub = tf_t.add_paragraph()
        p_sub.text = f"  {member.get('title')}"
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = TEXT_MUTED

    tb_r = s9.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Proposed Execution Timeline"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_ORANGE

    steps = [
        ("Phase 1: Treasury Deep Dive", "Confirm exact currency cashflow schedules and hedge policy boundaries."),
        ("Phase 2: Term Sheet Sign-off", "Calibrate collar strikes, forward tenors and ISDA schedules."),
        ("Phase 3: Staged Execution", "Implement tranche-by-tranche hedging orders via electronic desk."),
        ("Phase 4: Ongoing Monitoring", "Quarterly hedge rebalancing and mark-to-market accounting support.")
    ]
    for st, sd in steps:
        p = tf_r.add_paragraph()
        p.text = f"• {st}"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = ING_DARK_SLATE
        p.space_before = Pt(6)

        p_sub = tf_r.add_paragraph()
        p_sub.text = f"  {sd}"
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 10: Regulatory Disclaimers
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank)
    add_header(s10, "Important Regulatory Notice & Disclaimers")
    add_logo(s10)
    add_footer(s10)

    shp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_LIGHT
    shp.line.color.rgb = LINE_GRAY

    tb_d = s10.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.5))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True

    p = tf_d.paragraphs[0]
    p.text = "Regulatory Disclosures & Target Market Notice"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ING_DARK_SLATE

    disclaimers = ov.get("disclaimers", [
        "This document is prepared for illustrative and discussion purposes only and does not constitute an offer, solicitation, or recommendation to enter into any transaction.",
        "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II / UK MiFIR is eligible counterparties and professional clients only.",
        "This material has not been prepared in accordance with legal requirements designed to promote the independence of investment research.",
        "All rates, levels, spreads, and indicative terms shown are subject to change without notice and are not tradeable prices."
    ])

    for d in disclaimers:
        p = tf_d.add_paragraph()
        p.text = f"• {d}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(8)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
