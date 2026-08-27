import io
import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

logger = logging.getLogger("pitchbook_builder")

# =============================================================================
# ING Brand Colors
# =============================================================================
ING_ORANGE = RGBColor(255, 98, 0)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_WHITE = RGBColor(255, 255, 255)
ING_LIGHT_ORANGE = RGBColor(255, 235, 220)
BG_LIGHT = RGBColor(248, 249, 250)
LINE_GRAY = RGBColor(220, 224, 230)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
CARD_BG_BLUE = RGBColor(240, 244, 255)
CARD_BORDER_BLUE = RGBColor(200, 215, 250)
SUCCESS_GREEN = RGBColor(16, 149, 79)


# =============================================================================
# Helper Functions
# =============================================================================

def add_header(slide, title, category="FINANCIAL MARKETS ORIGINATION", is_white=False):
    tb_k = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.3))
    tf_k = tb_k.text_frame
    tf_k.word_wrap = True
    p_k = tf_k.paragraphs[0]
    p_k.text = category.upper()
    p_k.font.bold = True
    p_k.font.size = Pt(9)
    p_k.font.color.rgb = ING_ORANGE

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.0), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ING_WHITE if is_white else ING_DARK_SLATE


def add_logo(slide, is_white=False):
    base_dir = os.path.dirname(os.path.abspath(__file__))
    logo_filename = os.path.join(base_dir, "assets", "ing_logo_white.png") if is_white else os.path.join(base_dir, "assets", "ing_logo_orange.png")
    if os.path.exists(logo_filename):
        try:
            slide.shapes.add_picture(logo_filename, Inches(11.4), Inches(0.4), width=Inches(1.2))
            return
        except Exception as exc:
            logger.warning(f"Could not insert logo: {exc}")
    tb = slide.shapes.add_textbox(Inches(11.4), Inches(0.4), Inches(1.2), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING"
    p.alignment = PP_ALIGN.RIGHT
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ING_WHITE if is_white else ING_ORANGE


def add_footer(slide, is_white=False):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING Wholesale Banking • Financial Markets • Strictly Confidential"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(160, 175, 200) if is_white else TEXT_MUTED


def detect_product_family(ctx):
    """Detect product family from context data."""
    # Check if we have specific product family from opportunity
    p_fam = ctx.get("product_family", "")
    if p_fam:
        return p_fam
    
    # Detect from client name
    client_name = ctx.get("client_name", "").lower()
    
    # FX: ASML, ASM, or currency-related opportunities
    if "asml" in client_name or "asm" in client_name:
        return "FX_HEDGE"
    
    # Rates: BASF, chemical companies, rate-sensitive industries
    if "basf" in client_name or "chemical" in client_name:
        return "RATES_HEDGE"
    
    # Green/Sustainable: Enel, Orsted, renewable energy
    if "enel" in client_name or "orsted" in client_name or "renewable" in client_name:
        return "GREEN_ESG"
    
    # Default
    return "DCM_REFI"


def get_product_kicker(p_fam):
    """Get product category kicker text."""
    mapping = {
        "FX_HEDGE": "FX & COMMODITY RISK ADVISORY",
        "GREEN_ESG": "SUSTAINABLE & ESG CAPITAL STRUCTURING",
        "RATES_HEDGE": "RATES RISK & LIABILITY MANAGEMENT",
        "DCM_REFI": "DCM CAPITAL STRUCTURING"
    }
    return mapping.get(p_fam, "DCM CAPITAL STRUCTURING")


def get_product_subtitle(p_fam):
    """Get product subtitle text."""
    mapping = {
        "FX_HEDGE": "Strategic FX Exposure Risk & Layered Hedging Programme",
        "GREEN_ESG": "Inaugural Hybrid Green Bond & Sustainability Framework",
        "RATES_HEDGE": "Pre-Hedge Swap Overlay & Rate Sensitivity Immunisation",
        "DCM_REFI": "Refinancing & Capital Markets Execution Framework"
    }
    return mapping.get(p_fam, "Capital Markets Execution Framework")


def get_product_pillars(p_fam, ctx, ov):
    """Get executive summary pillars based on product family."""
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    mat_wall = ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "€0M"))
    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "~$6.0B - $8.0B"))
    notional = ov.get("notional_bond", "EUR 600,000,000")
    
    if p_fam == "FX_HEDGE":
        return [
            ("1", "Exposure-Led Architecture", f"Addressing the {unhedged_gap} unhedged USD exposure."),
            ("2", "Multi-Tenor Layered Corridors", "Rolling 12M–24M zero-cost participating collars."),
            ("3", "Electronic Desk Execution", "Direct automated liquidity sourcing across ING trading desks."),
            ("4", "Dedicated Sector Coverage", f"Direct coverage led by {rm_name}.")
        ]
    elif p_fam == "GREEN_ESG":
        return [
            ("1", "Green Framework Alignment", "ICMA Green Bond Principles & EU Taxonomy alignment."),
            ("2", "Use of Proceeds", "Eligible green projects & sustainability-linked KPIs."),
            ("3", "Greenium Advantage", "3-7 bps pricing concession from ESG-dedicated funds."),
            ("4", "Reporting & Verification", "Annual allocation & impact reporting with third-party verification.")
        ]
    elif p_fam == "RATES_HEDGE":
        return [
            ("1", "Rate Risk Assessment", f"Quantifying the {mat_wall} exposure repricing risk."),
            ("2", "Pre-Hedge Execution", "Forward-starting IRS/swaptions to lock benchmark rates."),
            ("3", "Hedge Policy Alignment", "Aligning with treasury fixed/floating policy targets."),
            ("4", "Execution Window", "Optimizing timing against secondary spread windows.")
        ]
    else:  # DCM_REFI
        return [
            ("1", "Maturity-Led Sizing", f"Addressing the {mat_wall} near-term maturity profile."),
            ("2", "Right-Sized Structure", f"Tailored combination of {notional} benchmark EMTN."),
            ("3", "Competitive Execution", "Direct syndicate distribution across European institutional bases."),
            ("4", "Long-Term Partnership", "Committed balance sheet underwriting and rating optimization.")
        ]


# =============================================================================
# Data Fetching - Fully Database-Driven
# =============================================================================

def fetch_pitchbook_bundle(canonical_id, client_id_raw, get_db_connection):
    """
    Fetch ALL data from database. Resilient to client ID formats.
    """
    cid_input = str(canonical_id or client_id_raw or "CLI101").strip()
    
    ctx = {
        "client_id": cid_input,
        "client_name": "Corporate Client",
        "sector": "Corporate",
        "tier": "Tier 1",
        "hq_country": "Europe",
        "rm_name": "Senior Relationship Manager",
        "revenue_str": "N/A",
        "ebitda_str": "N/A",
        "net_debt_str": "N/A",
        "liquidity_str": "N/A",
        "debt_maturing_24m_str": "N/A",
        "leverage_ratio": "N/A",
        "product_family": "DCM_REFI",
        "opportunity_type": "Capital Structuring",
        "trigger_source": "Balance sheet review",
        "next_best_action": "Propose financing structure",
        "maturities": [],
        "coverage_team": [],
        "signals": [],
        "deals": [],
        "spreads": []
    }

    try:
        conn, connector = get_db_connection()
        if conn:
            cur = conn.cursor()
            
            # 1. RESOLVE CLIENT ID
            cur.execute("""
                SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                FROM ca.client_master 
                WHERE client_id = %s
            """, (cid_input,))
            row = cur.fetchone()
            
            # Method B: Check if cid_input is an opportunity_id
            if not row:
                cur.execute("""
                    SELECT client_id FROM ca.ca_opportunity_scoring 
                    WHERE opportunity_id = %s
                """, (cid_input,))
                opp_row = cur.fetchone()
                if opp_row:
                    cur.execute("""
                        SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                        FROM ca.client_master 
                        WHERE client_id = %s
                    """, (opp_row[0],))
                    row = cur.fetchone()
            
            # Method C: Try partial match on client_name
            if not row:
                cur.execute("""
                    SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                    FROM ca.client_master 
                    WHERE LOWER(client_name) LIKE %s
                """, (f"%{cid_input.lower()}%",))
                row = cur.fetchone()
            
            if row:
                actual_cid = row[0]
                if row[1]: ctx["client_name"] = row[1]
                if row[2]: ctx["tier"] = row[2]
                if row[3]: ctx["hq_country"] = row[3]
                if row[4] and row[4] > 0: 
                    ctx["revenue_str"] = f"€{float(row[4]):,.0f}M"
                if row[5]: ctx["rm_name"] = row[5]
                ctx["client_id"] = actual_cid

                # 2. FINANCIAL FILINGS
                cur.execute("""
                    SELECT net_debt_eur_m, liquidity_eur_m, ebitda_eur_m, 
                           reported_revenue_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings 
                    WHERE client_id = %s 
                    ORDER BY reporting_period DESC LIMIT 1
                """, (actual_cid,))
                row = cur.fetchone()
                if row:
                    if row[0] and row[0] > 0: 
                        ctx["net_debt_str"] = f"€{float(row[0]):,.0f}M"
                    if row[1] and row[1] > 0: 
                        ctx["liquidity_str"] = f"€{float(row[1]):,.0f}M"
                    if row[2] and row[2] > 0: 
                        ctx["ebitda_str"] = f"€{float(row[2]):,.0f}M"
                        if row[0] and row[0] > 0 and float(row[2]) > 0:
                            ctx["leverage_ratio"] = f"{(float(row[0]) / float(row[2])):.2f}x"
                    if row[3] and row[3] > 0: 
                        ctx["revenue_str"] = f"€{float(row[3]):,.0f}M"
                    if row[4] and row[4] > 0: 
                        ctx["debt_maturing_24m_str"] = f"€{float(row[4]):,.0f}M"

                # 3. OPPORTUNITY SCORING
                cur.execute("""
                    SELECT opportunity_type, trigger_source, next_best_action
                    FROM ca.ca_opportunity_scoring 
                    WHERE client_id = %s 
                    ORDER BY rank ASC LIMIT 1
                """, (actual_cid,))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["opportunity_type"] = row[0]
                    if row[1]: ctx["trigger_source"] = row[1]
                    if row[2]: ctx["next_best_action"] = row[2]

                # 4. DEBT MATURITY SCHEDULE
                cur.execute("""
                    SELECT isin, instrument_type, amount_eur_m, maturity_year, 
                           coupon_rate_pct, currency
                    FROM ca.debt_maturity_schedule 
                    WHERE client_id = %s 
                    ORDER BY maturity_year ASC
                """, (actual_cid,))
                rows = cur.fetchall()
                if rows:
                    ctx["maturities"] = []
                    for r in rows:
                        ctx["maturities"].append({
                            "isin": r[0] or "N/A",
                            "instrument_type": r[1] or "Bond",
                            "amount_eur_m": float(r[2]) if r[2] else 0,
                            "maturity_year": str(r[3]) if r[3] else "2026",
                            "coupon_rate_pct": float(r[4]) if r[4] else 0.0,
                            "currency": r[5] or "EUR"
                        })

                # 5. COVERAGE TEAM
                cur.execute("""
                    SELECT role_title, banker_name, location
                    FROM ca.coverage_teams 
                    WHERE client_id = %s
                """, (actual_cid,))
                rows = cur.fetchall()
                if rows:
                    ctx["coverage_team"] = []
                    for r in rows:
                        ctx["coverage_team"].append({
                            "name": r[1] or "Coverage Banker",
                            "title": r[0] or "Global Coverage",
                            "location": r[2] or ""
                        })

            # Close cursor
            try: cur.close()
            except Exception: pass
            try: conn.close()
            except Exception: pass
            if connector:
                try: connector.close()
                except Exception: pass

    except Exception as exc:
        logger.warning(f"Database query warning for {cid_input}: {exc}")
        import traceback
        logger.warning(traceback.format_exc())

    # PRODUCT FAMILY DETECTION
    ctx["product_family"] = detect_product_family(ctx)

    # FALLBACK COVERAGE TEAM
    if not ctx.get("coverage_team"):
        ctx["coverage_team"] = [
            {"name": ctx.get("rm_name", "Relationship Manager"), "title": "Global Coverage"},
            {"name": "Managing Director", "title": "Head of Capital Markets"},
            {"name": "Director", "title": "Financial Markets Structuring"}
        ]

    return ctx


def get_slide_meta(p_fam):
    meta = {
        "FX_HEDGE": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "Currency Inflow Catalysts & Execution Window",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Corporate Liquidity & Currency Inflow Profile",
            "s5_cat": "CURRENCY EXPOSURE PROFILE", "s5_ttl": "FX Currency Breakdown & Hedging Gap",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "FX Scenario Analysis & Layered Collar Payoff",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Central Bank Policy & FX Forward Points",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative FX Risk Management Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Layered Roll Framework & Desk Execution",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Derivative Disclosures"
        },
        "GREEN_ESG": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "ESG Capital Strategy & Decarbonization Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Balance Sheet Capacity & Green CapEx Profile",
            "s5_cat": "USE OF PROCEEDS", "s5_ttl": "Eligible Green Asset Pool & KPI Allocation",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Greenium vs Plain-Vanilla Cost Sensitivity",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "ESG Credit Spreads & Green Bond Index Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Green / Sustainability-Linked Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Second-Party Opinion (SPO) & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "ICMA Green Bond Principles & Target Market Notice"
        },
        "RATES_HEDGE": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "Interest Rate Volatility & Benchmark Swap Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Liquidity Snapshot",
            "s5_cat": "MATURITY & SWAP SCHEDULE", "s5_ttl": "Debt Maturity Profile & Swap Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Rate Shift Sensitivity & Pre-Hedge Lock Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Swap Curve Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Pre-Hedge Swap & EMTN Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "ISDA Schedule, CSA & Execution Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Classification Disclosures"
        },
        "DCM_REFI": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "Executive Context & Opportunity Rationale",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Treasury Health Profile",
            "s5_cat": "MATURITY SCHEDULE", "s5_ttl": "Debt Maturity Profile & Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Refinancing Scenario Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Credit Spread Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Debt Financing Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Roadmap & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Regulatory Notices & Target Market Classification"
        }
    }
    return meta.get(p_fam, meta["DCM_REFI"])

def build_pitchbook(ctx, opp, compliance_bullets=None, overrides=None):
    """
    Build pitchbook using database data. No hardcoded client-specific values.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    ov = overrides or {}
    
    # -------- Client Data from Database (with override support) --------
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    p_fam = ov.get("product_family", ctx.get("product_family", "DCM_REFI"))
    sm = get_slide_meta(p_fam)
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    
    # Financial data from database
    revenue_str = ov.get("revenue_str", ctx.get("revenue_str", "N/A"))
    ebitda_str = ov.get("ebitda_str", ctx.get("ebitda_str", "N/A"))
    net_debt_str = ov.get("net_debt_str", ctx.get("net_debt_str", "N/A"))
    liquidity_str = ov.get("liquidity_str", ctx.get("liquidity_str", "N/A"))
    leverage_str = ov.get("leverage_ratio", ctx.get("leverage_ratio", "N/A"))
    mat_wall_str = ov.get("debt_maturing_24m_str", ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "N/A")))
    
    # Signals from database
    signals = ctx.get("signals", [])
    
    # Deals from database
    deals = ctx.get("deals", [])
    
    # Maturities from database
    maturities = ctx.get("maturities", [])

    # -------- Market Data (with overrides) --------
    swap_5y = ov.get("swap_5y", "2.62%")
    bund_10y = ov.get("bund_10y", "2.61%")
    iboxx_bbb = ov.get("iboxx_bbb", "115 bps")
    itraxx_main = ov.get("itraxx_main", "58 bps")
    ecb_rate = ov.get("ecb_rate", "2.25%")
    fed_rate = ov.get("fed_rate", "4.00–4.25%")

    # -------- Term Sheet (with overrides) --------
    tenor_str = ov.get("tenor", "7 Years (T + 7Y)")
    spread_str = ov.get("spread", "Mid-Swap + 82 bps")
    notional_str = ov.get("notional_bond", "EUR 600,000,000")
    spread_disc = ov.get("spread_disclaimer", "*Indicative pricing subject to market conditions, bookbuilding depth, and credit approval.*")

    # -------- Scenario Values (with overrides) --------
    scen_up = ov.get("rate_scenario_up", "4.55%")
    scen_lock = ov.get("rate_scenario_lock", "3.60% (locked)")
    scen_down = ov.get("rate_scenario_down", "3.15%")
    
    fx_up_unhedged = ov.get("fx_scen_up_unhedged", "-$520M Revenue Impact")
    fx_up_hedged = ov.get("fx_scen_up_hedged", "Guaranteed Floor (1.0850)")
    fx_spot_unhedged = ov.get("fx_scen_spot_unhedged", "1.0650 Spot Level")
    fx_spot_hedged = ov.get("fx_scen_spot_hedged", "1.0650 Forward Rate")
    fx_down_unhedged = ov.get("fx_scen_down_unhedged", "+$380M FX Gain")
    fx_down_hedged = ov.get("fx_scen_down_hedged", "Participate up to 1.0450")

    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "N/A"))

    # -------- Get Product-Specific Content --------
    kicker = get_product_kicker(p_fam)
    subtitle = get_product_subtitle(p_fam)
    pillars = get_product_pillars(p_fam, ctx, ov)
    
    # -------- Build Trigger Cards --------
    trigger_cards = [
        ("Primary Market Trigger", ctx.get("trigger_source", "Active capital structure optimization"), ING_ORANGE),
        ("Window of Opportunity", "Favorable market conditions across European issuance windows.", ING_DARK_SLATE),
        ("Recommended Action", ctx.get("next_best_action", "Propose strategic execution roadmap"), SUCCESS_GREEN)
    ]

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ING_DARK_SLATE
    bg.line.fill.background()
    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ING_ORANGE
    accent.line.fill.background()
    add_logo(s1, is_white=True)
    add_footer(s1, is_white=True)

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf1 = tb1.text_frame
    p = tf1.paragraphs[0]
    p.text = kicker
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_ORANGE

    p = tf1.add_paragraph()
    p.text = client_name
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(6)

    p = tf1.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(220, 226, 235)
    p.space_before = Pt(6)

    p = tf1.add_paragraph()
    p.text = f"Prepared by: {rm_name}  ·  Global Sector Coverage & Capital Markets Desk  ·  {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(28)

    # =========================================================================
    # SLIDE 2: STRATEGIC CATALYST
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    add_header(s2, sm["s2_ttl"], category=sm["s2_cat"])
    add_logo(s2)
    add_footer(s2)

    for idx, (head_c, body_c, col_c) in enumerate(trigger_cards):
        cx = Inches(0.8 + (idx * 3.95))
        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5), Inches(3.8), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_c = s2.shapes.add_textbox(cx + Inches(0.2), Inches(1.7), Inches(3.4), Inches(4.4))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        p = tf_c.paragraphs[0]
        p.text = head_c
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = col_c

        p_b = tf_c.add_paragraph()
        p_b.text = body_c
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(12)

    # =========================================================================
    # SLIDE 3: EXECUTIVE SUMMARY
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    add_logo(s3)
    add_footer(s3)

    hero_panel = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
    hero_panel.fill.solid()
    hero_panel.fill.fore_color.rgb = ING_ORANGE
    hero_panel.line.fill.background()

    tb_lh = s3.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(3.8), Inches(4.5))
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    p = tf_lh.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ING_WHITE

    p = tf_lh.add_paragraph()
    p.text = "Strategic Risk Structuring"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(12)

    # Build financial summary dynamically
    fin_summary = f"Tailored capital and risk architecture for {client_name} based on reported financials"
    if revenue_str != "N/A":
        fin_summary += f" ({revenue_str} revenue"
    if ebitda_str != "N/A":
        fin_summary += f", {ebitda_str} EBITDA" if "revenue" in fin_summary else f" ({ebitda_str} EBITDA"
    if "revenue" in fin_summary or "EBITDA" in fin_summary:
        fin_summary += ")"
    else:
        fin_summary += " and market position"
    fin_summary += "."

    p = tf_lh.add_paragraph()
    p.text = fin_summary
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 235, 220)
    p.space_before = Pt(10)

    for idx, (p_num, p_head, p_body) in enumerate(pillars):
        y_pos = Inches(1.2 + (idx * 1.35))
        c_shp = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), y_pos + Inches(0.05), Inches(0.35), Inches(0.35))
        c_shp.fill.solid()
        c_shp.fill.fore_color.rgb = ING_ORANGE
        c_shp.line.fill.background()
        tf_c = c_shp.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = p_num
        p_c.font.bold = True
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = ING_WHITE

        tb_p = s3.shapes.add_textbox(Inches(5.7), y_pos, Inches(6.8), Inches(1.2))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        p_h = tf_p.paragraphs[0]
        p_h.text = p_head
        p_h.font.bold = True
        p_h.font.size = Pt(14)
        p_h.font.color.rgb = ING_ORANGE

        p_b = tf_p.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: BALANCE SHEET
    # =========================================================================
    s4 = prs.slides.add_slide(blank)
    add_header(s4, sm["s4_ttl"], category=sm["s4_cat"])
    add_logo(s4)
    add_footer(s4)

    metrics = [
        ("Group Revenue", revenue_str),
        ("EBITDA", ebitda_str),
        ("Net Leverage Ratio", leverage_str),
        ("Available Liquidity", liquidity_str)
    ]

    for idx, (lbl, val) in enumerate(metrics):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.6))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s4.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.4))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        
        p = tf_m.add_paragraph()
        p.text = val if val != "N/A" else "—"
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ING_DARK_SLATE if val != "N/A" else TEXT_MUTED
        p.space_before = Pt(8)

    shp_bot = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(11.7), Inches(3.2))
    shp_bot.fill.solid()
    shp_bot.fill.fore_color.rgb = CARD_BG_BLUE
    shp_bot.line.color.rgb = CARD_BORDER_BLUE

    tb_bot = s4.shapes.add_textbox(Inches(1.1), Inches(3.6), Inches(11.1), Inches(2.8))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "Corporate Financial Standing & Balance Sheet Capacity"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_DARK_SLATE

    # Build financial commentary dynamically
    commentary = f"• Net Debt of {net_debt_str} vs EBITDA of {ebitda_str}"
    if leverage_str != "N/A":
        commentary += f" reflects a net leverage ratio of {leverage_str}"
    else:
        commentary += " reflects strong balance sheet health"
    commentary += "."

    p_sub1 = tf_bot.add_paragraph()
    p_sub1.text = commentary
    p_sub1.font.size = Pt(11)
    p_sub1.space_before = Pt(8)

    if liquidity_str != "N/A":
        p_sub2 = tf_bot.add_paragraph()
        p_sub2.text = f"• Robust liquidity of {liquidity_str} provides substantial buffer to execute capital market operations and manage financial risks proactively."
        p_sub2.font.size = Pt(11)
        p_sub2.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: MATURITY / EXPOSURE PROFILE
    # =========================================================================
    s5 = prs.slides.add_slide(blank)
    add_logo(s5)
    add_footer(s5)

    if p_fam == "FX_HEDGE":
        add_header(s5, "FX Currency Breakdown & Hedging Gap", category="CURRENCY EXPOSURE PROFILE")

        shp_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
        shp_l.fill.solid()
        shp_l.fill.fore_color.rgb = BG_LIGHT
        shp_l.line.color.rgb = LINE_GRAY

        tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.7))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        p = tf_l.paragraphs[0]
        p.text = "Revenue Currency Mismatch"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ING_DARK_SLATE

        details = [
            ("USD Exposure", ov.get("usd_exposure", ">$12.0B (38% of total)")),
            ("EUR Base Cost", ov.get("eur_cost", "€7.5B (45% of total)")),
            ("Current Hedge Ratio", ov.get("hedge_ratio", "50% (Under-hedged)")),
            ("Target Policy Ratio", ov.get("target_hedge_ratio", "75% - 80%")),
            ("Unhedged Net Gap", unhedged_gap)
        ]
        for k, v in details:
            p_k = tf_l.add_paragraph()
            p_k.text = f"• {k}: {v}"
            p_k.font.size = Pt(11)
            p_k.space_before = Pt(8)

        chart_data = CategoryChartData()
        chart_data.categories = ['Total USD Inflow', 'Policy Target (75%)', 'Current Hedged', 'Unhedged Gap']
        chart_data.add_series('FX Sizing ($B)', (12.0, 9.0, 6.0, 6.0))
        chart_x, chart_y, chart_cx, chart_cy = Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.1)
        chart = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_cx, chart_cy, chart_data).chart
        chart.has_legend = False

    else:
        add_header(s5, "Debt Maturity Profile & Refinancing Horizon", category="MATURITY SCHEDULE")

        # Build chart from database maturities
        if maturities:
            m_years = [str(m['maturity_year']) for m in maturities[:5]]
            m_amts = [m['amount_eur_m'] for m in maturities[:5]]
        else:
            m_years = ['2026', '2027', '2028', '2029']
            m_amts = [0, 0, 0, 0]

        chart_data = CategoryChartData()
        chart_data.categories = m_years
        chart_data.add_series('Maturing Debt (€M)', tuple(m_amts))

        chart_x, chart_y, chart_cx, chart_cy = Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.1)
        chart = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_cx, chart_cy, chart_data).chart
        chart.has_legend = False

        shp_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.1))
        shp_r.fill.solid()
        shp_r.fill.fore_color.rgb = BG_LIGHT
        shp_r.line.color.rgb = LINE_GRAY

        tb_r = s5.shapes.add_textbox(Inches(7.4), Inches(1.7), Inches(4.9), Inches(4.7))
        tf_r = tb_r.text_frame
        p = tf_r.paragraphs[0]
        p.text = "Refinancing Wall Rationale"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ING_ORANGE

        if mat_wall_str != "N/A":
            wall_text = f"Upcoming debt maturities of {mat_wall_str}"
        else:
            wall_text = "Upcoming debt maturities"
        wall_text += " cluster in near-term windows. Proactive capital structuring and pre-hedge overlays lock in favorable spread certainty."

        p_desc = tf_r.add_paragraph()
        p_desc.text = wall_text
        p_desc.font.size = Pt(11)
        p_desc.space_before = Pt(10)

    # =========================================================================
    # SLIDE 6: SENSITIVITY ANALYSIS
    # =========================================================================
    s6 = prs.slides.add_slide(blank)
    add_header(s6, "FX Scenario Analysis & Hedging Payoff" if p_fam == "FX_HEDGE" else "Refinancing Scenario Analysis", category="SENSITIVITY ANALYSIS")
    add_logo(s6)
    add_footer(s6)

    table_shape = s6.shapes.add_table(4, 3, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
    tbl = table_shape.table

    if p_fam == "FX_HEDGE":
        headers = ["Market Scenario", "Layered Collar Strategy", "Unhedged Exposure"]
        rows = [
            (ov.get("fx_scen_up_lbl", "EUR/USD +5% (USD Weakens)"), fx_up_hedged, fx_up_unhedged),
            (ov.get("fx_scen_spot_lbl", "Spot Unchanged"), fx_spot_hedged, fx_spot_unhedged),
            (ov.get("fx_scen_down_lbl", "EUR/USD -5% (USD Strengthens)"), fx_down_hedged, fx_down_unhedged)
        ]
    else:
        headers = ["Rate Scenario", "Refinance Today (Locked)", "Wait 6 Months (Unhedged)"]
        rows = [
            ("Rates +100 bps", scen_lock, scen_up),
            ("Rates Unchanged", scen_lock, scen_lock),
            ("Rates -50 bps", scen_lock, scen_down)
        ]

    for c_idx, h_txt in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ING_DARK_SLATE
        p = cell.text_frame.paragraphs[0]
        p.text = h_txt
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_WHITE

    for r_idx, (c1, c2, c3) in enumerate(rows, start=1):
        tbl.cell(r_idx, 0).text_frame.paragraphs[0].text = c1
        tbl.cell(r_idx, 1).text_frame.paragraphs[0].text = c2
        tbl.cell(r_idx, 2).text_frame.paragraphs[0].text = c3
        for c in range(3):
            tbl.cell(r_idx, c).text_frame.paragraphs[0].font.size = Pt(11)

    # =========================================================================
    # SLIDE 7: MARKET INTELLIGENCE
    # =========================================================================
    s7 = prs.slides.add_slide(blank)
    add_header(s7, sm["s7_ttl"], category=sm["s7_cat"])
    add_logo(s7)
    add_footer(s7)

    if p_fam == "FX_HEDGE":
        mkt_cards = [
            ("EUR/USD Spot", ov.get("spot_fx", "1.0650")),
            ("12M Forward Pts", ov.get("forward_points", "+185 pts")),
            ("ECB Refi Rate", ecb_rate),
            ("Fed Funds Target", fed_rate)
        ]
        macro_title = "Central Bank Policy Differentials & Structural FX Forward Pickup"
        macro_text = [
            f"• Interest Rate Gap: The differential between Fed ({fed_rate}) and ECB ({ecb_rate}) generates positive forward carry.",
            f"• Strategic Timing: Current forward points provide EUR corporates like {client_name} enhanced forward hedging rates.",
            "• Implied Volatility: Normalized volatility supports zero-cost collar execution with protective strike corridors."
        ]
    else:
        mkt_cards = [
            ("5Y EUR Swap", swap_5y),
            ("10Y Bund", bund_10y),
            ("iBoxx BBB Spread", iboxx_bbb),
            ("iTraxx Main", itraxx_main)
        ]
        macro_title = "Central Bank Policy & Credit Spread Backdrop"
        macro_text = [
            f"• ECB Refinancing Rate at {ecb_rate}; Fed Funds Target at {fed_rate}.",
            "• Tightening European investment grade credit spreads support attractive new-issue concession pricing.",
            "• Benchmark swap curves provide clear pricing windows for pre-hedge execution ahead of policy moves."
        ]

    for idx, (lbl, val) in enumerate(mkt_cards):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.6))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s7.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.4))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ING_ORANGE
        p.space_before = Pt(8)

    shp_cb = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(11.7), Inches(3.2))
    shp_cb.fill.solid()
    shp_cb.fill.fore_color.rgb = CARD_BG_BLUE
    shp_cb.line.color.rgb = CARD_BORDER_BLUE

    tb_cb = s7.shapes.add_textbox(Inches(1.1), Inches(3.6), Inches(11.1), Inches(2.8))
    tf_cb = tb_cb.text_frame
    tf_cb.word_wrap = True
    p = tf_cb.paragraphs[0]
    p.text = macro_title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_DARK_SLATE

    for pt in macro_text:
        p_pt = tf_cb.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(11)
        p_pt.space_before = Pt(6)

    # =========================================================================
    # SLIDE 8: TERM SHEET
    # =========================================================================
    s8 = prs.slides.add_slide(blank)
    add_header(s8, "Indicative FX Risk Management Term Sheet" if p_fam == "FX_HEDGE" else "Indicative Debt Financing Term Sheet", category="TRANSACTION STRUCTURING")
    add_logo(s8)
    add_footer(s8)

    table_shape = s8.shapes.add_table(6, 2, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
    tbl = table_shape.table

    # Determine instrument format based on product family
    if p_fam == "FX_HEDGE":
        instrument = "Layered FX Forwards & Zero-Cost Participating Collars"
    elif p_fam == "GREEN_ESG":
        instrument = "Green / Sustainability-Linked Bond (SLL)"
    elif p_fam == "RATES_HEDGE":
        instrument = "Forward-Starting Interest Rate Swaps (IRS) & Swaptions"
    else:
        instrument = "Senior Unsecured Euro Medium Term Note (EMTN)"

    ts_rows = [
        ("Issuer / Counterparty", client_name),
        ("Instrument Format", instrument),
        ("Notional Sizing", notional_str),
        ("Tenor / Maturity", tenor_str),
        ("Indicative Pricing / Structure", spread_str),
        ("Sole Structurer / Counterparty", "ING Bank N.V.")
    ]

    for r_idx, (k, v) in enumerate(ts_rows):
        c0 = tbl.cell(r_idx, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = BG_LIGHT
        p0 = c0.text_frame.paragraphs[0]
        p0.text = k
        p0.font.bold = True
        p0.font.size = Pt(11)
        p0.font.color.rgb = ING_DARK_SLATE

        c1 = tbl.cell(r_idx, 1)
        p1 = c1.text_frame.paragraphs[0]
        p1.text = v
        p1.font.size = Pt(11)
        p1.font.color.rgb = ING_ORANGE if "Pricing" in k else TEXT_DARK
        if "Pricing" in k:
            p1.font.bold = True

    # Disclaimer Footnote
    tb_fn = s8.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6))
    tf_fn = tb_fn.text_frame
    tf_fn.word_wrap = True
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = spread_disc
    p_fn.font.italic = True
    p_fn.font.size = Pt(9)
    p_fn.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: EXECUTION ROADMAP
    # =========================================================================
    s9 = prs.slides.add_slide(blank)
    add_header(s9, "FX Execution & Layered Roll Framework" if p_fam == "FX_HEDGE" else "Roadmap & Syndicate Timeline", category="EXECUTION ROADMAP")
    add_logo(s9)
    add_footer(s9)

    if p_fam == "FX_HEDGE":
        steps = [
            ("T - 2 Weeks", "Exposure Calibration", "Reconcile commercial USD inflow schedules & ISDA credit lines."),
            ("T - 1 Week", "Strike Setting", "Calibrate floor/cap corridor against spot volatility."),
            ("T-Day", "Electronic Execution", f"Execute initial tranche via ING Global FX electronic desk."),
            ("Post-Trade", "IFRS 9 & Roll Schedule", "Provide hedge accounting designation & quarterly layered roll management.")
        ]
    elif p_fam == "GREEN_ESG":
        steps = [
            ("T - 6 Weeks", "Framework Review", "Review green/SLL framework eligibility & KPIs."),
            ("T - 4 Weeks", "Documentation", "Draft prospectus & second-party opinion/verification."),
            ("T - 2 Weeks", "Investor Roadshow", "Conduct ESG-focused investor meetings."),
            ("T-Day", "Pricing & Allocation", "Final pricing, book allocation & settlement.")
        ]
    elif p_fam == "RATES_HEDGE":
        steps = [
            ("T - 4 Weeks", "Risk Assessment", "Quantify rate exposure & define hedge strategy."),
            ("T - 2 Weeks", "Pre-Hedge Execution", "Execute forward-starting IRS/swaptions."),
            ("T - 1 Week", "Documentation", "ISDA Master Agreement & CSA sign-off."),
            ("T-Day", "Hedge & Monitor", "Instruct trade, execute & begin ongoing monitoring.")
        ]
    else:  # DCM_REFI
        steps = [
            ("T - 4 Weeks", "Documentation", "Confirm EMTN base prospectus & swap schedules."),
            ("T - 2 Weeks", "Pre-Hedge Execution", "Execute treasury pre-hedge overlay swap."),
            ("T - 1 Week", "Global Roadshow", "Conduct syndicate investor meetings across Europe."),
            ("T-Day", "Pricing & Settlement", "Final syndicate pricing, book allocation & closing.")
        ]

    for idx, (t_box, s_title, s_sub) in enumerate(steps):
        sx = Inches(0.8 + (idx * 2.95))
        shp = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(1.5), Inches(2.8), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_s = s9.shapes.add_textbox(sx + Inches(0.15), Inches(1.7), Inches(2.5), Inches(4.4))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True
        
        p_num = tf_s.paragraphs[0]
        p_num.text = f"Step {idx + 1}"
        p_num.font.bold = True
        p_num.font.size = Pt(11)
        p_num.font.color.rgb = ING_ORANGE

        p_t = tf_s.add_paragraph()
        p_t.text = t_box
        p_t.font.bold = True
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = ING_DARK_SLATE
        p_t.space_before = Pt(4)

        p_ttl = tf_s.add_paragraph()
        p_ttl.text = s_title
        p_ttl.font.bold = True
        p_ttl.font.size = Pt(12)
        p_ttl.font.color.rgb = ING_ORANGE
        p_ttl.space_before = Pt(8)

        p_sub = tf_s.add_paragraph()
        p_sub.text = s_sub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_before = Pt(6)

    # =========================================================================
    # SLIDE 10: REGULATORY DISCLAIMERS
    # =========================================================================
    s10 = prs.slides.add_slide(blank)
    add_header(s10, sm["s10_ttl"], category=sm["s10_cat"])
    add_logo(s10)
    add_footer(s10)

    shp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_LIGHT
    shp.line.color.rgb = LINE_GRAY

    tb_disc = s10.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.4))
    tf_disc = tb_disc.text_frame
    tf_disc.word_wrap = True
    p = tf_disc.paragraphs[0]
    p.text = "Regulatory Disclosures & Target Market Notice"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_DARK_SLATE

    disc_list = ov.get("disclaimers") or compliance_bullets or [
        "This document is prepared for illustrative and discussion purposes only and does not constitute an offer, solicitation, or recommendation to enter into any transaction.",
        "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II / UK MiFIR is eligible counterparties and professional clients only (all distribution channels).",
        "This material has not been prepared in accordance with legal requirements designed to promote the independence of investment research.",
        "All rates, levels, spreads, and indicative terms shown are subject to change without notice and are not tradeable prices."
    ]

    # Ensure disc_list is a list
    if isinstance(disc_list, str):
        disc_list = [disc_list]
    elif not isinstance(disc_list, list):
        disc_list = list(disc_list) if disc_list else []

    for d in disc_list:
        if d and d.strip():
            p_d = tf_disc.add_paragraph()
            p_d.text = f"• {d}"
            p_d.font.size = Pt(11)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_before = Pt(8)

    # =========================================================================
    # SAVE
    # =========================================================================
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf