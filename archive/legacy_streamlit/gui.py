import html
import io
import json
import logging
import os
import re
from io import BytesIO

import feedparser
import urllib.parse
import requests
import streamlit as st
from pptx import Presentation
from pypdf import PdfReader

# =============================================================================
# Page configuration
# =============================================================================

st.set_page_config(
    page_title="ING AI Agentic Platform",
    page_icon="🦁",
    layout="wide",
)

# =============================================================================
# Configuration
# =============================================================================

BACKEND_URL = os.environ.get(
    "BACKEND_URL",
    "http://127.0.0.1:8080",
).rstrip("/")

REQUEST_TIMEOUT = int(
    os.environ.get(
        "BACKEND_TIMEOUT_SECONDS",
        "120",
    )
)

HEALTH_TIMEOUT = int(
    os.environ.get(
        "BACKEND_HEALTH_TIMEOUT_SECONDS",
        "5",
    )
)

DOCUMENT_PREVIEW_CHARS = 12000

# =============================================================================
# ING Service Catalog
# =============================================================================

SERVICE_CATALOG = [
    "Foreign Exchange",
    "Interest Rate",
    "Commodities",
    "Credit",
    "Equity Derivatives (GEP)",
    "Global Securities Finance",
    "Structured Financing (SPG)",
    "Money markets",
    "Financing/Capital Markets",
    "Sustainable Finance",
    "Cross-Asset & Discovery",
]

# =============================================================================
# Client Master - FIXED: Use canonical IDs (CLI101, CLI102, etc.)
# =============================================================================

CLIENTS = {
    # Canonical IDs (match seed_db.py)
    "CLI101": "Enel S.p.A. (Energy Utilities)",
    "CLI102": "ASML Holding N.V. (Semiconductors)",
    "CLI103": "BASF SE (Chemicals)",
    "CLI104": "Deutsche Lufthansa AG (Aviation)",
    "CLI105": "Bayer AG (Healthcare & Agri)",
    # Legacy aliases - kept for backward compatibility but map to canonical
    "CLI009_ENEL": "Enel SpA (Legacy Alias → CLI101)",
    "CLI010_BASF": "BASF SE (Legacy Alias → CLI103)",
    "CLI001": "Shell plc / Orsted A/S",
    "CLI002": "ASML Holding N.V.",
    "CLI003": "BMW / Stellantis",
    "CLI008": "A.P. Moller-Maersk A/S",
    "ORG_VODAFONE_UK": "Vodafone Group Plc (Telecom)",
}

# =============================================================================
# Client ID Normalization Helper (matches app.py)
# =============================================================================

def normalize_client_id_for_ui(client_id: str) -> str:
    """Map UI client IDs to canonical IDs expected by backend."""
    if not client_id:
        return client_id
    mapping = {
        "CLI009_ENEL": "CLI101",
        "CLI010_BASF": "CLI103",
        "CLI009": "CLI101",
        "CLI010": "CLI103",
        "ENEL": "CLI101",
        "BASF": "CLI103",
    }
    return mapping.get(str(client_id).strip(), str(client_id).strip())


# =============================================================================
# Demo Source Presets
#
# IMPORTANT:
# These are deliberately conservative.
# They demonstrate the evidence progression rather than manufacture
# a confirmed opportunity.
# =============================================================================

def get_demo_source(client_id: str, source_type: str):
    # Get the display name for the client
    display_name = CLIENTS.get(client_id, "Corporate Client")
    client_name = display_name.split(" (")[0]

    # Use canonical ID for logic
    canonical_id = normalize_client_id_for_ui(client_id)

    # -------------------------------------------------------------------------
    # 1. Document Disclosures / Filings / PDF Extracts
    # -------------------------------------------------------------------------
    if source_type == "document":
        if canonical_id == "CLI101" or client_id == "CLI009_ENEL":
            return (
                "Enel_Capital_Markets_Filing_2026.pdf",
                (
                    "Enel has a significant 2026-2027 debt maturity profile. "
                    "Approximately €10.13bn of debt matures across late 2026 "
                    "and 2027. Legacy fixed coupon debt includes approximately "
                    "1.20% coupons, while indicative current refinancing yields "
                    "are approximately 4.5%-5.0%. The Board has authorized up "
                    "to €12bn of financing capacity through March 2027 to "
                    "support maturities, liquidity and investment requirements. "
                    "These facts indicate a material refinancing and funding "
                    "review requirement, but do not by themselves establish "
                    "an unmet funding gap or confirmed financing mandate."
                ),
            )
        if canonical_id == "CLI103" or client_id == "CLI010_BASF":
            return (
                "BASF_Risk_Disclosure_2026.pdf",
                (
                    "BASF has a material interest-rate protection tranche "
                    "expiring within approximately 12 months. Post-roll "
                    "fixed coverage is expected to decline from 68% to 46% "
                    "against a 60% policy target. TTF gas and power price "
                    "volatility may create additional operating margin exposure. "
                    "The evidence supports further review but does not itself "
                    "constitute a confirmed hedge mandate."
                ),
            )
        # ... rest of document cases ...

        return (
            "Corporate_Filing_Extract.pdf",
            (
                f"{client_name} reports strategic capital allocation updates "
                "and is reviewing forward refinancing schedules and financial "
                "risk exposures."
            ),
        )

    # -------------------------------------------------------------------------
    # 2. News Wires / RSS Feeds
    # -------------------------------------------------------------------------
    if source_type == "news":
        if canonical_id == "CLI101" or client_id == "CLI009_ENEL":
            return (
                "News_RSS_Demo_Feed.txt",
                (
                    "Public market information indicates that Enel completed "
                    "a $2.5bn multi-tranche USD bond issuance in July 2026. "
                    "The transaction demonstrates continued capital-markets "
                    "access. The issuance may already address part of the "
                    "company's planned funding requirements, so the remaining "
                    "maturity profile and funding plan require reconciliation "
                    "before treating additional financing as an opportunity."
                ),
            )
        # ... rest of news cases ...

        return (
            "News_RSS_Demo_Feed.txt",
            (
                "Macro volatility and central-bank rate shifts are prompting "
                "corporate treasuries to review refinancing and financial "
                "risk-management exposures."
            ),
        )

    # -------------------------------------------------------------------------
    # 3. Microsoft Teams Discussions
    # -------------------------------------------------------------------------
    if source_type == "teams":
        if canonical_id == "CLI101" or client_id == "CLI009_ENEL":
            return (
                "MS Teams - European Utilities Coverage (#deal-coverage-enel)",
                (
                    "[09:08] Giulia Romano (RM): Public materials indicate an "
                    "active funding cycle, but the July transaction may have "
                    "addressed part of the requirement. We should reconcile "
                    "the transaction against the maturity profile.\n\n"
                    "[09:12] Luca Moretti (DCM): Agreed. Investment plans are "
                    "relevant context, not evidence of an unfunded amount. "
                    "We still need the residual maturity ladder, proceeds "
                    "allocation and mandate pipeline.\n\n"
                    "[09:15] Marta Nowak (Rates): We should not automatically "
                    "match a pre-hedge. That becomes relevant only if a "
                    "residual execution window exists and the fixed/floating "
                    "mix or policy leaves meaningful rate exposure.\n\n"
                    "[09:18] Elena Ferraro (Sustainable Finance): Green "
                    "financing remains conditional on eligible projects, "
                    "framework capacity and use-of-proceeds confirmation."
                ),
            )
        if canonical_id == "CLI103" or client_id == "CLI010_BASF":
            return (
                "MS Teams - Chemicals Coverage (#deal-coverage-basf)",
                (
                    "[09:02] Lena Hoffmann (Expert): Internal derivative "
                    "schedule shows a material rate-protection roll-off.\n\n"
                    "[09:06] Anna Keller (RM): Rates is the leading "
                    "hypothesis because fixed coverage falls from 68% to 46%.\n\n"
                    "[09:09] Roman Weiss (Rates): Recommend an interest-rate "
                    "hedge-maturity review before the 12-month expiry."
                ),
            )
        # ... rest of teams cases ...

        return (
            "MS Teams - Coverage Working Group",
            (
                "[09:15] Coverage Team: Reviewing upcoming refinancing "
                "ladder and derivative roll-off dates."
            ),
        )

    # -------------------------------------------------------------------------
    # 4. Treasury Inbound Emails
    # -------------------------------------------------------------------------
    if canonical_id == "CLI101" or client_id == "CLI009_ENEL":
        return (
            "Email: Group Treasurer (Enel SpA)",
            (
                "From: Group Treasurer, Enel SpA\n"
                "Subject: Treasury review 2026-2027\n\n"
                "We are reviewing our debt maturity profile and funding "
                "requirements following recent capital-markets activity. "
                "We would welcome an initial discussion covering funding "
                "sequencing, potential rate-risk management and the "
                "eligibility of future projects for Green financing. "
                "There is no immediate financing mandate at this stage."
            ),
        )
    if canonical_id == "CLI103" or client_id == "CLI010_BASF":
        return (
            "Email: BASF Treasury",
            (
                "From: BASF Treasury\n"
                "Subject: Treasury Risk & Derivatives Review\n\n"
                "We are evaluating our interest-rate hedge portfolio "
                "expiring over the next 12 months, along with power-price "
                "volatility. We would like to discuss hedge rollover "
                "alternatives and liquidity considerations."
            ),
        )
    # ... rest of email cases ...

    return (
        "Email: Group Treasury Contact",
        (
            f"Inquiry from {client_name} regarding funding sequencing and risk-management "
            "review for the upcoming planning cycle."
        ),
    )

# =============================================================================
# Backend Helpers
# =============================================================================

def backend_request(method: str, endpoint: str, **kwargs):
    """
    Calls the Flask backend with consistent timeouts, structured error handling,
    and automatic handling for both JSON payloads and binary responses (e.g., PPTX).
    """
    url = f"{BACKEND_URL}{endpoint}"
    kwargs.setdefault("timeout", REQUEST_TIMEOUT)

    try:
        response = requests.request(method, url, **kwargs)
    except requests.RequestException as exc:
        raise RuntimeError(f"Backend service unavailable at {BACKEND_URL}: {exc}") from exc

    # 1. Handle Binary Content Responses (e.g. PPTX / PDF binary downloads)
    content_type = response.headers.get("Content-Type", "").lower()
    is_binary = any(
        t in content_type for t in [
            "presentationml.presentation", 
            "octet-stream", 
            "application/pdf"
        ]
    )

    if is_binary and response.ok:
        return response.content, response

    # 2. Parse Standard JSON Payloads
    try:
        payload = response.json()
    except ValueError:
        payload = {
            "error_message": response.text[:1000] or "Backend returned an invalid non-JSON payload."
        }

    # 3. Raise Actionable Error Messages for HTTP 4xx / 5xx
    if not response.ok:
        if isinstance(payload, dict):
            message = payload.get("error_message") or payload.get("error") or f"HTTP {response.status_code}"
        else:
            message = str(payload)
        raise RuntimeError(f"{message} (HTTP {response.status_code})")

    return payload, response


def escape(value):
    """HTML escapes string values for clean Streamlit rendering."""
    return html.escape(str(value) if value is not None else "")


# =============================================================================
# UI Rendering Helpers
# =============================================================================

def render_signal_card(signal: dict):
    """
    Renders a styled corporate intelligence signal card displaying extracted triggers,
    service catalog mappings, urgency badges, and grounded metric exposures.
    """
    if not isinstance(signal, dict):
        return

    urgency = str(signal.get("urgency", "Medium"))
    evidence_type = str(signal.get("evidence_type", "Derived Signal"))

    urgency_class = {
        "High": "badge-high",
        "Medium": "badge-med",
        "Low": "badge-low",
    }.get(urgency, "badge-med")

    evidence_class = {
        "Fact": "badge-fact",
        "Derived Signal": "badge-derived",
        "Hypothesis": "badge-hypothesis",
        "Client-Validated Discovery": "badge-validated",
    }.get(evidence_type, "badge-derived")

    catalog = escape(signal.get("catalog_family", "Financing/Capital Markets"))
    signal_type = escape(signal.get("signal_type", "Signal"))
    trigger = escape(signal.get("trigger_summary", "Market or balance sheet event identified."))
    metric = escape(signal.get("metric_identified", "N/A"))
    
    # Handle confidence score formatting safely
    raw_conf = signal.get("confidence_pct", "85")
    try:
        confidence = f"{int(float(raw_conf))}%"
    except (ValueError, TypeError):
        confidence = f"{raw_conf}%" if raw_conf != "N/A" else "N/A"

    evidence_basis = escape(signal.get("evidence_basis", ""))

    html_content = (
        '<div class="signal-card">'
        '<div style="display:flex; justify-content:space-between; align-items:center; gap:12px;">'
        f'<span style="font-size:15px; font-weight:bold; color:#000066;">⚡ {signal_type}</span>'
        '<div>'
        f'<span class="badge-cat">{catalog}</span>&nbsp;'
        f'<span class="{urgency_class}">{escape(urgency)}</span>&nbsp;'
        f'<span class="{evidence_class}">{escape(evidence_type)}</span>'
        '</div>'
        '</div>'
        '<div style="margin-top:8px; font-size:13.5px; color:#2D3748; line-height:1.5;">'
        f'<b>Trigger:</b> {trigger}<br/>'
        f'<b>Identified Metric / Exposure:</b> <code>{metric}</code>&nbsp; | &nbsp;'
        f'<b>Confidence:</b> {confidence}'
    )

    if evidence_basis:
        html_content += f'<br/><b>Evidence Basis:</b> <span style="color:#4A5568;">{evidence_basis}</span>'

    html_content += '</div></div>'

    st.markdown(html_content, unsafe_allow_html=True)


def render_opportunity_status(status: str):
    """Renders a standard opportunity lifecycle pill banner."""
    status_str = str(status or "Hypothesis")
    if status_str == "Confirmed Mandate":
        st.success("🟢 **Confirmed Mandate** — Client authorized and active mandate engagement.")
    elif status_str == "Client-Validated Discovery":
        st.info("🔵 **Client-Validated Discovery** — Signal corroborated by client touchpoints or filings.")
    else:
        st.warning("🟠 **Opportunity Hypothesis** — Pre-validation opportunity derived from market data.")

# =============================================================================
# RSS Feed Helper
# =============================================================================

def fetch_rss_articles_as_text(url: str, client_target: str = "Corporate Client", max_entries: int = 8):
    """
    Fetches live RSS XML cleanly for both standard wires and Google News with browser headers,
    redirect handling, HTML cleanup, and robust error recovery.
    """
    if not url or url == "CUSTOM":
        return None, "Custom URL required"

    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
        }

        # Follow redirects and fetch with expanded browser headers
        response = requests.get(url, headers=headers, timeout=12, allow_redirects=True)
        feed = feedparser.parse(response.content)

        # If direct content parsing yielded no entries, attempt direct URL parsing
        if not feed.entries:
            feed = feedparser.parse(url)

        if not feed.entries:
            return None, "No entries found in feed."

        feed_title = feed.feed.get("title", f"Live Intelligence Wire ({client_target})")
        articles = []

        for idx, entry in enumerate(feed.entries[:max_entries]):
            title = entry.get("title", "Market Update")
            raw_summary = entry.get("summary", entry.get("description", ""))

            # Strip HTML tags and normalize entities
            summary = re.sub(r"<[^>]+>", "", raw_summary)
            summary = (
                summary.replace("&nbsp;", " ")
                .replace("&#39;", "'")
                .replace("&quot;", '"')
                .replace("&amp;", "&")
                .strip()
            )

            published = entry.get("published", entry.get("updated", "Recent"))
            link = entry.get("link", "")

            articles.append(
                f"[{idx+1}] HEADLINE: {title}\n"
                f"    PUBLISHED: {published}\n"
                f"    SUMMARY: {summary or 'Market update and capital allocation signal.'}\n"
                f"    URL: {link}\n"
            )

        formatted_text = (
            f"LIVE RSS INTELLIGENCE WIRE: {feed_title}\n"
            f"TARGET ENTITY: {client_target}\n\n" + "\n".join(articles)
        )
        return formatted_text, feed_title

    except Exception as e:
        return None, str(e)

# =============================================================================
# ING Corporate Styling & CSS System (unchanged)
# =============================================================================

st.markdown(
    """
    <style>
    /* Global Typography & Headers */
    .main-header {
        font-size: 28px;
        font-weight: 800;
        color: #000066;
        letter-spacing: -0.5px;
        margin-bottom: 2px;
    }

    .sub-header {
        font-size: 14.5px;
        color: #FF6200;
        font-weight: 600;
        margin-top: 0px;
        margin-bottom: 12px;
    }

    /* Primary ING Action Buttons */
    .stButton > button {
        background-color: #FF6200 !important;
        color: #FFFFFF !important;
        border-radius: 6px !important;
        border: 1px solid #E05500 !important;
        font-weight: 600 !important;
        font-size: 14px !important;
        width: 100% !important;
        height: 44px !important;
        transition: all 0.2s ease-in-out;
    }

    .stButton > button:hover {
        background-color: #E05500 !important;
        border-color: #C04400 !important;
        color: #FFFFFF !important;
        box-shadow: 0 2px 6px rgba(255, 98, 0, 0.25) !important;
    }

    /* Download Buttons */
    .stDownloadButton > button {
        background-color: #000066 !important;
        color: #FFFFFF !important;
        border-radius: 6px !important;
        border: 1px solid #000044 !important;
        font-weight: 600 !important;
        width: 100% !important;
        height: 44px !important;
    }

    .stDownloadButton > button:hover {
        background-color: #000044 !important;
        color: #FFFFFF !important;
        box-shadow: 0 2px 6px rgba(0, 0, 102, 0.25) !important;
    }

    /* Ingestion Channel Container */
    .channel-box {
        background-color: #F8F9FB;
        padding: 12px 16px;
        border-radius: 8px;
        border-left: 5px solid #FF6200;
        margin-bottom: 14px;
        color: #1A202C;
        font-size: 13.5px;
    }

    /* Signal Intelligence Card */
    .signal-card {
        background-color: #FFFFFF;
        border: 1px solid #E2E8F0;
        border-radius: 8px;
        padding: 14px 16px;
        margin-bottom: 12px;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.04);
    }

    /* Signal Badges - Urgency */
    .badge-high {
        background-color: #FDE8E8;
        color: #9B1C1C;
        font-weight: 700;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
        border: 1px solid #F8B4B4;
    }

    .badge-med {
        background-color: #FEF08A;
        color: #713F12;
        font-weight: 700;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
        border: 1px solid #FDE047;
    }

    .badge-low {
        background-color: #F3F4F6;
        color: #374151;
        font-weight: 600;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
    }

    /* Signal Badges - Category & Evidence Classification */
    .badge-cat {
        background-color: #EBF5FF;
        color: #1E429F;
        font-weight: 600;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
        border: 1px solid #BFDBFE;
    }

    .badge-fact {
        background-color: #DCFCE7;
        color: #166534;
        font-weight: 700;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
        border: 1px solid #BBF7D0;
    }

    .badge-derived {
        background-color: #E0F2FE;
        color: #075985;
        font-weight: 600;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
    }

    .badge-hypothesis {
        background-color: #FEF3C7;
        color: #92400E;
        font-weight: 600;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
    }

    .badge-validated {
        background-color: #DBEAFE;
        color: #1E40AF;
        font-weight: 700;
        padding: 3px 8px;
        border-radius: 4px;
        font-size: 11.5px;
        border: 1px solid #BFDBFE;
    }

    /* Pitchbook Highlight Pill */
    .pitchbook-metric {
        background-color: #FF6200;
        color: #FFFFFF;
        padding: 6px 12px;
        border-radius: 4px;
        font-weight: 700;
        font-size: 13.5px;
        display: inline-block;
    }

    /* Highlight Codes inside Cards */
    .signal-card code {
        background-color: #EDF2F7;
        color: #000066;
        font-weight: 600;
        padding: 2px 6px;
        border-radius: 3px;
        font-size: 12.5px;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# =============================================================================
# Header & Sidebar
# =============================================================================

st.markdown(
    '<p class="main-header">🦁 ING Financial Markets — AI Agentic Platform</p>',
    unsafe_allow_html=True,
)
st.markdown(
    '<p class="sub-header">Multi-Channel Ingestion, Context Fabric Signal Intelligence & Pitchbook Engine</p>',
    unsafe_allow_html=True,
)
st.divider()

st.sidebar.title("Institutional Client Coverage")

client_choice = st.sidebar.selectbox(
    "Select Active Corporate Entity:",
    options=list(CLIENTS.keys()),
    format_func=lambda cid: f"{cid} — {CLIENTS[cid]}",
)

client_id = client_choice
client_name = CLIENTS[client_id].split(" (")[0]

# Get canonical ID for backend calls
canonical_client_id = normalize_client_id_for_ui(client_id)

# Session Reset when active client changes
if st.session_state.get("active_client_id") != client_id:
    st.session_state["active_client_id"] = client_id
    st.session_state.pop("active_opportunity", None)
    st.session_state.pop("last_ingestion", None)
    st.session_state.pop("compliance_result", None)
    st.session_state.pop(f"compliance_bullets_{client_id}", None)
    st.session_state.pop(f"pitchbook_bytes_{client_id}", None)

st.sidebar.info(
    f"Target Counterparty: **{client_name}**\n\nClient ID: `{client_id}`\nCanonical ID: `{canonical_client_id}`"
)

if st.sidebar.button("🔌 Check Backend Health"):
    try:
        health, _ = backend_request(
            "GET", "/health", timeout=HEALTH_TIMEOUT
        )
        st.sidebar.success(
            f"Backend: {health.get('status', 'ok')} · {health.get('region', 'n/a')}"
        )
    except Exception as exc:
        st.sidebar.error(str(exc))

st.sidebar.markdown("---")
st.sidebar.markdown("### 🏛️ ING Service Catalog")
st.sidebar.caption("\n".join(f"• {item}" for item in SERVICE_CATALOG))
st.sidebar.caption(f"Backend: `{BACKEND_URL}`")

# =============================================================================
# Operational tabs
# =============================================================================

tab1, tab2, tab3, tab4 = st.tabs(
    [
        "1. Omni-Channel Ingestion & Signal Intelligence",
        "2. Opportunity Discovery (Catalog Mapped)",
        "3. Regulatory Compliance Gateway",
        "4. Pitchbook Presentation Rendering",
    ]
)

# =============================================================================
# TAB 1: Omni-Channel Signal Intelligence
# =============================================================================

with tab1:

    st.subheader("Omni-Channel Signal Ingestion & Intelligence Extraction")
    st.caption("Ingest structured and unstructured intelligence across enterprise touchpoints into Cloud SQL pgvector.")

    channel = st.radio(
        "Select Enterprise Ingestion Channel:",
        options=[
            "📄 Upload Document (PDF / PPTX)",
            "📰 News / RSS Feed",
            "💬 Teams Discussion",
            "✉️ Treasury Email / Context Fabric",
        ],
        horizontal=True,
        key=f"ingestion_channel_{client_id}",
    )

    extracted_text = ""
    source_name = ""
    source_channel = ""

    # -------------------------------------------------------------------------
    # 1. Upload Document
    # -------------------------------------------------------------------------
    if channel.startswith("📄"):

        st.markdown(
            '<div class="channel-box"><b>Upload Credit Reports, Corporate Filings, 10-K/20-F Disclosures or Pitch Decks</b></div>',
            unsafe_allow_html=True,
        )

        uploaded_file = st.file_uploader(
            "Upload PDF or PPTX File",
            type=["pdf", "pptx"],
            key=f"source_document_{client_id}",
        )

        if uploaded_file is not None:
            source_name = uploaded_file.name
            source_channel = "PDF_REPORT" if uploaded_file.name.lower().endswith(".pdf") else "PPTX_REPORT"

            try:
                if uploaded_file.name.lower().endswith(".pdf"):
                    reader = PdfReader(uploaded_file)
                    pages = [page.extract_text() for page in reader.pages if page.extract_text() and page.extract_text().strip()]
                    extracted_text = "\n".join(pages)
                else:
                    prs = Presentation(uploaded_file)
                    text_blocks = [shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
                    extracted_text = "\n".join(text_blocks)

            except Exception as exc:
                st.error(f"Could not extract text from file: {exc}")
                extracted_text = ""

            if extracted_text:
                st.text_area("Extracted Document Text:", value=extracted_text[:DOCUMENT_PREVIEW_CHARS], height=220, disabled=True)
            else:
                st.warning("The file was uploaded, but no extractable text was found.")
        else:
            source_name, extracted_text = get_demo_source(client_id, "document")
            st.info(f"Demo preset loaded: `{source_name}`.")
            extracted_text = st.text_area(
                "Document Content Preview / Edit",
                value=extracted_text,
                height=160,
                key=f"document_demo_text_{client_id}",
            )

    # -------------------------------------------------------------------------
    # 2. Live News & RSS Feeds
    # -------------------------------------------------------------------------
    elif channel.startswith("📰"):

        source_channel = "NEWS_RSS"

        st.markdown(
            '<div class="channel-box"><b>Live Syndicated News Wire & RSS Feed Intelligence</b></div>',
            unsafe_allow_html=True,
        )

        company_keyword = client_name.split()[0] if client_name else "Enel"
        raw_query = f'{company_keyword} (debt OR bonds OR refinancing OR "sustainable finance" OR hedging OR "credit facility")'
        encoded_query = urllib.parse.quote(raw_query)

        dynamic_feed_label = f"Google News Multi-Segment Wire ({client_name})"
        dynamic_google_url = f"https://news.google.com/rss/search?q={encoded_query}&hl=en-US&gl=US&ceid=US:en"

        PRESET_FEEDS = {
            dynamic_feed_label: dynamic_google_url,
            "CNBC Global Financial Markets": "https://www.cnbc.com/id/100003114/device/rss/rss.html",
            "CNBC Central Banks & Economy": "https://www.cnbc.com/id/20910258/device/rss/rss.html",
            "Yahoo Finance Markets": "https://finance.yahoo.com/news/rssindex",
            "BBC Business & Europe": "http://feeds.bbci.co.uk/news/business/rss.xml",
            "Custom RSS Feed URL": "CUSTOM",
        }

        selected_label = st.selectbox(
            "Select Live RSS Feed:",
            options=list(PRESET_FEEDS.keys()),
            index=0,
            key=f"rss_select_{client_id}",
        )

        if selected_label == "Custom RSS Feed URL":
            feed_url = st.text_input(
                "Enter Custom RSS URL:",
                value="https://www.cnbc.com/id/100003114/device/rss/rss.html",
                key=f"custom_rss_input_{client_id}",
            )
        else:
            feed_url = PRESET_FEEDS[selected_label]

        parsed_text, feed_name = fetch_rss_articles_as_text(feed_url, client_target=client_name)

        if parsed_text:
            source_name = f"RSS: {feed_name}"
            extracted_text = st.text_area(
                "Live Syndicated Articles Digest:",
                value=parsed_text,
                height=220,
                key=f"rss_text_digest_{client_id}_{selected_label}",
            )
            st.success(f"Connected to **{feed_name}**")
        else:
            source_name = f"RSS_Feed_{client_name}"
            extracted_text = feed_url
            st.warning(f"Could not parse live feed items directly ({feed_name}). Passing URL to backend.")

    # -------------------------------------------------------------------------
    # 3. Teams Discussion
    # -------------------------------------------------------------------------
    elif channel.startswith("💬"):

        source_channel = "TEAMS"

        st.markdown(
            '<div class="channel-box"><b>Internal Deal Team Communications & Working Group Chat Logs</b></div>',
            unsafe_allow_html=True,
        )

        source_name, extracted_text = get_demo_source(client_id, "teams")

        extracted_text = st.text_area(
            "Teams Discussion",
            value=extracted_text,
            height=220,
            key=f"teams_demo_text_{client_id}",
        )

    # -------------------------------------------------------------------------
    # 4. Treasury Email
    # -------------------------------------------------------------------------
    else:

        source_channel = "TREASURY_EMAIL"

        st.markdown(
            '<div class="channel-box"><b>Inbound Treasury Email / Context Fabric Trigger</b></div>',
            unsafe_allow_html=True,
        )

        source_name, extracted_text = get_demo_source(client_id, "email")

        extracted_text = st.text_area(
            "Treasury Email Content",
            value=extracted_text,
            height=220,
            key=f"email_demo_text_{client_id}",
        )

    # -------------------------------------------------------------------------
    # Ingestion Button & Execution Pipeline
    # -------------------------------------------------------------------------
    if st.button("🚀 Ingest Signal Across Omni-Channel Pipeline", type="primary", key=f"btn_ingest_signal_{client_id}"):

        if not extracted_text.strip():
            st.warning("Nothing to ingest. Provide source text or feed URL.")
        else:
            with st.spinner("Executing live feed retrieval, Vertex AI extraction (Gemini 2.5 Flash), 768-dim embedding, and Cloud SQL pgvector persistence..."):
                try:
                    payload = {
                        # Send canonical ID to backend
                        "client_id": canonical_client_id,
                        "source_channel": source_channel,
                        "source_name": source_name or "Manual_Context_Fabric_Input",
                        "text": extracted_text,
                    }

                    data, _ = backend_request("POST", "/ingest", json=payload)
                    st.session_state["last_ingestion"] = data

                    st.success(f"✅ Ingestion succeeded: `{data.get('source_name')}` processed and grounded in pgvector!")

                    col_v1, col_v2, col_v3 = st.columns(3)
                    col_v1.metric("Storage Vector DB", "Cloud SQL (pgvector)")
                    col_v2.metric("Chunk ID Assigned", str(data.get("chunk_id", "N/A")))
                    col_v3.metric("Embedding Dimension", "768-dim (dense)")

                    meta = data.get("extracted_metadata", {})
                    executive_summary = meta.get("executive_summary", "Signal detected across multi-channel stream.")
                    st.info(f"**Executive Synthesis:** {executive_summary}")

                    assessment = meta.get("overall_evidence_assessment")
                    if assessment:
                        st.caption(f"**Evidence Assessment:** {assessment}")

                    signals = data.get("detected_signals", [])
                    st.markdown("### 🎯 Captured Signals & Intelligence Details")

                    if signals:
                        for signal in signals:
                            render_signal_card(signal)
                    else:
                        st.info("No granular signals were extracted from this source.")

                    with st.expander("🔍 View Raw Extraction JSON"):
                        st.json(data)

                except Exception as exc:
                    st.error(str(exc))

    if st.session_state.get("last_ingestion"):
        with st.expander("Latest Ingestion Result"):
            st.json(st.session_state["last_ingestion"])

# =============================================================================
# TAB 2
# Opportunity Discovery
# =============================================================================

with tab2:

    st.subheader(
        "Opportunity Discovery & Wholesale Service Catalog Matching"
    )

    st.caption(
        "Matches structured Digital Twin signals with retrieved pgvector "
        "evidence and determines whether the result is a hypothesis, "
        "client-validated discovery or confirmed mandate."
    )

    col_t1, col_t2 = st.columns([3, 2])
    with col_t1:
        st.write(
            f"Target Counterparty: **{client_name}** (`{client_id}` → `{canonical_client_id}`)"
        )
    with col_t2:
        min_hurdle = st.slider(
            "🎯 Qualification Confidence Hurdle (%)",
            min_value=50,
            max_value=95,
            value=75,
            step=5,
            key=f"hurdle_slider_{client_id}",
            help="Opportunities with confidence below this hurdle require further discovery before pitchbook generation."
        )

    if st.button(
        "🔎 Discover & Prioritize Opportunities",
        type="primary",
        key=f"discover_opportunity_{client_id}",
    ):

        with st.spinner(
            "Retrieving relevant evidence from DB "
            "and running institutional opportunity reasoning..."
        ):

            try:
                opp, _ = backend_request(
                    "POST",
                    "/match-opportunity",
                    json={
                        "client_id": canonical_client_id
                    },
                )

                st.session_state[
                    "active_opportunity"
                ] = opp

                st.session_state.pop("compliance_result", None)
                st.session_state.pop(f"compliance_bullets_{client_id}", None)
                st.session_state.pop(f"pitchbook_bytes_{client_id}", None)

            except Exception as exc:
                st.error(str(exc))

    opp = st.session_state.get(
        "active_opportunity",
        {},
    )

    if opp:
        confidence = int(opp.get("confidence_pct", opp.get("avg_signal_confidence", 85)))
        score = opp.get("score", 0)
        source_count = opp.get("evidence_source_count", 0)

        st.markdown(
            "### 🏆 Primary Opportunity Assessment"
        )

        render_opportunity_status(
            opp.get(
                "opportunity_status",
                "Hypothesis",
            )
        )

        if confidence >= min_hurdle:
            st.success(
                f"🎯 **Institutional Opportunity Qualified**: "
                f"Confidence (**{confidence}%**) meets your **{min_hurdle}%** hurdle rate."
            )
        else:
            st.warning(
                f"⚠️ **Opportunity Below Qualification Hurdle**: "
                f"Confidence (**{confidence}%**) is below **{min_hurdle}%**. "
                f"Recommend pre-discovery before formal pitching."
            )

        col1, col2, col3, col4, col5 = st.columns(5)

        col1.metric(
            "Catalog Family",
            opp.get(
                "catalog_family",
                "N/A",
            ),
        )

        col2.metric(
            "Target Product",
            opp.get(
                "product",
                "N/A",
            ),
        )

        col3.metric(
            "Priority Score",
            f"{score}/100" if score is not None else "N/A",
        )

        col4.metric(
            "Confidence Score",
            f"{confidence}%",
            delta=f"{confidence - min_hurdle}% vs hurdle"
        )

        col5.metric(
            "Source Documents",
            source_count,
        )

        evidence_record_count = opp.get("evidence_record_count", 0)
        retrieval_mode = opp.get("retrieval_mode", "N/A")
        evidence_sources = opp.get("evidence_sources", [])

        st.caption(
            f"Evidence records retrieved: {evidence_record_count} · "
            f"Source documents: {source_count} · "
            f"Retrieval: `{retrieval_mode}`"
        )

        if evidence_sources:
            with st.expander(f"📚 Evidence Sources ({source_count})"):
                for source in evidence_sources:
                    if isinstance(source, dict):
                        source_channel = escape(source.get("source_channel", "Unknown"))
                        source_name = escape(source.get("source_name", "Unknown source"))
                        st.markdown(f"**{source_name}** `({source_channel})`")
                    else:
                        st.markdown(f"- {escape(source)}")

        rationale = opp.get("rationale", "N/A")
        trigger = opp.get("trigger_source", "N/A")
        urgency = opp.get("urgency", "Medium")
        validation_gap = opp.get("validation_gap", "N/A")

        st.markdown("---")
        st.markdown("#### 📋 Deal Rationale & Next Best Action")
        st.info(rationale)

        st.markdown(f"**Trigger Origin:** `{trigger}`")

        if urgency == "High":
            st.error(f"**Urgency:** {urgency}")
        elif urgency == "Medium":
            st.warning(f"**Urgency:** {urgency}")
        else:
            st.info(f"**Urgency:** {urgency}")

        st.markdown(f"**Validation Gap:** {validation_gap}")

        st.markdown(
            "### 🔄 Cross-Asset & Conditional Opportunity Discovery"
        )

        secs = opp.get("secondary_opportunities", [])

        if secs:
            for secondary in secs:
                if not isinstance(secondary, dict):
                    continue
                family = secondary.get("catalog_family", "N/A")
                product = secondary.get("product", "N/A")
                condition = secondary.get("condition", "N/A")

                with st.container():
                    col_left, col_right = st.columns([3, 1])
                    with col_left:
                        st.markdown(f"**• {family} → {product}**")
                    with col_right:
                        st.caption("Conditional Service")
                    st.markdown(f"**Qualification Requirement:** {condition}")
                    st.divider()
        else:
            st.info("No secondary or conditional opportunities were identified.")

        with st.expander("🔍 Opportunity JSON"):
            st.json(opp)

    else:
        st.info("Run opportunity discovery after ingesting one or more source signals.")

# =============================================================================
# TAB 3
# Regulatory Compliance Gateway - ENHANCED (FULLY FIXED)
# =============================================================================
with tab3:

    st.subheader(
        "FINRA Rule 2210 & MiFID II Pre-Render Inspection"
    )

    st.caption(
        "Automated screening guardrail powered by Gemini Flash. "
        "This is a POC control and not a substitute for formal legal or compliance approval."
    )

    # =====================================================================
    # Load opportunity data from Tab 2
    # =====================================================================

    opp = st.session_state.get("active_opportunity", {})

    if not opp:
        st.warning(
            "⚠️ No active opportunity found. Please run Opportunity Discovery in Tab 2 first."
        )
    else:
        # =====================================================================
        # Display opportunity summary with CORRECT client name
        # =====================================================================

        st.markdown(
            f"""
            <div style="
                background-color: #F8F9FB;
                padding: 12px 16px;
                border-radius: 8px;
                border-left: 4px solid #FF6200;
                margin-bottom: 16px;
            ">
                <b>📋 Opportunity Context</b><br/>
                <span style="font-size:14px;">
                    Client: <b>{client_name}</b> · 
                    Client ID: <b>{client_id}</b> · 
                    Canonical ID: <b>{canonical_client_id}</b> · 
                    Catalog: <b>{opp.get('catalog_family', 'N/A')}</b> · 
                    Product: <b>{opp.get('product', 'N/A')}</b> · 
                    Status: <b>{opp.get('opportunity_status', 'N/A')}</b> · 
                    Score: <b>{opp.get('score', 'N/A')}/100</b>
                </span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # =====================================================================
        # Editable pitchbook bullets (USER CAN EDIT HERE)
        # =====================================================================

        st.markdown("### ✏️ Pitchbook Bullets (Editable)")

        st.caption(
            "Edit the bullets below. These will be used for compliance checking "
            "and as content for the pitchbook generation."
        )

        # Get existing bullets from opportunity or create defaults
        default_bullets = []

        if opp:
            status = opp.get("opportunity_status", "Hypothesis")
            catalog_family = opp.get("catalog_family", "")
            product = opp.get("product", "")
            score = opp.get("score", 0)
            rationale = opp.get("rationale", "")
            validation_gap = opp.get("validation_gap", "")

            default_bullets = [
                f"Opportunity Status: {status}.",
                f"Primary Evidence Area: {catalog_family}.",
                f"Potential Solution: {product}.",
                f"Priority Score: {score}/100.",
            ]

            if rationale:
                default_bullets.append(f"Rationale: {rationale}")

            if validation_gap:
                default_bullets.append(f"Validation Requirement: {validation_gap}")

        # Store edited bullets in client-scoped session state
        bullets_key = f"compliance_bullets_{client_id}"
        if bullets_key not in st.session_state or not st.session_state[bullets_key]:
            st.session_state[bullets_key] = default_bullets

        # Editable text area for bullets
        bullet_text = st.text_area(
            "Edit Pitchbook Bullets",
            value="\n".join(st.session_state[bullets_key]),
            height=180,
            key=f"compliance_editable_bullets_{client_id}",
            help="Edit these bullets. Changes will be used for compliance check and pitchbook generation.",
        )

        # Update session state when user edits
        if bullet_text:
            st.session_state[bullets_key] = [
                line.strip()
                for line in bullet_text.splitlines()
                if line.strip()
            ]

        # =====================================================================
        # Run Compliance Check Button
        # =====================================================================

        if st.button(
            "🛡️ Run Compliance Audit",
            type="primary",
            key=f"btn_run_compliance_{client_id}",
        ):

            bullet_list = st.session_state[bullets_key]

            if not bullet_list:
                st.warning(
                    "Enter at least one narrative bullet before running the audit."
                )
            else:
                with st.spinner(
                    "Inspecting pitch narrative with Vertex AI Compliance Agent..."
                ):
                    try:
                        # Build payload with complete opportunity data
                        payload = {
                            "opportunity": opp,
                            "bullets": bullet_list,
                            "client_name": client_name,
                            "client_id": canonical_client_id,
                        }

                        result, _ = backend_request(
                            "POST",
                            "/check-compliance",
                            json=payload,
                        )

                        st.session_state[f"compliance_result_{client_id}"] = result

                    except Exception as exc:
                        st.error(str(exc))

        # =====================================================================
        # Display Compliance Results
        # =====================================================================

        result = st.session_state.get(f"compliance_result_{client_id}")

        if result:

            # -------- Overall Status --------
            if result.get("compliant") is False:
                st.error("⚠️ Regulatory Compliance Flags Intercepted")
            else:
                st.success("✅ Narrative passed the configured compliance screening checks.")

            # -------- Compliance Summary --------
            compliance_summary = result.get("compliance_summary", "")
            if compliance_summary:
                st.info(f"📋 **Summary:** {compliance_summary}")

            # -------- Risk Assessment --------
            risk_level = result.get("overall_risk_assessment", "LOW")
            risk_colors = {
                "LOW": "🟢",
                "MEDIUM": "🟡",
                "HIGH": "🔴",
            }
            st.metric(
                "Overall Risk Assessment",
                f"{risk_colors.get(risk_level, '')} {risk_level}",
            )

            # -------- Flags --------
            flags = result.get("flags", [])

            if flags:
                st.markdown("### 🚩 Compliance Flags")

                for i, flag in enumerate(flags):
                    flag_type = flag.get("flag_type", "Flag")
                    offending_text = flag.get("offending_text", "")
                    explanation = flag.get("explanation", "Review required.")
                    suggested_fix = flag.get("suggested_fix", "")

                    with st.expander(f"⚠️ {flag_type} - Issue #{i+1}"):
                        st.markdown(f"**Offending Text:** `{offending_text}`")
                        st.markdown(f"**Explanation:** {explanation}")
                        if suggested_fix:
                            st.markdown(f"**Suggested Fix:** {suggested_fix}")

            # -------- Suggested Edits (Auto-generate from flags if empty) --------
            suggested_edits = result.get("suggested_edits", {})
            
            # If suggested_edits is empty, auto-generate from flags
            if not suggested_edits and flags:
                suggested_edits = {}
                for i, flag in enumerate(flags):
                    suggested_edits[f"fix_{i+1}"] = {
                        "original_text": flag.get("offending_text", "N/A"),
                        "suggested_replacement": flag.get("suggested_fix", ""),
                        "explanation": flag.get("explanation", "")
                    }
                result["suggested_edits"] = suggested_edits
                st.session_state[f"compliance_result_{client_id}"] = result

            if suggested_edits:
                st.markdown("### 📝 Suggested Edits")

                # Handle both dict and list structures
                if isinstance(suggested_edits, list):
                    edits_items = [(i+1, e) for i, e in enumerate(suggested_edits) if isinstance(e, dict)]
                elif isinstance(suggested_edits, dict):
                    if "suggested_replacement" in suggested_edits:
                        edits_items = [(1, suggested_edits)]
                    else:
                        edits_items = list(enumerate(suggested_edits.values(), 1))
                else:
                    edits_items = []

                for idx, edit in edits_items:
                    if not isinstance(edit, dict):
                        continue
                    original = edit.get("original_text", "")
                    replacement = edit.get("suggested_replacement", "")
                    explanation = edit.get("explanation", "")

                    if original or replacement:
                        col1, col2 = st.columns(2)
                        with col1:
                            st.markdown(f"**Original:** `{original}`")
                        with col2:
                            st.markdown(f"**Suggested:** `{replacement}`")
                        if explanation:
                            st.caption(f"Reason: {explanation}")
                        st.divider()

            # -------- Required Risk Bullet --------
            required_risk = result.get("required_risk_bullet", "")
            if required_risk:
                st.success(f"**Suggested Risk Warning:** {required_risk}")

            # -------- Mandatory Disclaimer --------
            disclaimer = result.get("mandatory_disclaimer", "")
            if disclaimer:
                st.info(f"**Mandatory Legal Footer:** {disclaimer}")

            # -------- Raw JSON --------
            with st.expander("🔍 Compliance Result JSON"):
                st.json(result)

            # =====================================================================
            # Apply Compliance Recommendations - FULLY FIXED
            # =====================================================================
            
            st.markdown("### 🎯 Apply Compliance Recommendations")

            col1, col2 = st.columns(2)

            with col1:
                if st.button(
                    "📝 Apply Suggested Edits to Pitchbook Bullets",
                    key=f"btn_apply_edits_{client_id}",
                ):
                    current_bullets = st.session_state.get(bullets_key, [])
                    if isinstance(current_bullets, str):
                        updated_bullets = [b.strip() for b in current_bullets.split("\n") if b.strip()]
                    else:
                        updated_bullets = list(current_bullets) if current_bullets else []

                    added_count = 0

                    # 1. Normalize suggested_edits into an iterable list of dicts
                    edits_list = []
                    if isinstance(suggested_edits, list):
                        edits_list = [e for e in suggested_edits if isinstance(e, dict)]
                    elif isinstance(suggested_edits, dict):
                        if "suggested_replacement" in suggested_edits:
                            edits_list = [suggested_edits]
                        else:
                            edits_list = [v for v in suggested_edits.values() if isinstance(v, dict)]

                    for edit in edits_list:
                        original = edit.get("original_text", "").strip()
                        replacement = edit.get("suggested_replacement", "").strip()

                        if not replacement:
                            continue

                        # If original is N/A or flagged as missing, append as new bullet
                        if not original or original.upper().startswith("N/A") or "missing" in original.lower() or "lacks" in original.lower():
                            if not any(replacement.lower() in b.lower() or b.lower() in replacement.lower() for b in updated_bullets):
                                updated_bullets.append(replacement)
                                added_count += 1
                        else:
                            for i, bullet in enumerate(updated_bullets):
                                if original.lower() in bullet.lower():
                                    updated_bullets[i] = bullet.replace(original, replacement)
                                    added_count += 1

                    # 2. Check root-level compliance result fields and flags
                    res_obj = result if isinstance(result, dict) else {}
                    
                    # Mandatory disclaimer check
                    disclaimer_text = res_obj.get("mandatory_disclaimer", "FOR PROFESSIONAL CLIENTS ONLY. Indicative terms subject to market conditions and credit approval under MiFID II.")
                    if not any("FOR PROFESSIONAL CLIENTS ONLY" in b.upper() for b in updated_bullets):
                        updated_bullets.append(disclaimer_text)
                        added_count += 1

                    # Required risk bullet check
                    risk_text = res_obj.get("required_risk_bullet", "All proposed solutions and terms are indicative and subject to market conditions, credit approval, liquidity availability, and execution risk. Past performance is not indicative of future results.")
                    if not any("market, credit, liquidity" in b.lower() or "execution risk" in b.lower() for b in updated_bullets):
                        updated_bullets.append(risk_text)
                        added_count += 1

                    if added_count > 0:
                        if isinstance(st.session_state.get(bullets_key), str):
                            st.session_state[bullets_key] = "\n\n".join(updated_bullets)
                        else:
                            st.session_state[bullets_key] = updated_bullets
                        st.success(f"✅ Applied {added_count} compliance updates to pitchbook bullets!")
                        st.rerun()
                    else:
                        st.info("Bullets already reflect suggested compliance edits.")

            with col2:
                if st.button(
                    "📋 Use These Bullets for Pitchbook",
                    key=f"btn_use_bullets_{client_id}",
                ):
                    current_bullets = st.session_state.get(bullets_key, [])
                    st.session_state[f"pitchbook_bullets_{client_id}"] = current_bullets
                    st.success("✅ Bullets copied to pitchbook!")

# =============================================================================
# TAB 4: Institutional Master Pitchbook Generation & Export
# =============================================================================

with tab4:

    st.subheader("📑 Institutional Pitchbook Assembly & Export")
    st.caption(
        "Generates an ING-branded, compliance-cleared 16:9 widescreen PowerPoint pitchbook "
        "grounded in real-time signals, Cloud SQL master data, and corporate filings."
    )

    opp = st.session_state.get("active_opportunity", {})

    if not opp:
        st.info(
            "⚠️ No active opportunity loaded from Tab 2. You can still generate a standard "
            "institutional pitchbook, or run Opportunity Discovery first in Tab 2 for tailored intelligence."
        )

    # Display opportunity summary banner if available
    if opp:
        st.markdown(
            f"""
            <div style="
                background-color: #F8F9FB;
                padding: 14px 18px;
                border-radius: 8px;
                border-left: 4px solid #FF6200;
                margin-bottom: 16px;
            ">
                <b>📊 Opportunity Assessment Context</b><br/>
                <span style="font-size:14px; color: #1A202C;">
                    Client: <b>{client_name}</b> · 
                    Status: <b>{opp.get('opportunity_status', 'Discovery')}</b> · 
                    Priority Score: <b>{opp.get('score', 85)}/100</b> · 
                    Catalog Family: <b>{opp.get('catalog_family', 'Financing/Capital Markets')}</b> · 
                    Product: <b>{opp.get('product', 'Strategic Refinancing')}</b>
                </span>
            </div>
            """,
            unsafe_allow_html=True,
        )

    # Input parameters
    deck_client = st.text_input(
        "Client Legal Entity",
        value=client_name,
        key=f"deck_client_{client_id}",
    )

    opportunity_status = opp.get(
        "opportunity_status",
        "Client-Validated Discovery",
    ) if opp else "Client-Validated Discovery"
    
    default_product = opp.get(
        "product",
        "Funding Sequencing / Liability Management Advisory",
    ) if opp else "Funding Sequencing / Liability Management Advisory"
    
    default_title = f"{opportunity_status}: {default_product}"

    deck_title = st.text_input(
        "Proposal Headline",
        value=default_title,
        key=f"deck_title_{client_id}",
    )

    # Check for compliance bullets locked from Tab 3
    locked_bullets = st.session_state.get(f"pitchbook_bullets_{client_id}", [])
    if locked_bullets:
        st.success(f"🔒 Using **{len(locked_bullets)}** compliance-approved narrative bullets from Tab 3.")
        with st.expander("👁️ View Active Pitchbook Bullets"):
            for b in locked_bullets:
                st.markdown(f"- {escape(b)}")

    # Trigger Generation Section
    st.markdown("---")
    col_gen, col_status = st.columns([1, 1])

    with col_gen:
        if st.button(
            "📊 Generate Master Pitchbook Deck",
            type="primary",
            key=f"btn_gen_master_pitchbook_{client_id}",
        ):
            with st.spinner("Assembling institutional master pitchbook slides with python-pptx..."):
                try:
                    # Dynamically extract rationale and validation_gap from Tab 3 bullets if edited
                    rationale_to_send = opp.get(
                        "rationale",
                        "Active treasury refinancing sequencing window identified.",
                    ) if opp else "Active treasury refinancing sequencing window identified."
                    
                    validation_gap_to_send = opp.get(
                        "validation_gap",
                        "Confirm residual funding requirement post-recent bond issuance.",
                    ) if opp else "Confirm residual funding requirement post-recent bond issuance."

                    if locked_bullets:
                        for b in locked_bullets:
                            if "Rationale:" in b:
                                rationale_to_send = b.split("Rationale:", 1)[-1].strip()
                            if "Validation Requirement:" in b:
                                validation_gap_to_send = b.split("Validation Requirement:", 1)[-1].strip()

                    payload = {
                        "client_name": deck_client.strip() or client_name,
                        "client_id": canonical_client_id,
                        "title": deck_title.strip() or default_title,
                        "catalog_family": opp.get("catalog_family", "Financing/Capital Markets") if opp else "Financing/Capital Markets",
                        "product": opp.get("product", default_product) if opp else default_product,
                        "score": opp.get("score", 90) if opp else 90,
                        "opportunity_status": opportunity_status,
                        "rationale": rationale_to_send,
                        "validation_gap": validation_gap_to_send,
                        "urgency": opp.get("urgency", "High") if opp else "High",
                        "bullets": locked_bullets,
                        "evidence_source_count": opp.get("evidence_source_count", 0) if opp else 0,
                        "evidence_record_count": opp.get("evidence_record_count", 0) if opp else 0,
                        # Add these for completeness
                        "evidence_sources": opp.get("evidence_sources", []) if opp else [],
                        "secondary_opportunities": opp.get("secondary_opportunities", []) if opp else [],
                    }

                    res = requests.post(
                        f"{BACKEND_URL}/generate-pitchbook",
                        json=payload,
                        timeout=REQUEST_TIMEOUT,
                    )

                    if not res.ok:
                        try:
                            msg = res.json().get("error_message", "Pitchbook generation failed.")
                        except Exception:
                            msg = res.text[:500] or "Pitchbook generation failed."
                        raise RuntimeError(f"{msg} (HTTP {res.status_code})")

                    # Store binary presentation in session state
                    st.session_state[f"pitchbook_bytes_{client_id}"] = res.content
                    st.session_state[f"pitchbook_filename_{client_id}"] = f"ING_{client_id}_Pitchbook.pptx"
                    st.success("✅ Institutional Pitchbook generated successfully!")

                except Exception as exc:
                    st.error(f"Generation error: {exc}")

    # Download Button Area
    pitchbook_data = st.session_state.get(f"pitchbook_bytes_{client_id}")
    if pitchbook_data:
        st.markdown("---")
        col_msg, col_btn = st.columns([2, 1])

        with col_msg:
            st.success(f"📥 **Ready to Download:** `ING_{client_id}_Pitchbook.pptx`")
            file_kb = len(pitchbook_data) / 1024
            st.caption(f"Payload Size: {file_kb:.1f} KB · 16:9 Widescreen Presentation")

        with col_btn:
            st.download_button(
                label="⬇️ Download Deck (.pptx)",
                data=pitchbook_data,
                file_name=st.session_state.get(
                    f"pitchbook_filename_{client_id}",
                    f"ING_{client_id}_Pitchbook.pptx",
                ),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"download_btn_active_{client_id}",
                use_container_width=True,
            )
    else:
        st.caption("💡 Click the button above to generate and download the customized client presentation.")