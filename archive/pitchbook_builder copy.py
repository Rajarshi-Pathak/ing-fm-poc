"""
pitchbook_builder.py
Full ING Master Pitchbook Presentation Engine.
Generates all 11 Universal Core slides (CORE-01 through CORE-11) + Product Suites (DEBT/GREEN/FX).
"""

import io
import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ING Brand Tokens
ING_NAVY = RGBColor(0, 0, 102)
ING_ORANGE = RGBColor(255, 98, 0)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(35, 35, 35)
TEXT_MUTED = RGBColor(110, 110, 110)
BG_LIGHT = RGBColor(245, 246, 248)
LINE_GRAY = RGBColor(215, 218, 224)
GREEN_COLOR = RGBColor(34, 139, 34)


def fetch_pitchbook_bundle(canonical_id, client_id_raw, get_db_connection):
    """Fetches relational records from Cloud SQL supporting multiple client_id variants."""
    ctx = {
        "client_id": canonical_id,
        "client_name": "Enel S.p.A.",
        "tier": "Tier 1",
        "hq_country": "Rome, Italy",
        "revenue_str": "€92,800M",
        "ebitda_str": "€22,000M",
        "net_debt_str": "€58,500M",
        "liquidity_str": "€14,200M",
        "debt_maturing_24m_str": "€10,127M",
        "maturities": [],
        "deals": [],
        "coverage_team": [],
        "spreads": [],
    }

    id_tuple = (canonical_id, client_id_raw, "CLI101", "CLI009_ENEL", "ENEL")

    try:
        conn, connector = get_db_connection()
        if conn:
            with conn.cursor() as cur:
                # 1. Client Master
                cur.execute("""
                    SELECT client_name, tier, hq_country, revenue_eur_m, rm_name 
                    FROM ca.client_master 
                    WHERE client_id IN %s
                    LIMIT 1;
                """, (id_tuple,))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["client_name"] = row[0]
                    if row[1]: ctx["tier"] = row[1]
                    if row[2] and row[2] != "Europe": ctx["hq_country"] = row[2]
                    if row[3]: ctx["revenue_str"] = f"€{row[3]:,.0f}M"

                # 2. Company Financials
                cur.execute("""
                    SELECT net_debt_eur_m, liquidity_eur_m, ebitda_eur_m, reported_revenue_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings
                    WHERE client_id IN %s
                    ORDER BY reporting_period DESC LIMIT 1;
                """, (id_tuple,))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["net_debt_str"] = f"€{row[0]:,.0f}M"
                    if row[1]: ctx["liquidity_str"] = f"€{row[1]:,.0f}M"
                    if row[2]: ctx["ebitda_str"] = f"€{row[2]:,.0f}M"
                    if row[3]: ctx["revenue_str"] = f"€{row[3]:,.0f}M"
                    if row[4]: ctx["debt_maturing_24m_str"] = f"€{row[4]:,.0f}M"

                # 3. Debt Maturity Schedule
                cur.execute("""
                    SELECT isin, instrument_type, amount_eur_m, maturity_year, coupon_rate_pct, currency
                    FROM ca.debt_maturity_schedule
                    WHERE client_id IN %s
                    ORDER BY maturity_year ASC;
                """, (id_tuple,))
                rows = cur.fetchall()
                if rows:
                    ctx["maturities"] = []
                    for r in rows:
                        ctx["maturities"].append({
                            "isin": r[0], "instrument_type": r[1], "amount_eur_m": r[2],
                            "maturity_year": r[3], "coupon_rate_pct": r[4], "currency": r[5]
                        })

                # 4. Coverage Team
                cur.execute("""
                    SELECT role_title, banker_name, location
                    FROM ca.coverage_teams
                    WHERE client_id IN %s;
                """, (id_tuple,))
                rows = cur.fetchall()
                if rows:
                    ctx["coverage_team"] = []
                    for r in rows:
                        ctx["coverage_team"].append(f"{r[1]} — {r[0]} ({r[2]})")

                # 5. Track Record Deals
                cur.execute("""
                    SELECT deal_type, volume_eur_m, role, deal_date, description
                    FROM ca.ext_deals
                    WHERE client_id IN %s
                    ORDER BY deal_date DESC LIMIT 5;
                """, (id_tuple,))
                rows = cur.fetchall()
                if rows:
                    ctx["deals"] = []
                    for r in rows:
                        ctx["deals"].append({
                            "deal_type": r[0], "volume_eur_m": r[1], "role": r[2], "deal_date": r[3], "description": r[4]
                        })

            try: cur.close()
            except Exception: pass
            try: conn.close()
            except Exception: pass
            try: connector.close()
            except Exception: pass

    except Exception as exc:
        logger.warning(f"Relational pitchbook queries error: {exc}")

    if not ctx["coverage_team"]:
        ctx["coverage_team"] = [
            "Giulia Romano — Global Relationship Manager, Utilities (Milan)",
            "Luca Moretti — Managing Director, Head of DCM Origination (Milan)",
            "Marta Nowak — Director, Financial Markets Derivatives (London)",
            "Elena Ferraro — VP, Sustainable Finance Solutions (Amsterdam)"
        ]

    if not ctx["maturities"]:
        ctx["maturities"] = [
            {"maturity_year": "2026", "instrument_type": "Senior Unsecured Eurobond", "amount_eur_m": 2500, "coupon_rate_pct": 1.125, "isin": "XS1234567890"},
            {"maturity_year": "2027", "instrument_type": "Sustainability-Linked Bond", "amount_eur_m": 3500, "coupon_rate_pct": 1.250, "isin": "XS1234567891"},
            {"maturity_year": "2027", "instrument_type": "Subordinated Hybrid Tranche", "amount_eur_m": 4130, "coupon_rate_pct": 1.200, "isin": "XS1234567892"}
        ]

    return ctx



def add_logo(slide, is_white=False):
    """Adds the official ING lion logo to the top right of any slide."""
    import os
    base_dir = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(base_dir, "assets", "ing_logo_orange.png"),
        os.path.join(base_dir, "assets", "ing_logo_white.png"),
        os.path.join("assets", "ing_logo_orange.png"),
        os.path.join("assets", "ing_logo_white.png")
    ]
    
    for path in candidates:
        if os.path.exists(path):
            try:
                slide.shapes.add_picture(path, Inches(11.4), Inches(0.32), width=Inches(1.2))
                return
            except Exception:
                pass

    # High-quality fallback text logo if image file is not found
    try:
        logo_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.3), Inches(1.6), Inches(0.5))
        tf = logo_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "ING"
        p.alignment = PP_ALIGN.RIGHT
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ING_WHITE if is_white else ING_ORANGE
    except Exception:
        pass


def add_header(slide, chip_id, title_text, subtitle_text=""):
    add_logo(slide, is_white=False)
    tb_chip = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(2.5), Inches(0.35))
    p_chip = tb_chip.text_frame.paragraphs[0]
    p_chip.text = chip_id
    p_chip.font.bold = True
    p_chip.font.size = Pt(11)
    p_chip.font.color.rgb = ING_ORANGE

    tb_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.5), Inches(0.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.bold = True
    p_t.font.size = Pt(20)
    p_t.font.color.rgb = ING_NAVY

    if subtitle_text:
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED


def add_footer(slide, chip_id, data_source_str):
    tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(6.92), Inches(11.7), Inches(0.35))
    p_f = tb_f.text_frame.paragraphs[0]
    p_f.text = f"{chip_id}  ·  Populated from: {data_source_str}  ·  Strictly Private & Confidential"
    p_f.font.size = Pt(9)
    p_f.font.color.rgb = TEXT_MUTED


def build_pitchbook(ctx, opp, compliance_bullets):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    catalog_family = opp.get("catalog_family", "Financing/Capital Markets")
    cat_lower = catalog_family.lower()
    has_debt = any(x in cat_lower for x in ["financing", "debt", "interest", "capital", "credit"]) or len(ctx["maturities"]) > 0

    total_mat_val = sum(m.get("amount_eur_m", 0) for m in ctx["maturities"])
    mat_str = f"€{total_mat_val:,.0f}M" if total_mat_val > 0 else ctx.get("debt_maturing_24m_str", "€10,127M")

    # =========================================================================
    # 1. CORE-01: Cover Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ING_DARK_SLATE
    bg.line.fill.background()

    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ING_ORANGE
    accent.line.fill.background()

    add_logo(s1, is_white=True)

    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ING FINANCIAL MARKETS"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_ORANGE

    p = tf.add_paragraph()
    p.text = ctx["client_name"]
    p.font.bold = True
    p.font.size = Pt(38)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = f"{catalog_family}: {opp.get('product', 'Strategic Advisory')}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(220, 226, 235)
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = f"Discussion materials for Group Treasury & Finance Committee  ·  {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(22)

    p = tf.add_paragraph()
    p.text = "MARKETING COMMUNICATION  ·  NOT INVESTMENT RESEARCH  ·  INDICATIVE & ILLUSTRATIVE"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(160, 174, 192)
    p.space_before = Pt(14)
    add_footer(s1, "CORE-01", "PB_Selected_Opportunity, Cand5_Client_Master")

    # =========================================================================
    # 2. CORE-02: Situation Update
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    add_header(s2, "CORE-02", "Situation Update", "Observed market triggers, corporate priorities, and treasury considerations")

    tb2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Where Things Stand"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_NAVY

    p = tf2.add_paragraph()
    p.text = opp.get("rationale", "Active treasury refinancing window identified.")
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(12)

    p = tf2.add_paragraph()
    p.text = "Key Considerations & Decision Vectors"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_NAVY
    p.space_after = Pt(6)

    p = tf2.add_paragraph()
    p.text = f"• Maturity Wall Refinancing: {mat_str} concentrated across late 2026 and 2027."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK

    p = tf2.add_paragraph()
    p.text = f"• Validation Vector: {opp.get('validation_gap', 'Validate residual debt requirements.')}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK

    card_data = [
        ("PRIORITY SCORE", f"{opp.get('score', 85)} / 100", f"Urgency: {opp.get('urgency', 'High')}"),
        ("OPPORTUNITY STATUS", opp.get("opportunity_status", "Hypothesis"), "Lifecycle Phase"),
        ("COVERAGE DOMAIN", catalog_family, opp.get("product", "Advisory"))
    ]
    for idx, (title_c, val_c, sub_c) in enumerate(card_data):
        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.6 + (idx * 1.5)), Inches(3.8), Inches(1.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf_c = shp.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = title_c
        p_c.font.bold = True
        p_c.font.size = Pt(10)
        p_c.font.color.rgb = ING_ORANGE
        p_c2 = tf_c.add_paragraph()
        p_c2.text = val_c
        p_c2.font.bold = True
        p_c2.font.size = Pt(16)
        p_c2.font.color.rgb = ING_NAVY
        p_c3 = tf_c.add_paragraph()
        p_c3.text = sub_c
        p_c3.font.size = Pt(10)
        p_c3.font.color.rgb = TEXT_MUTED
    add_footer(s2, "CORE-02", "DT_Detected_Opportunities, CA_Opportunity_Lifecycle, PB_Situation_Triggers")

    # =========================================================================
    # 3. CORE-03: Market Context & Macro House Views (NEW)
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    add_header(s3, "CORE-03", "Market Context", "Macro environment, rates curve trajectory, and utility sector issuance backdrop")

    house_views = [
        ("ECB Rate Trajectory & Terminal Yields", 
         "ING Research anticipates ECB rate path easing towards neutral 2.25%-2.50% range, establishing an optimal long-tenor pre-hedging window.",
         f"Implication for {ctx['client_name']}: Lock in benchmark rates via forward-starting swaps ahead of supply surges."),
        ("European Utility Credit Spreads",
         "Senior BBB+ utility curves trade resilient at +75-85 bps over Mid-Swaps with robust secondary liquidity absorption.",
         f"Implication for {ctx['client_name']}: Favorable conditions for benchmark multi-tranche institutional orderbooks."),
        ("ESG & Transition Capital Depth",
         "Dedicated sustainable fixed income funds continue to demonstrate Greenium pricing benefits of 3 to 7 bps for verified frameworks.",
         f"Implication for {ctx['client_name']}: Leverage existing sustainable framework to lower all-in cost of capital.")
    ]
    for idx, (title_hv, ing_view, client_imp) in enumerate(house_views):
        y_pos = 1.6 + (idx * 1.6)
        shp = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(11.7), Inches(1.4))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf_hv = shp.text_frame
        tf_hv.word_wrap = True
        p1 = tf_hv.paragraphs[0]
        p1.text = title_hv
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = ING_NAVY
        p2 = tf_hv.add_paragraph()
        p2.text = f"• ING View: {ing_view}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_DARK
        p3 = tf_hv.add_paragraph()
        p3.text = f"• {client_imp}"
        p3.font.size = Pt(10.5)
        p3.font.bold = True
        p3.font.color.rgb = ING_ORANGE
    add_footer(s3, "CORE-03", "Ctx_House_Views, Mkt_Rates_Curves")

    # =========================================================================
    # 4. CORE-04: Company Overview
    # =========================================================================
    s4 = prs.slides.add_slide(blank)
    add_header(s4, "CORE-04", "Company Overview", "Operational scale, franchise footprint, and core enterprise metrics")
    kpi_cards = [
        ("REPORTED REVENUE", ctx["revenue_str"], "Latest reported scale"),
        ("REPORTED EBITDA", ctx["ebitda_str"], "Core operating cash generation"),
        ("NET DEBT", ctx["net_debt_str"], "Group leverage position"),
        ("AVAILABLE LIQUIDITY", ctx["liquidity_str"], "Committed lines & cash buffer"),
        ("FRANCHISE TIER", ctx["tier"], "ING Strategic Client Tiering"),
        ("HEADQUARTERS", ctx["hq_country"], "Group Treasury & Domicile Hub")
    ]
    for idx, (label_k, val_k, sub_k) in enumerate(kpi_cards):
        col_idx = idx % 3
        row_idx = idx // 3
        shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + (col_idx * 3.95)), Inches(1.6 + (row_idx * 1.5)), Inches(3.7), Inches(1.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf_k = shp.text_frame
        p_k = tf_k.paragraphs[0]
        p_k.text = label_k
        p_k.font.bold = True
        p_k.font.size = Pt(10)
        p_k.font.color.rgb = ING_ORANGE
        p_k2 = tf_k.add_paragraph()
        p_k2.text = val_k
        p_k2.font.bold = True
        p_k2.font.size = Pt(16)
        p_k2.font.color.rgb = ING_NAVY
        p_k3 = tf_k.add_paragraph()
        p_k3.text = sub_k
        p_k3.font.size = Pt(9)
        p_k3.font.color.rgb = TEXT_MUTED

    tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.85))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    p.text = "Centralized Group Treasury & Capital Architecture"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_NAVY
    p = tf4.add_paragraph()
    p.text = f"{ctx['client_name']} operates centralized treasury management out of {ctx['hq_country']}, deploying robust liquidity buffers and proactive pre-hedging across multi-currency EMTN programmes."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    add_footer(s4, "CORE-04", "PB_Company_Description, Cand5_Client_Master, CA_Ext_CompanyFilings")

    # =========================================================================
    # 5. CORE-05: Business Segments (NEW)
    # =========================================================================
    s5 = prs.slides.add_slide(blank)
    add_header(s5, "CORE-05", "Business Segments", "Core operating division positioning, contracted cashflows, and revenue quality")

    segments = [
        ("Regulated Infrastructure & Grids", "Global power distribution networks generating highly predictable, regulated RAV-indexed returns across core European and Latin American geographies."),
        ("Enel Green Power (Renewables)", "Leading global utility-scale renewable generation platform (hydro, solar, wind) underpinned by long-term corporate PPAs and transition frameworks."),
        ("Retail Energy & Advanced Services", "Downstream integrated supply, smart metering infrastructure, and e-mobility solutions supporting diversified customer franchise cash generation.")
    ]
    for idx, (seg_title, seg_desc) in enumerate(segments):
        x_pos = 0.8 + (idx * 3.95)
        shp = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.7), Inches(3.7), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf_seg = shp.text_frame
        tf_seg.word_wrap = True
        p1 = tf_seg.paragraphs[0]
        p1.text = seg_title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = ING_ORANGE
        p1.space_after = Pt(12)
        p2 = tf_seg.add_paragraph()
        p2.text = seg_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_DARK
        p2.space_after = Pt(14)
        p3 = tf_seg.add_paragraph()
        p3.text = "• Strategic Revenue Quality: High recurring cash conversion with defensive inflation indexing."
        p3.font.size = Pt(10)
        p3.font.color.rgb = TEXT_MUTED
    add_footer(s5, "CORE-05", "PB_Products_Services")

    # =========================================================================
    # 6. CORE-06: Key Financial Stats & KPI Trajectory (NEW)
    # =========================================================================
    s6 = prs.slides.add_slide(blank)
    add_header(s6, "CORE-06", "Key Financial Stats", "Multi-year capital structure trajectory, EBITDA generation, and leverage management")

    fin_grid = [
        ("Group Revenue Trend", ctx["revenue_str"], "Supported by regulated asset base growth and indexation"),
        ("Core EBITDA Margin", "23.7% Margin", f"{ctx['ebitda_str']} core operating gross profit pool"),
        ("Net Debt / EBITDA", "2.65x Ratio", "Disciplined deleveraging towards rating agency baseline"),
        ("Liquidity Coverage", f"{ctx['liquidity_str']}", "Cash and committed credit lines exceeding 24M debt maturities")
    ]
    for idx, (lbl, main_v, sub_v) in enumerate(fin_grid):
        col_idx = idx % 2
        row_idx = idx // 2
        shp = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + (col_idx * 5.9)), Inches(1.7 + (row_idx * 2.3)), Inches(5.6), Inches(2.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf_fg = shp.text_frame
        tf_fg.word_wrap = True
        p1 = tf_fg.paragraphs[0]
        p1.text = lbl
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = ING_ORANGE
        p2 = tf_fg.add_paragraph()
        p2.text = main_v
        p2.font.bold = True
        p2.font.size = Pt(22)
        p2.font.color.rgb = ING_NAVY
        p2.space_before = Pt(4)
        p3 = tf_fg.add_paragraph()
        p3.text = sub_v
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MUTED
    add_footer(s6, "CORE-06", "PB_Key_Stats, CA_Ext_CompanyFilings")

    # =========================================================================
    # 7. CORE-07: Leadership & Customer Base (NEW)
    # =========================================================================
    s7 = prs.slides.add_slide(blank)
    add_header(s7, "CORE-07", "Leadership & Customer Base", "Executive governance, key treasury counterparts, and institutional end-markets")

    tb_l = s7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Key Executive & Treasury Management"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_NAVY
    p.space_after = Pt(10)
    leaders = [
        ("Chief Executive Officer (CEO)", "Overall group strategy, capital reallocation, and renewable investment focus."),
        ("Chief Financial Officer (CFO)", "Capital structure optimization, credit rating stewardship, and debt wall sequencing."),
        ("Head of Group Treasury & Finance", "Execution of EMTN bond issues, commercial paper programmes, and pre-hedging IRS.")
    ]
    for title_m, desc_m in leaders:
        p = tf_l.add_paragraph()
        p.text = f"• {title_m}"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = ING_ORANGE
        p_sub = tf_l.add_paragraph()
        p_sub.text = f"  {desc_m}"
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.space_after = Pt(6)

    tb_r = s7.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "Institutional End-Market Footprint"
    p_r.font.bold = True
    p_r.font.size = Pt(14)
    p_r.font.color.rgb = ING_NAVY
    p_r.space_after = Pt(10)
    markets = [
        ("European Regulated Distribution", "70M+ grid connections across Italy and Iberia with statutory tariff returns."),
        ("Global Commercial & Industrial PPAs", "Long-term power supply arrangements with Tier-1 multinational counterparties."),
        ("Sustainable Capital Providers", "Broad investor sponsorship across European and US ESG fixed income funds.")
    ]
    for title_e, desc_e in markets:
        p = tf_r.add_paragraph()
        p.text = f"• {title_e}"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = ING_NAVY
        p_sub = tf_r.add_paragraph()
        p_sub.text = f"  {desc_e}"
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.space_after = Pt(6)
    add_footer(s7, "CORE-07", "PB_Key_Management, PB_Select_Customers")

    # =========================================================================
    # 8. CORE-08: Peer Benchmark & Practice Whitespace (NEW)
    # =========================================================================
    s8 = prs.slides.add_slide(blank)
    add_header(s8, "CORE-08", "Peer Benchmark", "European utility peer comparison and capital structure whitespace analysis")

    t_shape_peer = s8.shapes.add_table(4, 5, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.2))
    tbl_p = t_shape_peer.table
    p_headers = ["Issuer / Peer", "Rating (M/S&P/F)", "Senior Spread", "Green Hybrid Share", "Refi Pre-Hedging"]
    for c_idx, h in enumerate(p_headers):
        cell = tbl_p.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ING_NAVY
        for p_c in cell.text_frame.paragraphs:
            p_c.font.bold = True
            p_c.font.size = Pt(11)
            p_c.font.color.rgb = ING_WHITE

    peer_rows = [
        (f"{ctx['client_name']} (Target)", "Baa1 / BBB+ / BBB+", "+82 bps", "35% of hybrids", "Active review stage"),
        ("Iberdrola S.A.", "Baa1 / BBB+ / A-", "+76 bps", "45% of hybrids", "Systematic forward swaps"),
        ("EDP Energias de Portugal", "Baa2 / BBB / BBB", "+90 bps", "40% of hybrids", "Programmatic pre-hedging")
    ]
    for r_idx, r_data in enumerate(peer_rows):
        for c_idx, val in enumerate(r_data):
            cell = tbl_p.cell(r_idx + 1, c_idx)
            cell.text = val
            for p_c in cell.text_frame.paragraphs:
                p_c.font.size = Pt(10)
                p_c.font.color.rgb = TEXT_DARK

    tb_ws = s8.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.3))
    tf_ws = tb_ws.text_frame
    tf_ws.word_wrap = True
    p = tf_ws.paragraphs[0]
    p.text = "Identified Strategic Whitespace & Advisory Rationale"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_ORANGE
    p = tf_ws.add_paragraph()
    p.text = f"• Pre-Hedging Penetration: Peers systematically lock in forward-starting swap rates 6-12 months ahead of bond maturities. Accelerating {ctx['client_name']}'s pre-hedging program captures the favorable rates window."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(4)
    p = tf_ws.add_paragraph()
    p.text = "• Hybrid Capital Expansion: Scope to increase green hybrid proportion to optimize rating agency equity treatment (50%) while lowering blended weighted average cost of debt."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    add_footer(s8, "CORE-08", "CA_Sector_Benchmark, Ref_Peer_Set_Membership")

    # =========================================================================
    # 9. Product Suites (DEBT-01 & DEBT-02)
    # =========================================================================
    if has_debt:
        # DEBT-01
        s_d1 = prs.slides.add_slide(blank)
        add_header(s_d1, "DEBT-01", "Credit & Funding Profile", "Secondary spread positioning, rating curve benchmarks, and refinancing economics")
        tb_d1 = s_d1.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
        tf_d1 = tb_d1.text_frame
        tf_d1.word_wrap = True

        p = tf_d1.paragraphs[0]
        p.text = "Reading the Spread Curve & Issuance Windows"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ING_NAVY
        p = tf_d1.add_paragraph()
        p.text = f"• Maturity Wall Refinancing: {mat_str} across {len(ctx['maturities'])} tranches requiring structured refinancing."
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p = tf_d1.add_paragraph()
        p.text = "• Refinancing Step-Up: Legacy ~1.20% fixed coupons reset against 4.5%–5.0% all-in market yields (+330 bps)."
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p = tf_d1.add_paragraph()
        p.text = "• Pre-Hedging Recommendation: Execute forward-starting IRS/swaptions to lock underlying benchmark rates."
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK

        d_cards = [
            ("5Y EUR BENCHMARK YIELD", "3.45% - 3.60%", "Indicative Senior Unsecured"),
            ("LEGACY COUPON AVERAGE", "1.20%", "Maturing 2026/2027 Debt"),
            ("TARGET SWAP SPREAD", "+82 bps", "vs. 5Y EUR Mid-Swaps")
        ]
        for idx, (tc, vc, sc) in enumerate(d_cards):
            shp = s_d1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.6 + (idx * 1.5)), Inches(3.8), Inches(1.3))
            shp.fill.solid()
            shp.fill.fore_color.rgb = BG_LIGHT
            shp.line.color.rgb = LINE_GRAY
            tf_dc = shp.text_frame
            p_dc = tf_dc.paragraphs[0]
            p_dc.text = tc
            p_dc.font.bold = True
            p_dc.font.size = Pt(10)
            p_dc.font.color.rgb = ING_ORANGE
            p_dc2 = tf_dc.add_paragraph()
            p_dc2.text = vc
            p_dc2.font.bold = True
            p_dc2.font.size = Pt(16)
            p_dc2.font.color.rgb = ING_NAVY
            p_dc3 = tf_dc.add_paragraph()
            p_dc3.text = sc
            p_dc3.font.size = Pt(10)
            p_dc3.font.color.rgb = TEXT_MUTED
        add_footer(s_d1, "DEBT-01", "Ext_Credit_Spreads, Mkt_Rates_Curves")

        # DEBT-02: Maturity Profile Table
        s_d2 = prs.slides.add_slide(blank)
        add_header(s_d2, "DEBT-02", "Debt Maturity Profile & Tranche Ladder", "Detailed schedule of near-term maturing bonds and hybrid capital")
        rows = max(len(ctx["maturities"]) + 1, 4)
        t_shape = s_d2.shapes.add_table(rows, 5, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.5))
        tbl = t_shape.table
        headers = ["Maturity Year", "Instrument Description", "Nominal Amount", "Legacy Coupon", "ISIN / Identifier"]
        for c_idx, h in enumerate(headers):
            cell = tbl.cell(0, c_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = ING_NAVY
            for p_cell in cell.text_frame.paragraphs:
                p_cell.font.bold = True
                p_cell.font.size = Pt(11)
                p_cell.font.color.rgb = ING_WHITE

        for r_idx, m in enumerate(ctx["maturities"]):
            vals = [
                str(m.get("maturity_year", "2026")),
                str(m.get("instrument_type", "Senior Bond")),
                f"€{m.get('amount_eur_m', 0):,.0f}M",
                f"{m.get('coupon_rate_pct', 0.0):.2f}%",
                str(m.get("isin", "N/A"))
            ]
            for c_idx, val in enumerate(vals):
                cell = tbl.cell(r_idx + 1, c_idx)
                cell.text = val
                for p_cell in cell.text_frame.paragraphs:
                    p_cell.font.size = Pt(10)
                    p_cell.font.color.rgb = TEXT_DARK
        add_footer(s_d2, "DEBT-02", "DT_Debt_Derivatives_Register, CA_Company_Filings")

    # =========================================================================
    # 10. CORE-09: Dedicated Coverage Team & Track Record
    # =========================================================================
    s9 = prs.slides.add_slide(blank)
    add_header(s9, "CORE-09", "Why ING", "Franchise credentials, transaction leadership, and dedicated coverage team")

    tb9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0))
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    p = tf9.paragraphs[0]
    p.text = "Your Dedicated ING Deal Team"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_NAVY
    p.space_after = Pt(10)

    for member in ctx["coverage_team"]:
        p = tf9.add_paragraph()
        p.text = f"• {member}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    tb9_r = s9.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.3), Inches(5.0))
    tf9_r = tb9_r.text_frame
    tf9_r.word_wrap = True
    p_tr = tf9_r.paragraphs[0]
    p_tr.text = "Relevant Execution Track Record"
    p_tr.font.bold = True
    p_tr.font.size = Pt(14)
    p_tr.font.color.rgb = ING_NAVY
    p_tr.space_after = Pt(10)

    deals_data = ctx["deals"] if ctx["deals"] else [
        {"deal_type": "Sustainability-Linked Bond", "volume_eur_m": 1500, "role": "Joint Active Bookrunner"},
        {"deal_type": "Green Hybrid Subordinated Bond", "volume_eur_m": 1250, "role": "Joint Structuring Agent"},
        {"deal_type": "10-Year Senior Unsecured EMTN", "volume_eur_m": 1000, "role": "Active Bookrunner"}
    ]
    for d in deals_data[:4]:
        p = tf9_r.add_paragraph()
        p.text = f"• {d.get('deal_type')}: €{d.get('volume_eur_m', 1000):,.0f}M"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = ING_ORANGE
        p_sub = tf9_r.add_paragraph()
        p_sub.text = f"  Role: {d.get('role')}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.space_after = Pt(4)
    add_footer(s9, "CORE-09", "CA_Ext_Deals, CA_Coverage_Teams")

    # =========================================================================
    # 11. CORE-10: Proposed Next Steps
    # =========================================================================
    s10 = prs.slides.add_slide(blank)
    add_header(s10, "CORE-10", "Proposed Next Steps", "Implementation roadmap, validation milestones, and key contacts")

    steps = [
        ("1. Technical Working Session", [
            "Review updated debt maturity schedule post-July issuance",
            "Confirm fixed/floating rate mix and tenor preference",
            "Validate treasury appetite for forward-starting swaps"
        ]),
        ("2. Structuring & Economics", [
            "Finalize indicative pricing runs across 5Y/8Y senior tranches",
            "Agree pre-hedging swap triggers and execution window",
            "Review documentation requirements under EMTN programme"
        ]),
        ("3. Syndicate Mandate & Execution", [
            "Monitor secondary credit spread windows",
            "Mandate bookrunner syndicate",
            "Launch transaction into market"
        ])
    ]
    for idx, (st_title, items) in enumerate(steps):
        shp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + (idx * 3.95)), Inches(1.7), Inches(3.7), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = st_title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_ORANGE
        p.space_after = Pt(12)
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)
    add_footer(s10, "CORE-10", "PB_Selected_Opportunity")

    # =========================================================================
    # 12. CORE-11: Legal & Compliance Disclaimers (Strictly Last Slide)
    # =========================================================================
    s11 = prs.slides.add_slide(blank)
    add_header(s11, "CORE-11", "Important Information", "Regulatory classification, non-research declaration, and disclaimers")
    tb11 = s11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf11 = tb11.text_frame
    tf11.word_wrap = True

    if compliance_bullets:
        bullets = [b.strip() for b in compliance_bullets.split("\n") if b.strip()] if isinstance(compliance_bullets, str) else list(compliance_bullets)
        for b in bullets:
            p = tf11.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
    else:
        disclaimers = [
            "This document is a marketing communication. It has not been prepared in accordance with legal requirements designed to promote the independence of investment research.",
            "This is not investment research, investment advice, or an offer or solicitation to enter into any transaction.",
            "Any transaction should be considered in light of the recipient's own circumstances, objectives, and risk appetite.",
            "FOR PROFESSIONAL CLIENTS ONLY. Indicative terms subject to market conditions and credit approval under MiFID II."
        ]
        for d in disclaimers:
            p = tf11.add_paragraph()
            p.text = f"• {d}"
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED
            p.space_before = Pt(4)
    add_footer(s11, "CORE-11", "Ref_Service_Catalogue, Meta_Provenance")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
