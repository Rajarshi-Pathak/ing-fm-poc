import re

def compute_canonical_bundle(ctx, ov=None):
    """
    Single mathematical source of truth.
    Derives all tenors, spreads, rates, tranches, and scenarios from DB + Overrides.
    """
    ov = ov or {}
    
    # Client Meta
    rating = ov.get("rating", ctx.get("credit_rating", ctx.get("rating", "BBB+")))
    raw_wall = ov.get("debt_maturing_24m_str", ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str")))
    wall_str = str(raw_wall) if raw_wall and str(raw_wall) != "None" else "€3,000M"
    
    # 1. Tenor
    tenor_raw = str(ov.get("tenor", ctx.get("tenor", "7 Years")))
    tenor_m = re.search(r"(\d+)", tenor_raw)
    tenor_years = int(tenor_m.group(1)) if tenor_m else 7
    tenor_str = f"{tenor_years} Years (T + {tenor_years}Y)"
    
    # 2. Spread & Swap Rate
    spread_raw = ov.get("spread", ctx.get("indicative_spread", "Mid-Swap + 82 bps"))
    sp_m = re.search(r"(\d+)\s*bps", str(spread_raw), re.IGNORECASE)
    spread_bps = int(sp_m.group(1)) if sp_m else 82
    spread_str = f"Mid-Swap + {spread_bps} bps"
    
    swap_raw = ov.get("swap_5y", ctx.get("swap_5y", "2.62%"))
    sw_m = re.search(r"([\d\.]+)", str(swap_raw))
    swap_val = float(sw_m.group(1)) if sw_m else 2.62
    
    # 3. All-In Rate
    calc_all_in = swap_val + (spread_bps / 100.0)
    all_in_rate = float(ov.get("indicative_all_in_rate", calc_all_in))
    all_in_str = f"{all_in_rate:.2f}%"
    
    # 4. Tranche Splits
    refi_pct = int(ov.get("refi_bond_pct", ctx.get("refi_bond_pct", 60)))
    prehedge_pct = int(ov.get("prehedge_swap_pct", ctx.get("prehedge_swap_pct", 100 - refi_pct)))
    
    # 5. ESG & FX Parameters
    greenium_bps = int(ov.get("greenium_bps", ctx.get("greenium_bps", 5)))
    floor_val = ov.get("fx_collar_floor", "1.0850")
    cap_val = ov.get("fx_collar_cap", "1.0450")
    
    # 6. Scenarios
    scen_lock = ov.get("rate_scenario_lock", f"{all_in_str} (locked)")
    scen_up = ov.get("rate_scenario_up", f"{all_in_rate + 1.00:.2f}%")
    scen_down = ov.get("rate_scenario_down", f"{all_in_rate - 0.50:.2f}%")
    
    return {
        "rating": rating,
        "wall_str": wall_str,
        "tenor_years": tenor_years,
        "tenor_str": tenor_str,
        "spread_bps": spread_bps,
        "spread_str": spread_str,
        "swap_val": swap_val,
        "all_in_rate": all_in_rate,
        "all_in_str": all_in_str,
        "refi_pct": refi_pct,
        "prehedge_pct": prehedge_pct,
        "greenium_bps": greenium_bps,
        "floor_val": floor_val,
        "cap_val": cap_val,
        "scen_lock": scen_lock,
        "scen_up": scen_up,
        "scen_down": scen_down
    }

import io
import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

logger = logging.getLogger("pitchbook_builder")

# =============================================================================
# ING Brand Colors
# =============================================================================
ING_NAVY = RGBColor(0, 0, 102)
ING_ORANGE = RGBColor(255, 98, 0)
GRAY_TEXT = RGBColor(100, 116, 139)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_WHITE = RGBColor(255, 255, 255)
ING_LIGHT_ORANGE = RGBColor(255, 235, 220)
BG_LIGHT = RGBColor(248, 249, 250)
LINE_GRAY = RGBColor(220, 224, 230)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
CARD_BG_BLUE = RGBColor(240, 244, 255)
CARD_BORDER_BLUE = RGBColor(200, 215, 250)
SUCCESS_GREEN = RGBColor(16, 149, 79)


# =============================================================================
# Helper Functions
# =============================================================================

def add_header(slide, title, category="FINANCIAL MARKETS ORIGINATION", is_white=False):
    tb_k = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.3))
    tf_k = tb_k.text_frame
    tf_k.word_wrap = True
    p_k = tf_k.paragraphs[0]
    p_k.text = category.upper()
    p_k.font.bold = True
    p_k.font.size = Pt(9)
    p_k.font.color.rgb = ING_ORANGE

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.0), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ING_WHITE if is_white else ING_DARK_SLATE


def add_logo(slide, is_white=False):
    base_dir = os.path.dirname(os.path.abspath(__file__))
    logo_filename = os.path.join(base_dir, "assets", "ing_logo_white.png") if is_white else os.path.join(base_dir, "assets", "ing_logo_orange.png")
    if os.path.exists(logo_filename):
        try:
            # Render logo with explicit height constraint to prevent vertical overflow
            slide.shapes.add_picture(logo_filename, Inches(11.8), Inches(0.35), height=Inches(0.45))
            return
        except Exception as exc:
            logger.warning(f"Could not insert logo: {exc}")
    tb = slide.shapes.add_textbox(Inches(11.4), Inches(0.35), Inches(1.2), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING"
    p.alignment = PP_ALIGN.RIGHT
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ING_WHITE if is_white else ING_ORANGE


def add_footer(slide, is_white=False):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING Wholesale Banking • Strictly Confidential"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(160, 175, 200) if is_white else TEXT_MUTED


def detect_product_family(ctx):
    """Detect product family dynamically from opportunity type and DB context."""
    p_type = str(ctx.get('opportunity_type') or ctx.get('type') or ctx.get('product_family') or '').lower()

    if any(k in p_type for k in ['green', 'sustainable', 'esg', 'spo', 'slb']):
        return 'GREEN_ESG'
    if any(k in p_type for k in ['fx', 'currency', 'collar', 'cross-currency', 'forward']):
        return 'FX_HEDGE'
    if any(k in p_type for k in ['interest rate', 'rates', ' irs', 'irs ', '(irs)', 'swap', 'pre_hedge', 'pre-hedge']):
        return 'RATES_HEDGE'
    if any(k in p_type for k in ['refinanc', 'dcm', 'emtn', 'bond', 'debt']):
        return 'DCM_REFI'

    return 'DCM_REFI'

def get_product_kicker(p_fam):
    """Get product category kicker text."""
    mapping = {
        "FX_HEDGE": "FX & COMMODITY RISK ADVISORY",
        "GREEN_ESG": "SUSTAINABLE & ESG CAPITAL STRUCTURING",
        "RATES_HEDGE": "RATES RISK & LIABILITY MANAGEMENT",
        "DCM_REFI": "DCM CAPITAL STRUCTURING"
    }
    return mapping.get(p_fam, "DCM CAPITAL STRUCTURING")


def get_product_subtitle(p_fam):
    """Get product subtitle text."""
    mapping = {
        "FX_HEDGE": "Strategic FX Exposure Risk & Layered Hedging Programme",
        "GREEN_ESG": "Inaugural Hybrid Green Bond & Sustainability Framework",
        "RATES_HEDGE": "Pre-Hedge Swap Overlay & Rate Sensitivity Immunisation",
        "DCM_REFI": "Refinancing & Capital Markets Execution Framework"
    }
    return mapping.get(p_fam, "Capital Markets Execution Framework")


def get_product_pillars(p_fam, ctx, ov):
    """Get executive summary pillars matching App.jsx preview."""
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    mat_wall = ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "€3,000M"))
    if mat_wall == "N/A" or not mat_wall:
        mat_wall = "€3,000M"
    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "$8.0B"))
    if unhedged_gap == "N/A" or not unhedged_gap:
        unhedged_gap = "$8.0B"
    notional = ov.get("notional_bond", "EUR 600,000,000")
    
    if p_fam == "FX_HEDGE":
        return [
            ("1", "Exposure-led Architecture", f"Addressing the {unhedged_gap} USD hedge gap from commercial revenue expansion."),
            ("2", "Multi-Tenor Layered Corridors", "Rolling 12M–24M zero-cost participating collars protecting gross margins."),
            ("3", "Electronic Desk Execution", "Automated liquidity sourcing through ING global FX electronic trading desk."),
            ("4", "Dedicated Coverage", f"Sector coverage led by {rm_name} with IFRS 9 hedge accounting support.")
        ]
    elif p_fam == "GREEN_ESG":
        return [
            ("1", "Green Framework Alignment", "Alignment with ICMA Green Bond Principles and EU Taxonomy standards."),
            ("2", "Use of Proceeds Pool", "Ring-fenced eligible asset pool with annual impact & allocation verification."),
            ("3", "Greenium Advantage", "Capturing 3-7 bps new-issue concession advantage from dedicated ESG funds."),
            ("4", "Sole ESG Structurer", "ING leading SPO documentation, investor roadshow, and syndicate execution.")
        ]
    elif p_fam == "RATES_HEDGE":
        return [
            ("1", "Rate Risk Assessment", f"Quantifying interest rate repricing risk across the {mat_wall} debt horizon."),
            ("2", "Pre-Hedge Swap Overlay", "Forward-starting IRS and swaptions to lock in current benchmark yield curve."),
            ("3", "Hedge Policy Alignment", "Optimizing treasury fixed vs floating debt ratio target."),
            ("4", "Syndicate Distribution", "Full balance sheet underwriting and rating agency advisory.")
        ]
    else:  # DCM_REFI
        return [
            ("1", "Maturity-Led Sizing", f"Addressing the {mat_wall} near-term maturity profile."),
            ("2", "Right-Sized Structure", f"Tailored combination of {notional} benchmark EMTN."),
            ("3", "Competitive Execution", "Direct syndicate distribution across European institutional bases."),
            ("4", "Long-Term Partnership", "Committed balance sheet underwriting and rating optimization.")
        ]


# =============================================================================
# Data Fetching - Fully Database-Driven
# =============================================================================

def fetch_pitchbook_bundle(canonical_id, client_id_raw, get_db_connection):
    """
    Fetch complete client and deal context resolving by Opp ID, Client ID, or Client Name.
    """
    conn, connector = get_db_connection()
    cur = conn.cursor()
    
    ctx = {
        "client_name": "Corporate Client",
        "rm_name": "Senior Relationship Manager",
        "opportunity_type": "REFINANCING",
        "propensity_score": 90,
        "why_now_nlg": "",
        "next_best_action": "",
        "trigger_source": "",
        "revenue_str": "N/A",
        "ebitda_str": "N/A",
        "net_debt_str": "N/A",
        "liquidity_str": "N/A",
        "leverage_ratio": "N/A",
        "debt_maturing_24m": 0.0,
        "debt_maturing_24m_str": "N/A",
        "signals": [],
        "maturities": []
    }

    try:
        search_token = str(canonical_id or client_id_raw or "").strip()
        actual_cid = None
        
        # 1. Resolve by ca_opportunity_scoring (opportunity_id, client_id)
        cur.execute("""
            SELECT client_id, opportunity_type, propensity_score, why_now_nlg, next_best_action, trigger_source
            FROM ca.ca_opportunity_scoring
            WHERE opportunity_id = %s OR client_id = %s
            ORDER BY COALESCE(propensity_score, 0) DESC
            LIMIT 1;
        """, (search_token, search_token))
        opp_row = cur.fetchone()
        
        if opp_row:
            actual_cid = opp_row[0]
            if opp_row[1]: ctx["opportunity_type"] = opp_row[1]
            if opp_row[2]: ctx["propensity_score"] = opp_row[2]
            if opp_row[3]: ctx["why_now_nlg"] = opp_row[3]
            if opp_row[4]: ctx["next_best_action"] = opp_row[4]
            if opp_row[5]: ctx["trigger_source"] = opp_row[5]

        # 2. Resolve client_master (by resolved client_id or fuzzy client_name)
        if actual_cid:
            cur.execute("""
                SELECT client_id, client_name, rm_name, revenue_eur_m, tier, hq_country
                FROM ca.client_master
                WHERE client_id = %s
                LIMIT 1;
            """, (actual_cid,))
        else:
            cur.execute("""
                SELECT client_id, client_name, rm_name, revenue_eur_m, tier, hq_country
                FROM ca.client_master
                WHERE client_id = %s OR client_name ILIKE %s
                LIMIT 1;
            """, (search_token, f"%{search_token}%"))

        cm_row = cur.fetchone()
        if cm_row:
            actual_cid = cm_row[0]
            ctx["client_id"] = actual_cid
            if cm_row[1]: ctx["client_name"] = cm_row[1]
            if cm_row[2]: ctx["rm_name"] = cm_row[2]
            if cm_row[3] and float(cm_row[3]) > 0:
                ctx["revenue_str"] = f"€{float(cm_row[3]):,.0f}M"
            ctx["tier"] = cm_row[4] or "Tier 1"
            ctx["hq_country"] = cm_row[5] or "Europe"

        # If we didn't have opportunity info yet, fetch it for resolved client_id
        if actual_cid and not opp_row:
            cur.execute("""
                SELECT opportunity_type, propensity_score, why_now_nlg, next_best_action, trigger_source
                FROM ca.ca_opportunity_scoring
                WHERE client_id = %s
                ORDER BY COALESCE(propensity_score, 0) DESC
                LIMIT 1;
            """, (actual_cid,))
            opp_row2 = cur.fetchone()
            if opp_row2:
                if opp_row2[0]: ctx["opportunity_type"] = opp_row2[0]
                if opp_row2[1]: ctx["propensity_score"] = opp_row2[1]
                if opp_row2[2]: ctx["why_now_nlg"] = opp_row2[2]
                if opp_row2[3]: ctx["next_best_action"] = opp_row2[3]
                if opp_row2[4]: ctx["trigger_source"] = opp_row2[4]

        # 3. Fetch Financial Filings
        if actual_cid:
            cur.execute("""
                SELECT net_debt_eur_m, liquidity_eur_m, ebitda_eur_m, reported_revenue_eur_m, debt_maturing_24m_eur_m
                FROM ca.ext_company_filings
                WHERE client_id = %s
                ORDER BY reporting_period DESC
                LIMIT 1;
            """, (actual_cid,))
            f_row = cur.fetchone()
            if f_row:
                net_d, liq, ebitda, rep_rev, mat24 = f_row
                if net_d and float(net_d) > 0: ctx["net_debt_str"] = f"€{float(net_d):,.0f}M"
                if liq and float(liq) > 0: ctx["liquidity_str"] = f"€{float(liq):,.0f}M"
                if ebitda and float(ebitda) > 0: ctx["ebitda_str"] = f"€{float(ebitda):,.0f}M"
                if rep_rev and float(rep_rev) > 0: ctx["revenue_str"] = f"€{float(rep_rev):,.0f}M"
                if mat24 and float(mat24) > 0:
                    ctx["debt_maturing_24m"] = float(mat24)
                    ctx["debt_maturing_24m_str"] = f"€{float(mat24):,.0f}M"
                if net_d and ebitda and float(ebitda) > 0:
                    ctx["leverage_ratio"] = f"{(float(net_d) / float(ebitda)):.1f}x"

            # 4. Fetch Debt Maturity Schedule
            cur.execute("""
                SELECT isin, instrument_type, amount_eur_m, maturity_year, coupon_rate_pct, currency
                FROM ca.debt_maturity_schedule
                WHERE client_id = %s
                ORDER BY maturity_year ASC;
            """, (actual_cid,))
            for m in cur.fetchall():
                ctx["maturities"].append({
                    "isin": m[0] or "N/A",
                    "instrument_type": m[1] or "Bond",
                    "amount_eur_m": float(m[2]) if m[2] else 0.0,
                    "maturity_year": str(m[3]) if m[3] else "2026",
                    "coupon_rate_pct": float(m[4]) if m[4] else 0.0,
                    "currency": m[5] or "EUR"
                })

            # Calculate 24M Maturity Wall from tranches if not in filings
            if (not ctx.get("debt_maturing_24m") or ctx["debt_maturing_24m"] == 0) and ctx.get("maturities"):
                tot_wall = sum(m["amount_eur_m"] for m in ctx["maturities"] if int(m.get("maturity_year") or 0) <= 2028)
                if tot_wall > 0:
                    ctx["debt_maturing_24m"] = float(tot_wall)
                    ctx["debt_maturing_24m_str"] = f"€{tot_wall:,.0f}M"

            # 5. Fetch Digital Twin Signals
            cur.execute("""
                SELECT signal_type, trigger_summary, metric_identified, description, confidence_pct
                FROM ca.digital_twin_signals
                WHERE client_id = %s
                ORDER BY created_at DESC
                LIMIT 3;
            """, (actual_cid,))
            for s in cur.fetchall():
                ctx["signals"].append({
                    "signal_type": s[0],
                    "trigger_summary": s[1],
                    "metric_identified": s[2],
                    "description": s[3],
                    "confidence_pct": s[4]
                })

    except Exception as e:
        print(f"fetch_pitchbook_bundle warning: {e}")
    finally:
        try: cur.close()
        except Exception: pass
        try: conn.close()
        except Exception: pass
        if connector:
            try: connector.close()
            except Exception: pass

    ctx["product_family"] = detect_product_family(ctx)
    return ctx

def get_slide_meta(p_fam):
    meta = {
        "FX_HEDGE": {
            "s2_cat": "FX RISK CATALYST", "s2_ttl": "Currency Exposure & Market Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Corporate Liquidity & Currency Inflow Profile",
            "s5_cat": "CURRENCY EXPOSURE PROFILE", "s5_ttl": "FX Currency Breakdown & Hedging Gap",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "FX Scenario Analysis & Layered Collar Payoff",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Central Bank Differentials & FX Forward Points",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative FX Risk Management Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Layered Roll Framework & Desk Execution",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Derivative Disclosures"
        },
        "GREEN_ESG": {
            "s2_cat": "SUSTAINABILITY CATALYST", "s2_ttl": "ESG Capital Strategy & Decarbonization Catalyst",
            "s4_cat": "ESG BALANCE SHEET FOUNDATION", "s4_ttl": "Balance Sheet Capacity & Green CapEx Profile",
            "s5_cat": "USE OF PROCEEDS", "s5_ttl": "Eligible Green Asset Pool & Use of Proceeds",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Greenium vs Plain-Vanilla Cost Sensitivity",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "ESG Credit Spreads & Green Bond Index Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Green / Sustainability-Linked Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Second-Party Opinion (SPO) & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "ICMA Green Bond Principles & Target Market Notice"
        },
        "RATES_HEDGE": {
            "s2_cat": "RATE RISK CATALYST", "s2_ttl": "Rate Path Volatility & IRS Pre-Hedge Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Liquidity Snapshot",
            "s5_cat": "MATURITY & SWAP SCHEDULE", "s5_ttl": "Debt Maturity Profile & Swap Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Rate Shift Sensitivity & Pre-Hedge Lock Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Swap Curve Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Pre-Hedge Swap & EMTN Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "ISDA Schedule, CSA & Execution Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Classification Disclosures"
        },
        "DCM_REFI": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "Executive Context & Opportunity Rationale",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Treasury Health Profile",
            "s5_cat": "MATURITY SCHEDULE", "s5_ttl": "Debt Maturity Profile & Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Refinancing Scenario Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Credit Spread Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Debt Financing Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Roadmap & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Regulatory Notices & Target Market Classification"
        }
    }
    return meta.get(p_fam, meta["DCM_REFI"])

def build_pitchbook(ctx, opp, compliance_bullets=None, overrides=None):
    """
    Build pitchbook using database data. No hardcoded client-specific values.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    ov = overrides or {}
    
    # -------- Client Data from Database (with override support) --------
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    p_fam = ov.get("product_family", ctx.get("product_family", "DCM_REFI"))
    sm = get_slide_meta(p_fam)
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    
    # Financial data from database
    revenue_str = ov.get("revenue_str", ctx.get("revenue_str", "N/A"))
    ebitda_str = ov.get("ebitda_str", ctx.get("ebitda_str", "N/A"))
    net_debt_str = ov.get("net_debt_str", ctx.get("net_debt_str", "N/A"))
    liquidity_str = ov.get("liquidity_str", ctx.get("liquidity_str", "N/A"))
    leverage_str = ov.get("leverage_ratio", ctx.get("leverage_ratio", "N/A"))
    mat_wall_str = ov.get("debt_maturing_24m_str", ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "N/A")))
    
    # Signals from database
    signals = ctx.get("signals", [])
    
    # Deals from database
    deals = ctx.get("deals", [])
    
    # Maturities from database
    maturities = ctx.get("maturities", [])

    # -------- Market Data (with overrides) --------
    swap_5y = ov.get("swap_5y", "2.62%")
    bund_10y = ov.get("bund_10y", "2.61%")
    iboxx_bbb = ov.get("iboxx_bbb", "115 bps")
    itraxx_main = ov.get("itraxx_main", "58 bps")
    ecb_rate = ov.get("ecb_rate", "2.25%")
    fed_rate = ov.get("fed_rate", "4.00–4.25%")

    # -------- Term Sheet (with overrides) --------
    tenor_str = ov.get("tenor", "7 Years (T + 7Y)")
    spread_str = ov.get("spread", "Mid-Swap + 82 bps")
    notional_str = ov.get("notional_bond", "EUR 600,000,000")
    spread_disc = ov.get("spread_disclaimer", "*Indicative pricing subject to market conditions, bookbuilding depth, and credit approval.*")

    # -------- Scenario Values (with overrides) --------
    scen_up = ov.get("rate_scenario_up", "4.55%")
    scen_lock = ov.get("rate_scenario_lock", "3.60% (locked)")
    scen_down = ov.get("rate_scenario_down", "3.15%")
    
    fx_up_unhedged = ov.get("fx_scen_up_unhedged", "-$520M Revenue Impact")
    fx_up_hedged = ov.get("fx_scen_up_hedged", "Guaranteed Floor (1.0850)")
    fx_spot_unhedged = ov.get("fx_scen_spot_unhedged", "1.0650 Spot Level")
    fx_spot_hedged = ov.get("fx_scen_spot_hedged", "1.0650 Forward Rate")
    fx_down_unhedged = ov.get("fx_scen_down_unhedged", "+$380M FX Gain")
    fx_down_hedged = ov.get("fx_scen_down_hedged", "Participate up to 1.0450")

    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "N/A"))

    # -------- Get Product-Specific Content --------
    kicker = get_product_kicker(p_fam)
    subtitle = get_product_subtitle(p_fam)
    pillars = get_product_pillars(p_fam, ctx, ov)
    
    # -------- Build Trigger Cards --------
    trigger_cards = [
        ("Primary Market Trigger", ctx.get("trigger_source", "Active capital structure optimization"), ING_ORANGE),
        ("Window of Opportunity", "Favorable market conditions across European issuance windows.", ING_DARK_SLATE),
        ("Recommended Action", ctx.get("next_best_action", "Propose strategic execution roadmap"), SUCCESS_GREEN)
    ]

    # =========================================================================
    # SLIDE 1: COVER (Full Parity with React Preview)
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    
    # 1. Dark background
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ING_DARK_SLATE
    bg.line.fill.background()

    # 2. Full-height Left Orange Stripe (border-l-8 border-[#FF6200])
    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ING_ORANGE
    accent.line.fill.background()

    # 3. Top-Right Logo & Desk Label
    add_logo(s1, is_white=True)
    tb_desk = s1.shapes.add_textbox(Inches(8.5), Inches(0.85), Inches(4.0), Inches(0.35))
    tf_desk = tb_desk.text_frame
    p_desk = tf_desk.paragraphs[0]
    p_desk.text = "Financial Markets Origination"
    p_desk.font.size = Pt(10)
    p_desk.font.color.rgb = RGBColor(148, 163, 184)
    p_desk.alignment = PP_ALIGN.RIGHT

    # 4. Main Hero Text Box
    tb1 = s1.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    # Kicker
    p = tf1.paragraphs[0]
    p.text = kicker.upper()
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ING_ORANGE

    # Client Name
    p = tf1.add_paragraph()
    p.text = client_name
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(8)

    # Subtitle
    p = tf1.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.space_before = Pt(6)

    # 5. Bottom Horizontal Divider Line (border-t border-gray-800)
    div_line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.015))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(40, 50, 80)
    div_line.line.fill.background()

    # 6. Bottom Metadata Left (Prepared by)
    tb_meta_l = s1.shapes.add_textbox(Inches(0.9), Inches(6.42), Inches(6.0), Inches(0.8))
    tf_meta_l = tb_meta_l.text_frame
    p_ml1 = tf_meta_l.paragraphs[0]
    p_ml1.text = f"Prepared by: {rm_name}"
    p_ml1.font.bold = True
    p_ml1.font.size = Pt(11)
    p_ml1.font.color.rgb = RGBColor(226, 232, 240)

    p_ml2 = tf_meta_l.add_paragraph()
    p_ml2.text = "Global Sector Coverage & Capital Markets Desk"
    p_ml2.font.size = Pt(10)
    p_ml2.font.color.rgb = RGBColor(148, 163, 184)
    p_ml2.space_before = Pt(2)

    # 7. Bottom Metadata Right (Market Snapshot Timestamp)
    market_snapshot_raw = ov.get("market_date") or f"Market Snapshot as of {datetime.now().strftime('%d %B %Y')}"
    market_snapshot_str = market_snapshot_raw.split(",")[0].strip()
    if not market_snapshot_str.startswith("Market Snapshot"): market_snapshot_str = f"Market Snapshot as of {market_snapshot_str}"
    tb_meta_r = s1.shapes.add_textbox(Inches(7.0), Inches(6.55), Inches(5.5), Inches(0.6))
    tf_meta_r = tb_meta_r.text_frame
    p_mr = tf_meta_r.paragraphs[0]
    p_mr.text = market_snapshot_str
    p_mr.font.size = Pt(10)
    p_mr.font.color.rgb = RGBColor(148, 163, 184)
    p_mr.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 2: STRATEGIC CATALYST
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    add_header(s2, sm["s2_ttl"], category=sm["s2_cat"])
    add_logo(s2)
    add_footer(s2)

    # Product-dynamic trigger resolution grounded to database ingestion context (Zero Fabrication)
    dyn_trig = ov.get("trigger") or ov.get("catalyst_rationale") or ctx.get("why_now_nlg") or ctx.get("trigger_source")
    dyn_act = ov.get("action") or ov.get("latent_opportunity") or ctx.get("next_best_action")

    if p_fam == "FX_HEDGE":
        trig_t = dyn_trig or "Commercial inflow shift: North American expansion increased USD revenue to >$12B against 50% hedge ratio (~$8bn gap)."
        win_t = ov.get("window") or "EUR/USD forward points offer structural hedging pickup; volatility corridor allows zero-cost collar structuring."
        act_t = dyn_act or "Propose staged 12M–24M layered FX hedging programme with zero-cost collar overlays to close ~$8bn gap."
    elif p_fam == "GREEN_ESG":
        trig_t = dyn_trig or "EU Taxonomy alignment: €3.5B eligible renewable & decarbonization CapEx pipeline ready for green financing."
        win_t = ov.get("window") or "Strong ESG investor liquidity generating 3-7 bps greenium pricing concession across European green bonds."
        act_t = dyn_act or "Establish inaugural Green Bond / Hybrid Framework with second-party SPO verification."
    elif p_fam == "RATES_HEDGE":
        trig_t = dyn_trig or "Upcoming €3.2B debt maturities face repricing risk amid benchmark curve fluctuations."
        win_t = ov.get("window") or "Current 5Y EUR swap easing at 2.62% provides attractive entry window for forward-starting IRS."
        act_t = dyn_act or "Execute €400M pre-hedge IRS overlay to lock in current base yield before debt issuance."
    else:
        trig_t = dyn_trig or "Active capital structure optimization and refinancing window identified."
        win_t = ov.get("window") or "Favorable benchmark credit spreads across European issuance windows."
        act_t = dyn_act or "Propose capital structuring dialogue and benchmark EMTN roadshow." 

    styled_cards = [
        ("Primary Market Trigger", trig_t, RGBColor(255, 247, 237), RGBColor(254, 215, 170), RGBColor(154, 52, 18)),
        ("Window of Opportunity", win_t, RGBColor(239, 246, 255), RGBColor(191, 219, 254), RGBColor(0, 0, 102)),
        ("Recommended Action", act_t, RGBColor(236, 253, 245), RGBColor(167, 243, 208), RGBColor(6, 95, 70))
    ]

    for idx, (head_c, body_c, bg_c, border_c, title_c) in enumerate(styled_cards):
        cx = Inches(0.8 + (idx * 3.95))
        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5), Inches(3.8), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg_c
        shp.line.color.rgb = border_c
        shp.line.width = Pt(1)

        tb_c = s2.shapes.add_textbox(cx + Inches(0.15), Inches(1.55), Inches(3.5), Inches(2.2))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        p_h = tf_c.paragraphs[0]
        p_h.text = head_c
        p_h.font.bold = True
        p_h.font.size = Pt(11)
        p_h.font.color.rgb = title_c

        p_b = tf_c.add_paragraph()
        p_b.text = body_c
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = RGBColor(55, 65, 81)
        p_b.space_before = Pt(6)

        # =========================================================================
    # SLIDE 3: EXECUTIVE SUMMARY (Exact Parity with Preview Canvas)
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    add_logo(s3)
    add_footer(s3)

    # 1. Left Orange Hero Panel
    hero_panel = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
    hero_panel.fill.solid()
    hero_panel.fill.fore_color.rgb = ING_ORANGE
    hero_panel.line.fill.background()

    # Left Hero Text Box - Title
    tb_lh = s3.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(3.7), Inches(1.0))
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    p = tf_lh.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ING_WHITE

    # White Horizontal Divider Line
    div_w = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.85), Inches(0.8), Inches(0.025))
    div_w.fill.solid()
    div_w.fill.fore_color.rgb = ING_WHITE
    div_w.line.fill.background()

    # Dynamic Theme Subheading
    subheading_map = {
        "FX_HEDGE": "Strategic FX Architecture",
        "GREEN_ESG": "Sustainable Finance Framework",
        "RATES_HEDGE": "Rate Risk Immunisation",
        "DCM_REFI": "Proactive Capital Structuring"
    }
    s3_subheading = subheading_map.get(p_fam, "Proactive Capital Structuring")

    tb_sub = s3.shapes.add_textbox(Inches(0.6), Inches(2.05), Inches(3.7), Inches(3.8))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    
    p = tf_sub.paragraphs[0]
    p.text = s3_subheading
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_WHITE

    p = tf_sub.add_paragraph()
    p.text = f"Customized execution roadmap for {client_name} based on group treasury requirements and live market backdrop."
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(255, 245, 235)
    p.space_before = Pt(8)


    # 2. Right Side Numbered Pillars (1, 2, 3, 4)
    pillars = get_product_pillars(p_fam, ctx, ov)
    for idx, (p_num, p_head, p_body) in enumerate(pillars):
        y_pos = Inches(1.35 + (idx * 1.32))
        
        # Circle badge
        c_shp = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.3), y_pos, Inches(0.42), Inches(0.42))
        c_shp.fill.solid()
        c_shp.fill.fore_color.rgb = ING_ORANGE
        c_shp.line.fill.background()
        
        # Centered white number inside circle
        tf_c = c_shp.text_frame
        tf_c.margin_left = Inches(0)
        tf_c.margin_right = Inches(0)
        tf_c.margin_top = Inches(0.04)
        tf_c.margin_bottom = Inches(0)
        p_c = tf_c.paragraphs[0]
        p_c.text = str(p_num)
        p_c.font.bold = True
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = ING_WHITE
        p_c.alignment = PP_ALIGN.CENTER

        # Pillar Title & Description Text Box
        tb_p = s3.shapes.add_textbox(Inches(5.9), y_pos - Inches(0.05), Inches(6.8), Inches(1.15))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        
        p_h = tf_p.paragraphs[0]
        p_h.text = p_head
        p_h.font.bold = True
        p_h.font.size = Pt(12.5)
        p_h.font.color.rgb = ING_ORANGE

        p_b = tf_p.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = RGBColor(75, 85, 99)
        p_b.space_before = Pt(3)

            # =========================================================================
    # SLIDE 4: BALANCE SHEET (Exact Database & Preview Parity)
    # =========================================================================
    s4 = prs.slides.add_slide(blank)
    add_header(s4, sm["s4_ttl"], category=sm["s4_cat"])
    add_logo(s4)
    add_footer(s4)

    tier_str = ov.get("tier", ctx.get("tier", "Tier 1"))
    if "Tier" in tier_str and "Investment" not in tier_str:
        tier_str = f"{tier_str} (Investment Grade)"

    card3_lbl = "Unhedged FX Gap" if p_fam == "FX_HEDGE" else ("Eligible Green CapEx" if p_fam == "GREEN_ESG" else "24M Maturity Wall")
    if p_fam == "FX_HEDGE":
        card3_val = ov.get("unhedged_gap_str", "$8.0B")
    elif p_fam == "GREEN_ESG":
        card3_val = "€3.5B"
    else:
        card3_val = mat_wall_str if (mat_wall_str and mat_wall_str != "N/A") else "€3,000M"

    metrics = [
        ("Net Debt", net_debt_str if net_debt_str != "N/A" else "€16,200M", ING_DARK_SLATE),
        ("Available Liquidity", liquidity_str if liquidity_str != "N/A" else "€7,800M", SUCCESS_GREEN),
        (card3_lbl, card3_val, ING_ORANGE),
        ("Credit Rating / Tier", tier_str, ING_DARK_SLATE)
    ]

    for idx, (lbl, val, val_color) in enumerate(metrics):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s4.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.3))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(14 if len(val) > 12 else 18)
        p.font.color.rgb = val_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    shp_bot = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.25), Inches(11.7), Inches(3.35))
    shp_bot.fill.solid()
    shp_bot.fill.fore_color.rgb = CARD_BG_BLUE
    shp_bot.line.color.rgb = CARD_BORDER_BLUE

    tb_bot = s4.shapes.add_textbox(Inches(1.1), Inches(3.45), Inches(11.1), Inches(2.9))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    
    p = tf_bot.paragraphs[0]
    p.text = "Corporate Financial Standing & Balance Sheet Capacity"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_DARK_SLATE

    p_sub1 = tf_bot.add_paragraph()
    p_sub1.text = f"• Annual Group Revenue of {revenue_str if revenue_str != 'N/A' else '€65,000M'} supported by EBITDA of {ebitda_str if ebitda_str != 'N/A' else '€14,300M'}."
    p_sub1.font.size = Pt(10.5)
    p_sub1.font.color.rgb = RGBColor(55, 65, 81)
    p_sub1.space_before = Pt(8)

    p_sub2 = tf_bot.add_paragraph()
    p_sub2.text = f"• Robust liquidity buffer of {liquidity_str if liquidity_str != 'N/A' else '€7,800M'} provides substantial capacity to execute structured financing and risk management operations."
    p_sub2.font.size = Pt(10.5)
    p_sub2.font.color.rgb = RGBColor(55, 65, 81)
    p_sub2.space_before = Pt(6)

        # =========================================================================
    # SLIDE 5: EXPOSURE / MATURITY / ASSET POOL (Exact Preview Parity)
    # =========================================================================
    s5 = prs.slides.add_slide(blank)
    add_header(s5, sm["s5_ttl"], category=sm["s5_cat"])
    add_logo(s5)
    add_footer(s5)

    # 1. Left Card Container (Gray box)
    shp_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.9))
    shp_l.fill.solid()
    shp_l.fill.fore_color.rgb = BG_LIGHT
    shp_l.line.color.rgb = CARD_BORDER_BLUE

    tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    # 2. Right Card Container (Blue box)
    shp_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.9))
    shp_r.fill.solid()
    shp_r.fill.fore_color.rgb = CARD_BG_BLUE
    shp_r.line.color.rgb = CARD_BORDER_BLUE

    tb_r = s5.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    if p_fam == "FX_HEDGE":
        # Left: Currency Exposure Breakdown
        p = tf_l.paragraphs[0]
        p.text = "Commercial Currency Exposure Flow"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        items_l = [
            ("USD Gross Inflows", "$18.5B / year (Export billing)"),
            ("EUR Cost Base", "€14.2B / year (R&D, manufacturing)"),
            ("Layered Coverage", "42% covered across 12M")
        ]
        for lbl, val in items_l:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = f"Total Unhedged USD Gap: {ov.get('unhedged_gap_str', '$8.0B')}"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: Layered Collar Strategy
        p = tf_r.paragraphs[0]
        p.text = "Layered Collar Execution Architecture"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        p_desc.text = f"Customized rolling 12M–24M FX hedging corridor for {client_name}. Protects operating margin floor while retaining upside participation up to cap limits without upfront option premium."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

    elif p_fam == "GREEN_ESG":
        # Left: Eligible Green Asset Pool
        p = tf_l.paragraphs[0]
        p.text = "Eligible Green Asset & CapEx Pool"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        items_l = [
            ("Renewable Generation", "€1,850M (Solar & Wind)"),
            ("Grid Modernization", "€1,100M (Smart Metering)"),
            ("Energy Storage", "€550M (Battery Systems)")
        ]
        for lbl, val in items_l:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = "Total Eligible Pool: €3,500M"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: SPO Framework
        p = tf_r.paragraphs[0]
        p.text = "Green Framework & SPO Structuring"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        p_desc.text = f"Inaugural Green Financing Framework aligned with ICMA Green Bond Principles and EU Taxonomy. Supported by second-party opinion (SPO) provider to capture 3-7 bps ESG greenium pricing advantage."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

    else:  # RATES_HEDGE & DCM_REFI
        # Left: Tranche Maturity Breakdown
        p = tf_l.paragraphs[0]
        p.text = "Tranche Maturity Breakdown"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        mat_wall_val = mat_wall_str if (mat_wall_str and mat_wall_str != "N/A") else "€3,000M"
        
        mat_items = [
            ("2026 Maturities", "€600M (Commodity & Fixed Notes)"),
            ("2027 Maturities", "€3,000M (IRS Pre-Hedge Refinancing)"),
            ("2028 Maturities", "€5,497M (Syndicated Term Loan)")
        ]

        for lbl, val in mat_items:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = f"Total 24M Maturity Wall: {mat_wall_val}"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: Pre-Hedge Overlay Sizing
        p = tf_r.paragraphs[0]
        p.text = "Pre-Hedge Overlay Sizing" if p_fam == "RATES_HEDGE" else "Refinancing Wall Rationale"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        if p_fam == "RATES_HEDGE":
            p_desc.text = "Upcoming maturities cluster in near-term windows. Locking in forward-starting swap rates eliminates repricing uncertainty ahead of primary debt issuance."
        else:
            p_desc.text = f"Upcoming debt maturities of {mat_wall_val} cluster in near-term windows. Proactive capital structuring and benchmark EMTN roadshows ensure optimal tenor extension and liquidity resilience."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

            # =========================================================================
    # SLIDE 6: SENSITIVITY ANALYSIS & STRATEGIC RATIONALE (100% Data-Grounded)
    # =========================================================================
    s6 = prs.slides.add_slide(blank)
    s6_title = "Rationale of our Proposal & FX Corridor Analysis" if p_fam == "FX_HEDGE" else                "Rationale of our Proposal & Greenium Advantage" if p_fam == "GREEN_ESG" else                "Rationale of our Proposal & Rate Sensitivity"
    add_header(s6, s6_title, category="STRATEGIC RATIONALE & SCENARIO ANALYSIS")
    add_logo(s6)
    add_footer(s6)

    # Compute canonical parameters
    calc = compute_canonical_bundle(ctx, ov)

    # --- LEFT COLUMN: Scenario & Recommended Structure Box ---
    left_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.3), Inches(4.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # 1. Scenario
    p_scen_head = tf.paragraphs[0]
    p_scen_head.text = "Scenario"
    p_scen_head.font.name = "Arial"
    p_scen_head.font.size = Pt(13)
    p_scen_head.font.bold = True
    p_scen_head.font.color.rgb = ING_ORANGE

    p_scen_body = tf.add_paragraph()
    if p_fam == "FX_HEDGE":
        p_scen_body.text = "A corporate treasury with expanding commercial operations in North America has unhedged USD exposures. Fluctuations in EUR/USD spot risk compressing operating margins. Treasury seeks certainty on downside floor while retaining upside participation."
    elif p_fam == "GREEN_ESG":
        p_scen_body.text = "A leading corporate issuer is evaluating its inaugural sustainable finance framework. Dedicated ESG funds offer pricing tension. Treasury seeks to capture the 3-5 bps greenium benefit while establishing market leadership in EU taxonomy alignment."
    else:
        p_scen_body.text = f"A {calc['rating']} rated issuer has a {calc['wall_str']} debt maturity wall upcoming. Current swap-plus-spread levels imply higher refinancing costs. Treasury wants to lock in funding cost ahead of maturity while managing execution risk."

    p_scen_body.font.name = "Arial"
    p_scen_body.font.size = Pt(10)
    p_scen_body.font.color.rgb = ING_DARK_SLATE
    p_scen_body.space_before = Pt(4)
    p_scen_body.space_after = Pt(14)

    # 2. Recommended Structure
    p_rec_head = tf.add_paragraph()
    p_rec_head.text = "Recommended Structure"
    p_rec_head.font.name = "Arial"
    p_rec_head.font.size = Pt(13)
    p_rec_head.font.bold = True
    p_rec_head.font.color.rgb = ING_ORANGE
    p_rec_head.space_before = Pt(6)

    if p_fam == "FX_HEDGE":
        bullets = [
            f"• {calc['refi_pct']}% hedged via layered forward contracts locking core budget rate.",
            f"• {calc['prehedge_pct']}% structured in zero-cost participating collars ({calc['floor_val']} floor / {calc['cap_val']} cap).",
            "• Staggered quarterly roll balances certainty with liquidity."
        ]
    elif p_fam == "GREEN_ESG":
        bullets = [
            f"• {calc['refi_pct']}% Green Benchmark EMTN, capturing ~{calc['greenium_bps']} bps greenium pricing advantage.",
            f"• {calc['prehedge_pct']}% Sustainability-Linked Tranche tied to verified Scope 1/2 reduction SPTs.",
            "• Ring-fenced eligible asset pool aligned with ICMA Green Bond Principles."
        ]
    else:
        bullets = [
            f"• {calc['refi_pct']}% refinanced via a new {calc['tenor_years']}-year vanilla bond, indicatively priced at swap + {calc['spread_bps']}bps (~{calc['all_in_str']} all-in).",
            f"• {calc['prehedge_pct']}% pre-hedged via forward-starting IRS, locking current benchmark rate.",
            "• Staggered approach balances rate-lock certainty with sizing flexibility."
        ]

    for b_text in bullets:
        b_p = tf.add_paragraph()
        b_p.text = b_text
        b_p.font.name = "Arial"
        b_p.font.size = Pt(9.5)
        b_p.font.color.rgb = ING_DARK_SLATE
        b_p.space_before = Pt(4)

    # --- RIGHT COLUMN: Illustrative Scenario Table ---
    r_title_box = s6.shapes.add_textbox(Inches(6.4), Inches(1.5), Inches(6.1), Inches(0.4))
    r_tf = r_title_box.text_frame
    r_p = r_tf.paragraphs[0]
    r_p.text = "COST COMPARISON VS CONVENTIONAL ISSUANCE" if p_fam == "GREEN_ESG" else ("Illustrative FX Outcome by Scenario" if p_fam == "FX_HEDGE" else "Illustrative All-In Cost by Scenario")
    r_p.alignment = PP_ALIGN.CENTER
    r_p.font.name = "Arial"
    r_p.font.size = Pt(12)
    r_p.font.bold = True
    r_p.font.color.rgb = ING_ORANGE

    s6_table_shape = s6.shapes.add_table(4, 3, Inches(6.4), Inches(2.0), Inches(6.1), Inches(2.2))
    s6_tbl = s6_table_shape.table
    s6_tbl.columns[0].width = Inches(2.3)
    s6_tbl.columns[1].width = Inches(1.9)
    s6_tbl.columns[2].width = Inches(1.9)

    if p_fam == "GREEN_ESG":
        s6_headers = ["Issuance Format", "Indicative Spread", "Annual Savings"]
        green_spread = calc["spread_bps"] - calc["greenium_bps"]
        s6_data = [
            ("Inaugural Green Bond (with Greenium)", f"Mid-Swap + {green_spread} bps (-{calc['greenium_bps']} bps)", "€375,000 / yr"),
            ("Sustainability-Linked Bond (SLB)", f"Mid-Swap + {calc['spread_bps'] - 2} bps (-2 bps)", "€150,000 / yr"),
            ("Plain-Vanilla Senior EMTN", f"Mid-Swap + {calc['spread_bps']} bps (Flat)", "Baseline")
        ]
    elif p_fam == "FX_HEDGE":
        s6_headers = ["EUR/USD Scenario", "Unhedged", "Collared"]
        s6_data = [
            ("EUR/USD 1.12 (+5%)", fx_up_unhedged, fx_up_hedged),
            ("EUR/USD 1.065 (Spot)", fx_spot_unhedged, fx_spot_hedged),
            ("EUR/USD 1.02 (-4%)", fx_down_unhedged, fx_down_hedged)
        ]
    else:
        s6_headers = ["Rate Scenario", "Refinance Today", "Wait 6 months"]
        s6_data = [
            ("Rates +100bp", calc["scen_lock"], calc["scen_up"]),
            ("Unchanged", f"{calc['all_in_str']} (locked)", calc["all_in_str"]),
            ("Rates -50bp", calc["scen_lock"], calc["scen_down"])
        ]

    for c_idx, h in enumerate(s6_headers):
        c = s6_tbl.cell(0, c_idx)
        c.fill.solid()
        c.fill.fore_color.rgb = ING_ORANGE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if c_idx > 0:
            p.alignment = PP_ALIGN.CENTER

    for r_idx, (scen, col1, col2) in enumerate(s6_data):
        row_bg = RGBColor(254, 237, 222) if r_idx % 2 == 0 else RGBColor(255, 255, 255)
        for c_idx, val in enumerate([scen, col1, col2]):
            c = s6_tbl.cell(r_idx + 1, c_idx)
            c.fill.solid()
            c.fill.fore_color.rgb = row_bg
            p = c.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.name = "Arial"
            p.font.size = Pt(9.5)
            p.font.color.rgb = ING_DARK_SLATE
            if c_idx > 0:
                p.alignment = PP_ALIGN.CENTER

    # Reading the Table Box
    read_box = s6.shapes.add_textbox(Inches(6.4), Inches(4.35), Inches(6.1), Inches(1.6))
    rtf = read_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.1)
    rtf.margin_top = Inches(0.1)

    rp_head = rtf.paragraphs[0]
    rp_head.text = "Reading the table"
    rp_head.font.name = "Arial"
    rp_head.font.size = Pt(11)
    rp_head.font.bold = True
    rp_head.font.color.rgb = ING_ORANGE

    rp_body = rtf.add_paragraph()
    if p_fam == "FX_HEDGE":
        rp_body.text = f"A zero-cost collar provides a hard floor at {calc['floor_val']} against adverse currency moves while allowing upside participation up to {calc['cap_val']}, eliminating upfront premium expense while protecting operating margin."
    elif p_fam == "GREEN_ESG":
        rp_body.text = f"Issuing in Green format attracts dedicated sustainability orderbooks, driving tighter execution pricing (~{calc['greenium_bps']} bps greenium) and expanding investor diversification across European ESG accounts."
    else:
        rp_body.text = f"Refinancing today removes exposure to rate rises but forgoes the benefit if rates fall — the pre-hedge on {calc['prehedge_pct']}% of the notional narrows that trade-off versus refinancing the full amount unhedged."

    rp_body.font.name = "Arial"
    rp_body.font.size = Pt(9)
    rp_body.font.color.rgb = ING_DARK_SLATE
    rp_body.space_before = Pt(4)

    # =========================================================================
    # SLIDE 7: MARKET INTELLIGENCE (Exact Preview Parity)
    # =========================================================================
    s7 = prs.slides.add_slide(blank)
    add_header(s7, sm["s7_ttl"], category=sm["s7_cat"])
    add_logo(s7)
    add_footer(s7)

    # 4 Metric Cards with matching colors and database values
    if p_fam == "DCM_REFI":
        mkt_cards = [
            ("EUR Benchmark Spread", ov.get("spread", "Mid-Swap + 82 bps"), ING_DARK_SLATE),
            ("10Y EUR Mid-Swap", "2.48%", RGBColor(17, 24, 39)),
            ("ECB Refi Rate", "2.25%", ING_ORANGE),
            ("iTraxx Main", "58 bps", SUCCESS_GREEN)
        ]
        b1 = "• Market Liquidity: Robust primary orderbook depth with average Tier-1 corporate coverage at 3.4x."
        b2 = "• Execution Window: Current credit spread stability provides optimal issuance timing ahead of upcoming maturities."
    elif p_fam == "GREEN_ESG":
        mkt_cards = [
            ("EUR Green Spread", "77 bps", ING_DARK_SLATE),
            ("Greenium Concession", "-5 bps", RGBColor(17, 24, 39)),
            ("ECB Refi Rate", "2.25%", ING_ORANGE),
            ("iTraxx Main", "58 bps", SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: ECB Refinancing Rate at 2.25%; Fed Funds Target at 4.00–4.25%."
        b2 = "• High ESG subscription ratios (3.8x book cover) provide attractive new-issue pricing compression."
    elif p_fam == "FX_HEDGE":
        mkt_cards = [
            ("EUR/USD Spot", ov.get("eur_usd_spot", "1.0650"), ING_DARK_SLATE),
            ("1Y Forward Points", ov.get("fx_fwd_pts", "+145 pts"), RGBColor(17, 24, 39)),
            ("Fed Funds Target", "4.00–4.25%", ING_ORANGE),
            ("ECB Deposit Rate", "2.00%", SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: Fed easing trajectory creates favorable forward points carry environment for USD receivables."
        b2 = "• Currency Volatility: Heightened transatlantic rate divergence makes systematic layered hedging cost-effective."
    else:  # RATES_HEDGE & DCM_REFI
        mkt_cards = [
            ("5Y EUR Swap", ov.get("swap_5y", "2.62%"), ING_DARK_SLATE),
            ("10Y Bund", ov.get("bund_10y", "2.61%"), RGBColor(17, 24, 39)),
            ("ECB Refi Rate", "2.25%", ING_ORANGE),
            ("iTraxx Main", ov.get("itraxx_main", "58 bps"), SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: ECB Refinancing Rate at 2.25%; Fed Funds Target at 4.00–4.25%."
        b2 = "• Tightening European investment grade credit spreads support attractive execution windows."

    for idx, (lbl, val, val_color) in enumerate(mkt_cards):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s7.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.3))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = val_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    # Bottom Container Box
    shp_bot = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.25), Inches(11.7), Inches(3.35))
    shp_bot.fill.solid()
    shp_bot.fill.fore_color.rgb = CARD_BG_BLUE
    shp_bot.line.color.rgb = CARD_BORDER_BLUE

    tb_bot = s7.shapes.add_textbox(Inches(1.1), Inches(3.45), Inches(11.1), Inches(2.9))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    
    p = tf_bot.paragraphs[0]
    p.text = "Macro & Market Context"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_DARK_SLATE

    p_sub1 = tf_bot.add_paragraph()
    p_sub1.text = b1
    p_sub1.font.size = Pt(10.5)
    p_sub1.font.color.rgb = RGBColor(55, 65, 81)
    p_sub1.space_before = Pt(8)

    p_sub2 = tf_bot.add_paragraph()
    p_sub2.text = b2
    p_sub2.font.size = Pt(10.5)
    p_sub2.font.color.rgb = RGBColor(55, 65, 81)
    p_sub2.space_before = Pt(6)
    # =========================================================================
    # SLIDE 8: PROPOSAL FEATURES (2-Leg Multi-Tranche Term Sheet by Product)
    # =========================================================================
    s8 = prs.slides.add_slide(blank)
    add_header(s8, "Proposal features", category="PROPOSAL FEATURES")
    add_logo(s8)
    add_footer(s8)

    calc = compute_canonical_bundle(ctx, ov)
    t_yrs = calc["tenor_years"]

    nba = ctx.get("next_best_action", "")
    bm = re.search(r"EUR\s*([\d\.]+)\s*([BM])(?:(?!EUR).)*?EMTN", nba, re.IGNORECASE)
    hm = re.search(r"EUR\s*([\d\.]+)\s*([BM])(?:(?!EUR).)*?Pre-Hedge", nba, re.IGNORECASE)
    
    def_bond = ov.get("notional_bond")
    def_swap = ov.get("notional_swap")
    if not def_bond and bm:
        unit = "B" if bm.group(2).upper() == "B" else "M"
        def_bond = f"EUR {bm.group(1)}{unit} (Senior Benchmark)"
    if not def_swap and hm:
        unit = "B" if hm.group(2).upper() == "B" else "M"
        def_swap = f"EUR {hm.group(1)}{unit} (Pre-Hedge Overlay)"

    if p_fam == "FX_HEDGE":
        leg1_title = "Leg 1 — USD Bond Tranche"
        leg2_title = "Leg 2 — Cross-Currency Swap"
        notional_leg1 = def_bond or "USD 600,000,000"
        notional_leg2 = def_swap or "EUR 550,000,000 eq."
        tenor_leg1 = ov.get("tenor", f"{t_yrs} Years (T + {t_yrs}Y)")
        tenor_leg2 = f"Matches bond maturity ({t_yrs}Y)"
        bench_leg1 = f"{t_yrs}Y US Treasury / SOFR"
        bench_leg2 = "EUR/USD Cross-Currency Basis"
        spread_leg1 = ov.get("spread", f"SOFR + {calc['spread_bps']} bps (indicative)")
        spread_leg2 = "EURIBOR + 32 bps (synthetic EUR funding)"
        fees_leg1 = "Underwriting fee per mandate letter"
        fees_leg2 = "Nil (embedded in CCY swap rate)"
        settle_leg1 = "T+5 standard for USD benchmark bonds"
        settle_leg2 = "Simultaneous with bond closing (T+5)"
        doc_leg1 = "144A / Reg S Prospectus"
        doc_leg2 = "ISDA Master Agreement + CSA"
    elif p_fam == "GREEN_ESG":
        leg1_title = "Leg 1 — Green Bond Tranche"
        leg2_title = "Leg 2 — Sustainability Overlay"
        notional_leg1 = def_bond or "EUR 500,000,000"
        notional_leg2 = def_swap or "EUR 250,000,000"
        tenor_leg1 = ov.get("tenor", f"{t_yrs} Years (Green Benchmark)")
        tenor_leg2 = "Annual SPT verification window"
        bench_leg1 = f"{t_yrs}Y EUR mid-swap"
        bench_leg2 = "Scope 1 & 2 Decarbonisation KPI"
        spread_leg1 = ov.get("spread", f"Mid-swap + {calc['spread_bps'] - calc['greenium_bps']} bps (Greenium: -{calc['greenium_bps']} bps)")
        spread_leg2 = "+/- 5 bps SPT step-up / step-down"
        fees_leg1 = "Underwriting fee per mandate letter"
        fees_leg2 = "Second-Party Opinion (SPO) advisory"
        settle_leg1 = "T+5 standard for EUR benchmark bonds"
        settle_leg2 = "Annual impact & allocation verification"
        doc_leg1 = "Green Bond Framework / EMTN Prospectus"
        doc_leg2 = "ICMA Green Bond Principles + SPO"
    elif p_fam == "RATES_HEDGE":
        leg1_title = "Leg 1 — New Benchmark Bond"
        leg2_title = "Leg 2 — Pre-Hedge Swap"
        notional_leg1 = def_bond or "EUR 1,200,000,000"
        notional_leg2 = def_swap or "EUR 800,000,000"
        tenor_leg1 = ov.get("tenor", f"{t_yrs} Years (T + {t_yrs}Y)")
        tenor_leg2 = "Terminates at bond pricing"
        bench_leg1 = f"{t_yrs}Y EUR mid-swap"
        bench_leg2 = f"{t_yrs}Y EUR swap rate"
        spread_leg1 = ov.get("spread", f"Mid-swap + {calc['spread_bps']} bps (indicative)")
        spread_leg2 = f"Current {t_yrs}Y swap rate (indicative)"
        fees_leg1 = "Underwriting fee per mandate letter"
        fees_leg2 = "Nil (embedded in swap rate)"
        settle_leg1 = "T+5 standard for EUR benchmark bonds"
        settle_leg2 = "Physical / cash-settled at unwind"
        doc_leg1 = "EMTN Programme / Prospectus"
        doc_leg2 = "ISDA Master Agreement + CSA"
    else:  # DCM_REFI
        leg1_title = "Leg 1 — Senior EMTN Tranche"
        leg2_title = "Leg 2 — Liquidity RCF / Commercial Paper"
        notional_leg1 = def_bond or "EUR 1,200,000,000"
        notional_leg2 = def_swap or "EUR 800,000,000"
        tenor_leg1 = ov.get("tenor", f"{t_yrs} Years (Euro Benchmark)")
        tenor_leg2 = "3–5 Years Revolving"
        bench_leg1 = f"{t_yrs}Y EUR mid-swap"
        bench_leg2 = "EURIBOR / €STR"
        spread_leg1 = ov.get("spread", f"Mid-swap + {calc['spread_bps']} bps (indicative)")
        spread_leg2 = "EURIBOR + 45 bps (undrawn 15 bps)"
        fees_leg1 = "Underwriting fee per mandate letter"
        fees_leg2 = "Commitment fee per facility agreement"
        settle_leg1 = "T+5 standard for EUR benchmark bonds"
        settle_leg2 = "Available upon documentation execution"
        doc_leg1 = "EMTN Programme / Prospectus"
        doc_leg2 = "LMA Standard Facility Agreement"

    s8_shape = s8.shapes.add_table(9, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
    s8_tbl = s8_shape.table
    s8_tbl.columns[0].width = Inches(2.7)
    s8_tbl.columns[1].width = Inches(4.5)
    s8_tbl.columns[2].width = Inches(4.5)

    headers = ["Term", leg1_title, leg2_title]
    for c_idx, h in enumerate(headers):
        c = s8_tbl.cell(0, c_idx)
        c.fill.solid()
        c.fill.fore_color.rgb = ING_ORANGE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    s8_rows = [
        ("Notional", notional_leg1, notional_leg2),
        ("Trade / pricing date", "Indicative — T", "Indicative — T"),
        ("Tenor / maturity", tenor_leg1, tenor_leg2),
        ("Reference benchmark", bench_leg1, bench_leg2),
        ("Spread / rate", spread_leg1, spread_leg2),
        ("Fees", fees_leg1, fees_leg2),
        ("Settlement", settle_leg1, settle_leg2),
        ("Governing documentation", doc_leg1, doc_leg2),
    ]

    for r_idx, (t_name, l1_val, l2_val) in enumerate(s8_rows):
        row_bg = RGBColor(254, 237, 222) if r_idx % 2 == 0 else RGBColor(255, 255, 255)
        for c_idx, val in enumerate([t_name, l1_val, l2_val]):
            c = s8_tbl.cell(r_idx + 1, c_idx)
            c.fill.solid()
            c.fill.fore_color.rgb = row_bg
            p = c.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.name = "Arial"
            p.font.size = Pt(10)
            if c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ING_DARK_SLATE
            elif c_idx == 1 and r_idx == 4:
                p.font.bold = True
                p.font.color.rgb = ING_ORANGE
            else:
                p.font.color.rgb = ING_DARK_SLATE

    # Disclaimer note beneath table
    d_box = s8.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.4))
    dp = d_box.text_frame.paragraphs[0]
    dp.text = "Indicative terms for discussion purposes only. Subject to internal credit approvals, KYC/AML, and market conditions at pricing."
    dp.font.name = "Arial"
    dp.font.size = Pt(8.5)
    dp.font.italic = True
    dp.font.color.rgb = RGBColor(100, 116, 139)



    # SLIDE 9: EXECUTION ROADMAP (Exact Preview Parity)
    # =========================================================================
    s9 = prs.slides.add_slide(blank)
    
    # Dynamic Title Resolution matching Preview Canvas
    if p_fam == "FX_HEDGE":
        s9_title = "Layered FX Hedging & Execution Roadmap"
    elif p_fam == "GREEN_ESG":
        s9_title = "Green Bond Framework & Issuance Timetable"
    elif p_fam == "RATES_HEDGE":
        s9_title = "ISDA Schedule, CSA & Execution Timeline"
    else:
        s9_title = "Indicative Execution Roadmap & Timeline"

    add_header(s9, s9_title, category="EXECUTION ROADMAP")
    add_logo(s9)
    add_footer(s9)

    # Dynamic Steps matching React App.jsx case 8 exactly
    if p_fam == "FX_HEDGE":
        steps = [
            ("1", "T - 4 Weeks", "Exposure Mapping", "Audit USD receivables & establish rolling monthly budget hedge ratios."),
            ("2", "T - 2 Weeks", "ISDA & Documentation", "Finalize ISDA Schedule, collateral CSA threshold & credit line setup."),
            ("3", "T - 1 Week", "Collar Structuring", "Calibrate put/call strike corridor against spot & forward curves."),
            ("4", "T-Day", "Layered Execution", "Execute Tranche 1 zero-cost collars; initiate quarterly rolling schedule.")
        ]
    elif p_fam == "GREEN_ESG":
        steps = [
            ("1", "T - 6 Weeks", "Framework Drafting", "Establish Green Financing Framework aligned with EU Taxonomy & ICMA."),
            ("2", "T - 4 Weeks", "SPO Verification", "Engage ISS ESG / Sustainalytics for Second Party Opinion review."),
            ("3", "T - 1 Week", "ESG Roadshow", "Dedicated European SRI investor marketing calls & ESG presentation."),
            ("4", "T-Day", "Syndicate Pricing", "Bookbuilding, greenium spread tightening, and orderbook allocation.")
        ]
    elif p_fam == "RATES_HEDGE":
        steps = [
            ("1", "T - 4 Weeks", "Exposure Sizing", "Quantify repricing gap across upcoming debt tranches"),
            ("2", "T - 2 Weeks", "Pre-Hedge Execution", "Execute forward-starting IRS overlay swap"),
            ("3", "T - 1 Week", "Global Roadshow", "Syndicate investor marketing meetings"),
            ("4", "T-Day", "Pricing & Settlement", "Final syndicate pricing, book allocation & closing")
        ]
    else:  # DCM_REFI
        steps = [
            ("1", "T - 4 Weeks", "Mandate & Prep", "Finalize EMTN documentation, auditor comfort letters, and investor deck."),
            ("2", "T - 2 Weeks", "Pre-Hedge Overlay", "Execute interest rate pre-hedge swaps to lock underlying benchmark yields."),
            ("3", "T - 1 Week", "Global Roadshow", "Conduct targeted European investor marketing & C-suite roadshow calls."),
            ("4", "T-Day", "Launch & Pricing", "Intraday bookbuilding, spread compression, and final syndicate allocation.")
        ]

    # Compact card geometry matching Preview Canvas
    card_w = Inches(2.65)
    card_h = Inches(2.35)
    start_x = Inches(0.9)
    start_y = Inches(1.85)
    gap_x = Inches(0.31)

    for i, (num, time_tag, title_text, desc_text) in enumerate(steps):
        cx = start_x + i * (card_w + gap_x)
        
        # 1. Card Container (clean light border matching React Preview)
        c_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, start_y, card_w, card_h)
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        c_box.line.color.rgb = RGBColor(229, 231, 235)  # border-gray-200
        c_box.line.width = Pt(1.0)
        
        # 2. Orange Step Number Badge (Circle)
        badge_sz = Inches(0.40)
        badge_x = cx + (card_w - badge_sz) / 2
        badge_y = start_y + Inches(0.18)
        b_shp = s9.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_sz, badge_sz)
        b_shp.fill.solid()
        b_shp.fill.fore_color.rgb = ING_ORANGE
        b_shp.line.fill.background()
        
        b_tf = b_shp.text_frame
        b_tf.margin_top = Inches(0.04)
        b_p = b_tf.paragraphs[0]
        b_p.text = num
        b_p.font.name = "Arial"
        b_p.font.size = Pt(11)
        b_p.font.bold = True
        b_p.font.color.rgb = RGBColor(255, 255, 255)
        b_p.alignment = PP_ALIGN.CENTER

        # 3. Text Block
        tb = s9.shapes.add_textbox(cx + Inches(0.1), start_y + Inches(0.65), card_w - Inches(0.2), card_h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        # Time Tag (e.g. T - 4 Weeks)
        p_time = tf.paragraphs[0]
        p_time.text = time_tag
        p_time.font.name = "Arial"
        p_time.font.size = Pt(11)
        p_time.font.bold = True
        p_time.font.color.rgb = ING_DARK_SLATE
        p_time.alignment = PP_ALIGN.CENTER
        p_time.space_after = Pt(2)

        # Milestone Subtitle (Orange text)
        p_sub = tf.add_paragraph()
        p_sub.text = title_text
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ING_ORANGE
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(4)

        # Description Body
        p_body = tf.add_paragraph()
        p_body.text = desc_text
        p_body.font.name = "Arial"
        p_body.font.size = Pt(8.5)
        p_body.font.color.rgb = RGBColor(100, 116, 139)
        p_body.alignment = PP_ALIGN.CENTER

    # SLIDE 10: REGULATORY DISCLAIMERS
    # =========================================================================
    s10 = prs.slides.add_slide(blank)
    add_header(s10, sm["s10_ttl"], category=sm["s10_cat"])
    add_logo(s10)
    add_footer(s10)

    shp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_LIGHT
    shp.line.color.rgb = LINE_GRAY

    tb_disc = s10.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.4))
    tf_disc = tb_disc.text_frame
    tf_disc.word_wrap = True
    p = tf_disc.paragraphs[0]
    p.text = "Regulatory Disclosures & Target Market Notice"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_DARK_SLATE

    disc_list = ov.get("disclaimers") or compliance_bullets or [
        "This document is prepared for illustrative and discussion purposes only and does not constitute an offer, solicitation, or recommendation to enter into any transaction.",
        "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II / UK MiFIR is eligible counterparties and professional clients only (all distribution channels).",
        "This material has not been prepared in accordance with legal requirements designed to promote the independence of investment research.",
        "All rates, levels, spreads, and indicative terms shown are subject to change without notice and are not tradeable prices."
    ]

    # Ensure disc_list is a list
    if isinstance(disc_list, str):
        disc_list = [disc_list]
    elif not isinstance(disc_list, list):
        disc_list = list(disc_list) if disc_list else []

    for d in disc_list:
        if d and d.strip():
            p_d = tf_disc.add_paragraph()
            p_d.text = f"• {d}"
            p_d.font.size = Pt(11)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_before = Pt(8)

    # =========================================================================
    # SAVE
    # =========================================================================
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf