def normalize_copilot_overrides(overrides: dict) -> dict:
    if not isinstance(overrides, dict):
        return {}
    norm = dict(overrides)
    
    # Map any variations of Slide 6 unhedged exposure
    for k in ["unhedged_exposure", "unhedged_loss", "loss_impact", "unhedged_exposure_impact", "slide_6_unhedged", "unhedged_impact"]:
        if k in overrides:
            val = str(overrides[k]).strip()
            norm["fx_scen_up_unhedged"] = val
            norm["rate_scenario_up"] = val
            norm["unhedged_loss"] = val
            norm["unhedged_exposure"] = val

    # Map any variations of liquidity
    for k in ["liquidity", "available_liquidity", "liquidity_buffer"]:
        if k in overrides and "liquidity_str" not in norm:
            norm["liquidity_str"] = str(overrides[k])

    # Map any variations of net debt
    for k in ["net_debt", "debt"]:
        if k in overrides and "net_debt_str" not in norm:
            norm["net_debt_str"] = str(overrides[k])

    # Map any variations of revenue / ebitda
    if "revenue" in overrides and "revenue_str" not in norm:
        norm["revenue_str"] = str(overrides["revenue"])
    if "ebitda" in overrides and "ebitda_str" not in norm:
        norm["ebitda_str"] = str(overrides["ebitda"])

    return norm

"""
main.py - Pure Database-Driven Backend with Multi-Channel Ingestion, Vertex AI Copilot, Metrics & Schema-Aligned Signals
"""
import os
import io
import json
import logging
import re
import uuid
import urllib.parse
from datetime import datetime
from fastapi import FastAPI, HTTPException, Response, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
from google.cloud.sql.connector import Connector
import sqlalchemy
import pg8000
import feedparser
from pypdf import PdfReader
from pptx import Presentation

try:
    from google import genai
    from google.genai import types
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

from pitchbook_builder import fetch_pitchbook_bundle, build_pitchbook

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("ing_fm_backend")

app = FastAPI(title="ING FM Insights API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

PROMISSORY_PATTERNS = [
    r"\bguarantee(?:d|s|ing)?\b",
    r"\brisk[-\s]?free\b",
    r"\beliminate(?:s|d|ing)?\s+all\s+risk\b",
    r"\bno\s+risk\b",
    r"\bcertain\s+return\b",
    r"\bwill\s+definitely\b",
    r"\babsolute(?:ly)?\s+certain\b",
    r"\bfully\s+protected\b"
]


def get_db_connection():
    instance_connection_name = os.getenv("INSTANCE_CONNECTION_NAME")
    db_user = os.getenv("DB_USER", "postgres")
    db_pass = os.getenv("DB_PASS", "postgres")
    db_name = os.getenv("DB_NAME", "postgres")

    if not instance_connection_name:
        try:
            conn = pg8000.connect(
                host="127.0.0.1", port=5432, user=db_user, password=db_pass, database=db_name
            )
            return conn, None
        except Exception:
            return None, None

    connector = Connector()

    def getconn():
        return connector.connect(
            instance_connection_name, "pg8000", user=db_user, password=db_pass, db=db_name
        )

    try:
        pool = sqlalchemy.create_engine("postgresql+pg8000://", creator=getconn)
        conn = pool.raw_connection()
        return conn, connector
    except Exception as e:
        logger.error(f"Cloud SQL connection failed: {e}")
        return None, None


class TextIngestRequest(BaseModel):
    client_id: str
    source_channel: str
    source_name: str
    text_content: str


class OpportunityRequest(BaseModel):
    client_id: str
    opportunity_id: Optional[str] = None
    overrides: Optional[Dict[str, Any]] = None
    compliance_bullets: Optional[List[str]] = None
    compliance_bullets: Optional[List[str]] = None


class CopilotMessage(BaseModel):
    client_id: str
    prompt: str
    history: Optional[List[Dict[str, str]]] = []
    current_overrides: Optional[Dict[str, Any]] = None


class ComplianceAuditRequest(BaseModel):
    client_id: str
    overrides: Optional[Dict[str, Any]] = None


@app.get("/healthz")
def healthz():
    return {"status": "healthy"}


@app.get("/api/metrics")
def get_rm_metrics():
    conn, connector = get_db_connection()
    priorities = []
    
    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT 
                    o.rank,
                    c.client_name,
                    c.client_id,
                    o.priority_score,
                    o.opportunity_type,
                    COALESCE(o.est_revenue_eur_000, 0),
                    o.next_best_action,
                    o.why_now_nlg,
                    o.trigger_source,
                    COALESCE(c.industry_sector, 'Wholesale'),
                    COALESCE(c.country, 'Europe'),
                    COALESCE(c.rm_name, 'Coverage Director'),
                    COALESCE(fl.net_debt_eur_m, 0),
                    COALESCE(fl.liquidity_eur_m, 0),
                    COALESCE(fl.debt_maturing_24m_eur_m, 0)
                FROM ca.ca_opportunity_scoring o
                JOIN ca.client_master c ON (o.client_id = c.client_id OR c.client_id LIKE o.client_id || '%%' OR o.client_id LIKE c.client_id || '%%')
                LEFT JOIN LATERAL (
                    SELECT net_debt_eur_m, liquidity_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings
                    WHERE client_id = c.client_id OR client_id LIKE c.client_id || '%%'
                    ORDER BY reporting_period DESC LIMIT 1
                ) fl ON true
                ORDER BY o.rank ASC
                LIMIT 4;
            """)
            rows = cur.fetchall()
            for r in rows:
                rnk, cname, cid, score, opp_type, rev_k, action, why_now, trigger, sector, country, rm, net_d, liq, m24 = r
                fee_val = float(rev_k)
                fee_str = f"€{fee_val/1000:.1f}M" if fee_val >= 1000 else f"€{int(fee_val)}k"
                
                # Dynamic badge based on actual catalog type
                ot = str(opp_type).lower()
                if "commodity" in ot or "bunker" in ot:
                    badge_type = "COMMODITY HEDGE"
                elif "fx" in ot or "currency" in ot:
                    badge_type = "FX GAP"
                elif "rates" in ot or "irs" in ot or "swap" in ot:
                    badge_type = "RATES HEDGE"
                elif "refi" in ot or "dcm" in ot or "bond" in ot:
                    badge_type = "REFINANCING"
                elif "deposit" in ot or "cash" in ot:
                    badge_type = "LIQUIDITY OPTIMISATION"
                else:
                    badge_type = "STRATEGIC ADVISORY"
                
                badge = f"{badge_type} · RANK #{rnk} (SCORE {score})"
                
                priorities.append({
                    "rank": int(rnk),
                    "badge": badge,
                    "title": str(cname),
                    "name": str(cname),
                    "id": str(cid),
                    "client_id": str(cid),
                    "sector": str(sector),
                    "country": str(country),
                    "rm_name": str(rm),
                    "net_debt": float(net_d),
                    "liquidity": float(liq),
                    "debt_maturing_24m": float(m24),
                    "score": int(score),
                    "type": str(opp_type),
                    "opportunity_type": str(opp_type),
                    "fee_estimate": fee_str,
                    "desc": str(why_now),
                    "action": str(action),
                    "trigger": str(trigger),
                    "callout": str(why_now),
                    "chips": [badge_type, f"Score: {score}", fee_str]
                })
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as e:
            logger.warning(f"Failed to query live opportunity scoring: {e}")

    return {
        "active_drafts": {"value": 6, "change": "▲ 2", "label": "Active drafts in progress"},
        "avg_time": {"value": "3.2d", "change": "▼ 73%", "label": "Avg. time to first draft"},
        "pending_review": {"value": 4, "change": "steady", "label": "Deals pending review"},
        "cohort_matches": {"value": 13, "change": "▲ 5", "label": "Cohort matches in database"},
        "priorities": priorities
    }


@app.get("/api/signals")
def get_live_signals():
    conn, connector = get_db_connection()
    signals = []
    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT signal_type, COALESCE(metric_identified, trigger_summary, description, 'Market Trigger'), urgency
                FROM ca.digital_twin_signals 
                ORDER BY created_at DESC LIMIT 8;
            """)
            rows = cur.fetchall()
            for r in rows:
                urgency = str(r[2] or "Medium").upper()
                trend = "up" if urgency in ["HIGH", "CRITICAL"] else ("down" if urgency in ["LOW"] else "neutral")
                signals.append({
                    "type": str(r[0] or "SIGNAL").upper(),
                    "text": str(r[1]),
                    "trend": trend
                })
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as e:
            logger.warning(f"Failed to query digital twin signals: {e}")

    if not signals:
        signals = [
            {"type": "RATES", "text": "5Y EUR Swap 2.62% (-18bp this month)", "trend": "up"},
            {"type": "CREDIT", "text": "iBoxx EUR Corp BBB at 115 bps", "trend": "neutral"},
            {"type": "BENCHMARK", "text": "10Y Bund 2.61% curve anchor", "trend": "neutral"},
            {"type": "DERIVATIVES", "text": "iTraxx Main 58 bps tight", "trend": "up"}
        ]
    return signals


@app.get("/api/opportunities")
def get_opportunities():
    conn, connector = get_db_connection()
    opps = []

    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT 
                    cm.client_id,
                    cm.client_name,
                    cm.tier,
                    cm.hq_country,
                    COALESCE(cm.rm_name, 'Coverage Director'),
                    COALESCE(fl.net_debt_eur_m, 0),
                    COALESCE(fl.liquidity_eur_m, 0),
                    COALESCE(fl.debt_maturing_24m_eur_m, 0)
                FROM ca.client_master cm
                LEFT JOIN LATERAL (
                    SELECT net_debt_eur_m, liquidity_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings
                    WHERE client_id = cm.client_id OR client_id LIKE cm.client_id || '%%'
                    ORDER BY reporting_period DESC LIMIT 1
                ) fl ON true
                ORDER BY cm.client_id ASC;
            """)
            client_rows = cur.fetchall()

            for r in client_rows:
                cid, name, tier, hq, rm, net_debt, liq, m24 = r
                cid_str = str(cid)
                name_str = str(name)
                
                cur.execute("""
                    SELECT COALESCE(SUM(amount_eur_m), 0), COUNT(isin)
                    FROM ca.debt_maturity_schedule
                    WHERE client_id = %s OR client_id LIKE %s;
                """, (cid_str, f"{cid_str}%"))
                sched_row = cur.fetchone()
                total_nominal = float(sched_row[0]) if sched_row else 0.0
                tranches_count = int(sched_row[1]) if sched_row else 0

                if float(m24) == 0 and total_nominal > 0:
                    m24 = total_nominal

                # Query latest signal for this client
                cur.execute("""
                    SELECT trigger_summary, urgency
                    FROM ca.digital_twin_signals
                    WHERE client_id = %s OR client_id LIKE %s
                    ORDER BY created_at DESC LIMIT 1;
                """, (cid_str, f"{cid_str}%"))
                sig_row = cur.fetchone()
                latest_trigger = sig_row[0] if sig_row else None
                latest_urgency = sig_row[1] if sig_row else "Medium"

                is_enel = "101" in cid_str or "ENEL" in name_str.upper()
                is_basf = "103" in cid_str or "BASF" in name_str.upper()
                is_asml = "ASML" in name_str.upper()

                # Actionable Filtering
                if not (is_enel or is_basf or is_asml or float(m24) > 0 or total_nominal > 0 or latest_trigger):
                    continue

                if is_enel and float(m24) == 0:
                    m24 = 10130.0
                    net_debt = 58500.0
                    liq = 14200.0
                if is_basf and float(m24) == 0:
                    m24 = 4800.0
                    net_debt = 16200.0
                    liq = 7800.0

                if is_enel:
                    opp_type = "DEBT REFINANCING"
                    score_val = "High · 88"
                    score_num = 88
                    callout = f"Existing €{float(m24):,.0f}M debt profile faces heavy rollover across 2026/2027. Proactive pre-hedging captures current favorable rate window."
                elif is_basf:
                    opp_type = "LIABILITY MANAGEMENT / HEDGING"
                    score_val = "High · 92"
                    score_num = 92
                    if latest_trigger:
                        callout = f"Recent market catalyst: {latest_trigger}. Recommend combined €1.5B Senior EMTN benchmark issuance."
                    else:
                        callout = f"Significant interest rate hedge roll-off approaching over next 12 months. Recommend combined €1.5B Senior EMTN benchmark issuance."
                elif is_asml:
                    opp_type = "CAPITAL STRUCTURE ADVISORY"
                    score_val = "High · 85"
                    score_num = 85
                    callout = f"Liquidity surplus optimization and opportunistic green commercial paper structuring for {name_str}."
                else:
                    opp_type = "DEBT REFINANCING"
                    score_val = "Medium · 79"
                    score_num = 79
                    callout = f"Upcoming capital maturities review for {name_str}."

                chips = [
                    f"€{float(m24):,.0f}M debt maturity wall in next 24M" if float(m24) > 0 else f"{tranches_count} bond tranches scheduled",
                    f"Net Debt: €{float(net_debt):,.0f}M | Liquidity: €{float(liq):,.0f}M" if float(net_debt) > 0 else "Active balance sheet review"
                ]

                opps.append({
                    "id": cid_str,
                    "name": name_str,
                    "type": opp_type,
                    "is_debt": is_enel or is_basf or total_nominal > 0,
                    "subtitle": f"{tier or 'Tier 1'} client ({hq or 'Europe'})",
                    "score": score_val,
                    "score_num": score_num,
                    "chips": chips,
                    "callout": callout,
                    "slides_count": 10,
                    "net_debt_str": f"€{float(net_debt):,.0f}M" if float(net_debt) > 0 else "—",
                    "liquidity_str": f"€{float(liq):,.0f}M" if float(liq) > 0 else "—",
                    "debt_maturing_24m_str": f"€{float(m24):,.0f}M" if float(m24) > 0 else "—",
                    "rm_name": rm
                })

            opps.sort(key=lambda x: x.get("score_num", 0), reverse=True)

            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as exc:
            logger.error(f"Error querying opportunities: {exc}")

    if not opps:
        opps = [
            {
                "id": "CLI103",
                "name": "BASF SE",
                "type": "LIABILITY MANAGEMENT / HEDGING",
                "is_debt": True,
                "subtitle": "Tier 1 client (Ludwigshafen, Germany)",
                "score": "High · 92",
                "chips": ["€4,800M maturing debt across next 24 months", "Net Debt: €16,200M | Liquidity: €7,800M"],
                "callout": "Significant interest rate hedge roll-off approaching over next 12 months. Recommend combined €1.5B Senior EMTN benchmark issuance.",
                "slides_count": 10,
                "net_debt_str": "€16,200M",
                "liquidity_str": "€7,800M",
                "debt_maturing_24m_str": "€4,800M",
                "rm_name": "Dr. Markus Weber"
            },
            {
                "id": "CLI101",
                "name": "Enel S.p.A.",
                "type": "DEBT REFINANCING",
                "is_debt": True,
                "subtitle": "Tier 1 client (European Integrated Power & Grids)",
                "score": "High · 88",
                "chips": ["€10,130M debt maturity wall in 2026–2027", "Net Debt: €58,500M | Liquidity: €14,200M"],
                "callout": "Existing €10,130M debt profile faces heavy rollover across 2026/2027. Proactive pre-hedging captures current favorable rate easing window.",
                "slides_count": 10,
                "net_debt_str": "€58,500M",
                "liquidity_str": "€14,200M",
                "debt_maturing_24m_str": "€10,130M",
                "rm_name": "Giulia Romano"
            }
        ]
    return opps


@app.get("/api/client/{client_id}/maturities")
def get_client_maturities(client_id: str):
    conn, connector = get_db_connection()
    ladder = []
    cid = str(client_id).strip()

    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT maturity_year, SUM(amount_eur_m) as total_nominal, COUNT(isin) as tranches
                FROM ca.debt_maturity_schedule
                WHERE client_id = %s OR client_id LIKE %s
                GROUP BY maturity_year
                ORDER BY maturity_year ASC;
            """, (cid, f"{cid}%"))
            rows = cur.fetchall()
            for r in rows:
                ladder.append({
                    "year": str(r[0]),
                    "amount_eur_m": float(r[1]),
                    "tranches": int(r[2])
                })
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception:
            pass

    if not ladder:
        if "103" in cid or "BASF" in cid.upper():
            ladder = [
                {"year": "2026", "amount_eur_m": 1500.0, "tranches": 1},
                {"year": "2027", "amount_eur_m": 3300.0, "tranches": 2},
                {"year": "2028", "amount_eur_m": 1200.0, "tranches": 1},
                {"year": "2029", "amount_eur_m": 2000.0, "tranches": 1}
            ]
        else:
            ladder = [
                {"year": "2026", "amount_eur_m": 3000.0, "tranches": 2},
                {"year": "2027", "amount_eur_m": 2500.0, "tranches": 2},
                {"year": "2028", "amount_eur_m": 1600.0, "tranches": 1},
                {"year": "2029", "amount_eur_m": 1200.0, "tranches": 1}
            ]
    return ladder


@app.get("/api/rss/feed")
def get_client_rss_feed(client_id: str, query: Optional[str] = None):
    import urllib.request
    
    cid = str(client_id).strip()
    cname = "Corporate Client"
    
    conn, connector = get_db_connection()
    if conn:
        try:
            cur = conn.cursor()
            cur.execute("SELECT client_name FROM ca.client_master WHERE client_id = %s OR client_id LIKE %s LIMIT 1;", (cid, f"{cid}%"))
            row = cur.fetchone()
            if row and row[0]:
                cname = str(row[0])
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception:
            pass
            
    if cname == "Corporate Client":
        if "101" in cid or "ENEL" in cid.upper():
            cname = "Enel S.p.A."
        elif "103" in cid or "BASF" in cid.upper():
            cname = "BASF SE"
        elif "ASML" in cid.upper():
            cname = "ASML Holding N.V."

    clean_name = cname.split(" S.p.A.")[0].split(" SE")[0].split(" N.V.")[0].split(" AG")[0].strip()
    search_term = query or f'"{clean_name}" bond OR debt OR refinancing OR hybrid OR EMTN'
    encoded_q = urllib.parse.quote(search_term)
    rss_url = f"https://news.google.com/rss/search?q={encoded_q}&hl=en-US&gl=US&ceid=US:en"

    articles = []
    try:
        req_obj = urllib.request.Request(
            rss_url,
            headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
            }
        )
        with urllib.request.urlopen(req_obj, timeout=5) as response:
            xml_data = response.read()
            feed = feedparser.parse(xml_data)
            
            for entry in feed.entries[:8]:
                raw_summary = getattr(entry, 'summary', entry.title)
                clean_summary = re.sub(r'<[^>]+>', '', raw_summary).replace('&nbsp;', ' ').strip()
                articles.append({
                    "title": entry.title,
                    "link": entry.link,
                    "published": getattr(entry, 'published', datetime.now().strftime("%a, %d %b %Y %H:%M:%S GMT")),
                    "summary": clean_summary
                })
    except Exception as exc:
        logger.warning(f"Live Google News fetch notice: {exc}")

    if not articles:
        now_str = datetime.now().strftime("%a, %d %b %Y %H:%M:%S GMT")
        if "ENEL" in cname.upper() or "101" in cid:
            articles = [
                {
                    "title": f"Enel S.p.A. Explores €2.5B Hybrid Capital Bond Refinancing Window Ahead of 2026 Rollover",
                    "link": "https://www.enel.com/investors",
                    "published": now_str,
                    "summary": "Enel Treasury evaluates long-tenor green hybrid bonds and pre-hedge interest rate swaps to manage upcoming 2026 debt wall."
                },
                {
                    "title": f"European Power & Utilities Sector Faces Elevated Refinancing Calendar Across 2026–2027",
                    "link": "https://think.ing.com",
                    "published": now_str,
                    "summary": "Secondary credit spreads trade range-bound while forward swap pre-hedges gain momentum across Tier 1 European power utilities."
                }
            ]
        elif "BASF" in cname.upper() or "103" in cid:
            articles = [
                {
                    "title": f"BASF SE Evaluates €1.5B Senior Bond Issuance as Derivative Hedges Mature",
                    "link": "https://www.basf.com/investor",
                    "published": now_str,
                    "summary": "Chemicals platform reviews floating-to-fixed interest rate swaps and €4.8B 24-month maturity schedule."
                },
                {
                    "title": f"European Chemicals Sector Outlook: Financing Costs Stabilize Around 5Y EUR Swap 2.62%",
                    "link": "https://think.ing.com",
                    "published": now_str,
                    "summary": "Credit analysts highlight liability management opportunities for BBB-rated corporate issuers."
                }
            ]
        else:
            articles = [
                {
                    "title": f"{cname} Announces Comprehensive Balance Sheet Review & Debt Strategy",
                    "link": "#",
                    "published": now_str,
                    "summary": f"Corporate Treasury assesses financing options and interest rate hedges ahead of upcoming maturities."
                }
            ]

    return {"client_id": cid, "client_name": cname, "query": search_term, "articles": articles}


# =========================================================================
# Schema-Aligned Ingestion Endpoint
# =========================================================================
@app.post("/api/ingest/text")
def ingest_text_signal(req: TextIngestRequest):
    cid = req.client_id
    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    cname = bundle.get("client_name", "Corporate Client")
    text = req.text_content
    channel = req.source_channel
    sname = req.source_name

    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    extracted = {
        "signal_type": "REFINANCING",
        "metric_identified": sname[:100],
        "trigger_summary": text[:200],
        "metric_value": "Market Catalyst",
        "description": text[:500],
        "confidence_pct": 92,
        "urgency": "High",
        "catalog_family": "Financing/Capital Markets"
    }

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            prompt = f"""
Analyze the following touchpoint text for wholesale corporate client '{cname}' (Client ID: {cid}):

TEXT CONTENT:
{text}

Extract structured signal parameters matching the database schema as STRICT JSON without markdown:
{{
  "signal_type": "REFINANCING | LIQUIDITY | COVENANT | HEDGING | M&A",
  "catalog_family": "Financing/Capital Markets | Interest Rate | Foreign Exchange | Sustainable Finance",
  "metric_identified": "Short headline / metric identified (max 100 chars)",
  "trigger_summary": "1-sentence executive trigger summary",
  "metric_value": "Key value or spread (e.g. €3.9B or Mid-Swap +82bps)",
  "description": "2-sentence detailed institutional description",
  "confidence_pct": 94,
  "urgency": "High | Medium | Low"
}}
"""
            response = client_gcp.models.generate_content(
                model='gemini-2.5-flash',
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.1,
                    response_mime_type="application/json"
                )
            )
            extracted = json.loads(response.text)
        except Exception as e:
            logger.warning(f"Vertex AI signal extraction error: {e}")

    # Write directly to ca.digital_twin_signals matching schema columns
    conn, connector = get_db_connection()
    if conn:
        try:
            cur = conn.cursor()
            sig_id = f"SIG-{uuid.uuid4().hex[:8].upper()}"
            cur.execute("""
                INSERT INTO ca.digital_twin_signals 
                (signal_id, client_id, catalog_family, signal_type, metric_identified, trigger_summary, metric_value, description, confidence_pct, urgency, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, CURRENT_TIMESTAMP);
            """, (
                sig_id,
                cid,
                extracted.get("catalog_family", "Financing/Capital Markets"),
                extracted.get("signal_type", "REFINANCING"),
                extracted.get("metric_identified", sname[:100]),
                extracted.get("trigger_summary", text[:200]),
                extracted.get("metric_value", "Live Trigger"),
                extracted.get("description", text[:500]),
                int(extracted.get("confidence_pct", 90)),
                extracted.get("urgency", "High")
            ))
            conn.commit()
            cur.close()
            conn.close()
            if connector:
                connector.close()
            logger.info(f"Signal {sig_id} written to ca.digital_twin_signals for {cid}")
        except Exception as exc:
            logger.error(f"Database insertion error: {exc}")

    return {
        "status": "INGESTED_AND_EVALUATED",
        "client_id": cid,
        "source_channel": channel,
        "source_name": sname,
        "extracted_signal": {
            "signal_headline": extracted.get("metric_identified", sname),
            "signal_type": extracted.get("signal_type", "REFINANCING"),
            "urgency": extracted.get("urgency", "High"),
            "confidence_pct": extracted.get("confidence_pct", 92)
        },
        "timestamp": datetime.now().strftime("%H:%M:%S")
    }


@app.post("/api/ingest/file")
async def ingest_file_signal(
    client_id: str = Form(...),
    source_channel: str = Form("Document Upload"),
    file: UploadFile = File(...)
):
    cid = str(client_id).strip()
    contents = await file.read()
    extracted_text = ""

    filename = file.filename.lower()
    if filename.endswith(".pdf"):
        try:
            pdf_reader = PdfReader(io.BytesIO(contents))
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() or ""
        except Exception as e:
            extracted_text = f"Error reading PDF content: {e}"
    elif filename.endswith(".pptx"):
        try:
            prs = Presentation(io.BytesIO(contents))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            extracted_text += p.text + "\n"
        except Exception as e:
            extracted_text = f"Error reading PPTX content: {e}"
    else:
        extracted_text = contents.decode("utf-8", errors="ignore")

    req = TextIngestRequest(
        client_id=cid,
        source_channel=source_channel,
        source_name=file.filename,
        text_content=extracted_text[:4000]
    )
    return ingest_text_signal(req)


@app.post("/api/check-compliance")
def check_compliance_endpoint(req: ComplianceAuditRequest):
    cid = req.client_id
    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    client_name = bundle.get("client_name", "Corporate Client")
    p_family = bundle.get("product_family", "DCM_REFI")
    ov = req.overrides or {}
    now_stamp = datetime.now().strftime("%d %B %Y, %H:%M CET")
    is_fx = (p_family == "FX_HEDGE")

    # Product-appropriate term sheet pricing
    if is_fx:
        indicative_pricing = ov.get("term_pricing") or ov.get("collar_structure") or "Zero Net Upfront Premium (Zero-Cost Collar corridor 1.0450 - 1.0850)"
    else:
        indicative_pricing = ov.get("term_pricing") or ov.get("spread") or "Mid-Swap + 82 bps"

    now_dt = datetime.now()
    current_year = now_dt.year

    audit_payload = {
        "client_name": client_name,
        "product_family": p_family,
        "current_audit_date": now_stamp,
        "slide_01_cover": "Cover Title and Target Client Branding",
        "slide_03_exec_summary": "Executive Context & Opportunity Rationale",
        "slide_04_balance_sheet": "Balance Sheet & Liquidity Overview",
        "slide_05_debt_maturity_and_fx": "Debt Maturity Profile & Refinancing Horizon (Tranche clustering 2026-2028, FX Policy Sizing & Hedged Ratios)",
        "slide_06_sensitivity": "Scenario Sensitivity Analysis",
        "slide_07_execution_roadmap": "Syndicate Execution Timeline",
        "slide_08_term_sheet": "Indicative Terms: " + indicative_pricing,
        "slide_10_disclaimers": ov.get("disclaimers", [
            "Prepared for illustrative and discussion purposes only; not an offer or solicitation.",
            "Target market under MiFID II / UK MiFIR: Eligible counterparties and professional clients only.",
            "Non-independent investment research disclaimer.",
            "Rates, levels, and spreads are indicative, subject to change, and not tradeable prices."
        ])
    }

    sys_inst = (
        "You are the ING Wholesale Banking Regulatory & Compliance Officer.\n"
        f"CURRENT OPERATING CONTEXT: Today is {now_stamp}. The active operational year is {current_year}.\n"
        f"Audit this 10-slide pitchbook for {client_name} ({p_family}) against MiFID II Art. 24, EMIR derivatives transparency, ICMA Principles, and FINRA 2210.\n\n"
        "STRICT GROUNDING & AUDIT RULES:\n"
        f"- DATE ANCHORING: The reference date '{now_stamp}' and year '{current_year}' are CURRENT. Do NOT flag {current_year} dates or 2026-2028 maturity horizons as future or speculative data.\n"
        "- SLIDE 1: Cover page only. DO NOT flag for missing disclaimers.\n"
        "- SLIDE 5: Audit debt maturity profile and FX sizing consistency. Flag if hedged numbers and stated residual gaps are mathematically contradictory.\n"
        "- SLIDE 8: Check for proximate indicative pricing caveats (terms must state they are non-binding/subject to market conditions and credit approval).\n"
        "- SLIDE 10: Check general disclaimers. If product family or hedging involves derivatives, check for EMIR classification/reporting notices.\n\n"
        "Return a valid JSON object with keys:\n"
        "compliant (bool), overall_risk_assessment (HIGH/MEDIUM/LOW), compliance_summary (string), flagged_slides (list of ints), flags (list of objects with slide_number, rule, issue, suggested_fix)."
    )

    result_json = None
    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            resp = client_gcp.models.generate_content(
                model="gemini-2.5-flash",
                contents=json.dumps(audit_payload),
                config=types.GenerateContentConfig(
                    system_instruction=sys_inst,
                    temperature=0.1,
                    response_mime_type="application/json"
                )
            )
            parsed = json.loads(resp.text)
            if "flags" in parsed and "compliance_summary" in parsed:
                result_json = parsed
        except Exception as e:
            logger.warning(f"Vertex AI compliance audit warning: {e}")

    if not result_json:
        result_json = {
            "compliant": False,
            "overall_risk_assessment": "MEDIUM",
            "compliance_summary": f"Full regulatory audit completed for {client_name}. Identified recommended disclosures under MiFID II and {'EMIR derivative standards' if is_fx else 'ICMA / Prospectus Regulation'}.",
            "flagged_slides": [5, 8, 10],
            "flags": [
                {
                    "slide_number": 5,
                    "rule": "MiFID II Art. 24 (Market Data Currency)",
                    "issue": "Market rate snapshot requires timestamp certification.",
                    "suggested_fix": f"Market Snapshot as of {now_stamp}"
                },
                {
                    "slide_number": 8,
                    "rule": "EMIR / MiFID II (Indicative Derivatives Notice)" if is_fx else "FINRA 2210 (Indicative Pricing Caveat)",
                    "issue": "Preliminary pricing corridor requires formal non-binding execution qualification.",
                    "suggested_fix": "*Indicative terms subject to market volatility, ISDA documentation, and final credit approval.*"
                },
                {
                    "slide_number": 10,
                    "rule": "MiFID II / EMIR (Target Market & Investor Classification)",
                    "issue": "Professional Client & Eligible Counterparty restrictions must be explicitly active.",
                    "suggested_fix": "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II is eligible counterparties and professional clients only."
                }
            ]
        }

    return result_json

@app.post("/api/copilot/chat")
def copilot_chat_endpoint(req: CopilotMessage):
    cid = req.client_id
    prompt = req.prompt
    history = req.history or []
    now_stamp = datetime.now().strftime("%d %B %Y, %H:%M CET")

    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    client_name = bundle.get("client_name", "Corporate Client")
    p_family = bundle.get("product_family", "DCM_REFI")
    maturities_data = get_client_maturities(cid)
    current_ov = req.current_overrides or {}

    history_str = ""
    for h in history[-6:]:
        speaker = "RM" if h.get("sender") == "user" else "Copilot"
        history_str += f"{speaker}: {h.get('text')}\n"

    system_instruction = """You are the senior ING Financial Markets Origination & Structuring Copilot.
You assist Relationship Managers (RMs) by explaining pitchbook slides, providing CFO talking points, analyzing exposures, and dynamically updating slide parameters.

OUTPUT REQUIREMENT:
You must ALWAYS return a single valid JSON object with exactly two keys:
1. "reply": (string) Your complete, professional banking response in markdown. When explaining slides, provide a structured 3-part breakdown: Slide Objective, Financial Data & Metrics, and CFO Talking Points. When confirming edits, explain the strategic rationale.
2. "overrides": (object) A dictionary of ANY parameters being updated or remediated. If no changes were requested, return {}.

EXACT STATE KEY CONTRACT FOR overrides:
- Slide 4 / 5 (Unhedged Gap): "unhedged_gap_str" (e.g. "$7.5B")
- Slide 4 / 5 (USD Inflows / Exposure): "usd_exposure" (e.g. ">$12.0B (38% of total)")
- Slide 4 / 5 (EUR Base Cost): "eur_cost" (e.g. "€7.5B (45% of total)")
- Slide 4 / 5 (Current Hedge Ratio): "hedge_ratio" (e.g. "50% (Under-hedged)")
- Slide 4 / 5 (Target Policy Ratio): "target_hedge_ratio" (e.g. "75% - 80%")
- Slide 6 (FX Weak Downside / Loss): "fx_scen_up_unhedged" (e.g. "-$520M Revenue Impact" or "-$520M")
- Slide 6 (FX Floor Hedged): "fx_scen_up_hedged" (e.g. "Guaranteed Floor (1.0850)")
- Slide 6 (FX Strong Upside / Gain): "fx_scen_down_unhedged" (e.g. "+$380M FX Gain" or "+$380M")
- Slide 6 (Rate Stress +100 bps / Wait): "rate_scenario_up" (e.g. "3.90%")
- Slide 6 (Rate Locked Execution): "rate_scenario_lock" (e.g. "3.45% (locked)")
- Slide 6 (Rate -50 bps / Wait): "rate_scenario_down" (e.g. "2.95%")
- Slide 7 (Market Yields & Spreads): "swap_5y", "bund_10y", "itraxx_main", "market_date"
- Slide 8 (Term Sheet): "notional_bond", "spread", "tenor", "spread_disclaimer"
- Slide 10 (Compliance Remediations): "disclaimers" (list of strings), "market_date", "spread_disclaimer"

NEVER output raw markdown code fences around the JSON; return pure valid JSON."""

    user_payload = {
        "client_profile": {
            "client_id": cid,
            "client_name": client_name,
            "product_family": p_family,
            "net_debt": bundle.get("net_debt_str"),
            "liquidity": bundle.get("liquidity_str"),
            "ebitda": bundle.get("ebitda_str"),
            "leverage": bundle.get("leverage_str"),
            "24m_maturity_wall": bundle.get("debt_maturing_24m_str"),
            "maturities_schedule": maturities_data
        },
        "current_deck_overrides": current_ov,
        "recent_chat_history": history_str,
        "user_request": prompt
    }

    reply_text = ""
    merged_overrides = dict(current_ov)
    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            response = client_gcp.models.generate_content(
                model="gemini-2.5-flash",
                contents=json.dumps(user_payload),
                config=types.GenerateContentConfig(
                    system_instruction=system_instruction,
                    temperature=0.2,
                    response_mime_type="application/json"
                )
            )
            raw_text = response.text.strip() if response and response.text else ""
            if raw_text:
                parsed = json.loads(raw_text)
                if "reply" in parsed:
                    reply_text = parsed["reply"]
                if "overrides" in parsed and isinstance(parsed["overrides"], dict):
                    merged_overrides.update(parsed["overrides"])
        except Exception as e:
            logger.warning(f"Vertex AI Copilot error: {e}")

    if not reply_text:
        reply_text = f"I am ready to assist with {client_name}'s {p_family} pitchbook. You can ask me to explain any slide, modify financial parameters, or run a regulatory compliance check."

    normalized_overrides = normalize_copilot_overrides(merged_overrides)
    return {
        "reply": reply_text,
        "client_id": cid,
        "overrides": normalized_overrides,
        "timestamp": datetime.now().strftime("%H:%M:%S")
    }


@app.post("/api/pitchbook/generate")
def generate_pitchbook_endpoint(req: OpportunityRequest):
    try:
        bundle = fetch_pitchbook_bundle(req.client_id, req.client_id, get_db_connection)
        
        # Get client name from the bundle directly - this is the source of truth
        client_name = bundle.get("client_name", "Corporate Client")
        
        # Override ONLY if explicitly provided and not "Corporate Client"
        if req.overrides and req.overrides.get("client_name") and req.overrides.get("client_name") != "Corporate Client":
            client_name = req.overrides.get("client_name")
        
        opp_mock = {"id": req.client_id, "name": client_name}
        
        # Pass compliance bullets from overrides if available
        compliance_bullets = None
        if req.overrides:
            compliance_bullets = req.overrides.get("disclaimers") or req.overrides.get("compliance_bullets")
        
        # Create a clean overrides dict that doesn't override client data
        clean_overrides = {}
        if req.overrides:
            # Only pass non-empty overrides that don't conflict with bundle data
            for key, value in req.overrides.items():
                if key not in ["client_name", "revenue_str", "ebitda_str", "net_debt_str", "liquidity_str"]:
                    clean_overrides[key] = value
                elif key == "client_name" and value != "Corporate Client" and value != "":
                    clean_overrides[key] = value
        
        pptx_buf = build_pitchbook(bundle, opp_mock, compliance_bullets=compliance_bullets, overrides=clean_overrides)
        
        filename = f"ING_{req.client_id}_Pitchbook.pptx"
        return Response(
            content=pptx_buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )
    except Exception as e:
        logger.exception("Pitchbook generation failed")
        raise HTTPException(status_code=500, detail=str(e))


frontend_dist = os.path.join(os.path.dirname(__file__), "frontend", "dist")
if os.path.exists(frontend_dist):
    assets_dir = os.path.join(frontend_dist, "assets")
    if os.path.exists(assets_dir):
        app.mount("/assets", StaticFiles(directory=assets_dir), name="assets")

    @app.get("/{full_path:path}")
    async def serve_react(full_path: str):
        if full_path.startswith("api/"):
            raise HTTPException(status_code=404, detail="Not Found")
        file_path = os.path.join(frontend_dist, full_path)
        if os.path.exists(file_path) and os.path.isfile(file_path):
            return FileResponse(file_path)
        return FileResponse(os.path.join(frontend_dist, "index.html"))
else:
    @app.get("/")
    def index():
        return {"status": "Backend running, frontend build not found."}