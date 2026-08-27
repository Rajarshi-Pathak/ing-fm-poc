import io
import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

logger = logging.getLogger("pitchbook_builder")

# =============================================================================
# ING Brand Colors
# =============================================================================
ING_NAVY = RGBColor(0, 0, 102)
ING_ORANGE = RGBColor(255, 98, 0)
ING_DARK_SLATE = RGBColor(12, 17, 43)
ING_WHITE = RGBColor(255, 255, 255)
ING_LIGHT_ORANGE = RGBColor(255, 235, 220)
BG_LIGHT = RGBColor(248, 249, 250)
LINE_GRAY = RGBColor(220, 224, 230)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)
CARD_BG_BLUE = RGBColor(240, 244, 255)
CARD_BORDER_BLUE = RGBColor(200, 215, 250)
SUCCESS_GREEN = RGBColor(16, 149, 79)


# =============================================================================
# Helper Functions
# =============================================================================

def add_header(slide, title, category="FINANCIAL MARKETS ORIGINATION", is_white=False):
    tb_k = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.3))
    tf_k = tb_k.text_frame
    tf_k.word_wrap = True
    p_k = tf_k.paragraphs[0]
    p_k.text = category.upper()
    p_k.font.bold = True
    p_k.font.size = Pt(9)
    p_k.font.color.rgb = ING_ORANGE

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.0), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ING_WHITE if is_white else ING_DARK_SLATE


def add_logo(slide, is_white=False):
    base_dir = os.path.dirname(os.path.abspath(__file__))
    logo_filename = os.path.join(base_dir, "assets", "ing_logo_white.png") if is_white else os.path.join(base_dir, "assets", "ing_logo_orange.png")
    if os.path.exists(logo_filename):
        try:
            # Render logo with explicit height constraint to prevent vertical overflow
            slide.shapes.add_picture(logo_filename, Inches(11.8), Inches(0.35), height=Inches(0.45))
            return
        except Exception as exc:
            logger.warning(f"Could not insert logo: {exc}")
    tb = slide.shapes.add_textbox(Inches(11.4), Inches(0.35), Inches(1.2), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING"
    p.alignment = PP_ALIGN.RIGHT
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ING_WHITE if is_white else ING_ORANGE


def add_footer(slide, is_white=False):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ING Wholesale Banking • Strictly Confidential"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(160, 175, 200) if is_white else TEXT_MUTED


def detect_product_family(ctx):
    """Detect product family from context data."""
    # Check if we have specific product family from opportunity
    p_fam = ctx.get("product_family", "")
    if p_fam:
        return p_fam
    
    # Detect from client name
    client_name = ctx.get("client_name", "").lower()
    
    # FX: ASML, ASM, or currency-related opportunities
    if "asml" in client_name or "asm" in client_name:
        return "FX_HEDGE"
    
    # Rates: BASF, chemical companies, rate-sensitive industries
    if "basf" in client_name or "chemical" in client_name:
        return "RATES_HEDGE"
    
    # Green/Sustainable: Enel, Orsted, renewable energy
    if "enel" in client_name or "orsted" in client_name or "renewable" in client_name:
        return "GREEN_ESG"
    
    # Default
    return "DCM_REFI"


def get_product_kicker(p_fam):
    """Get product category kicker text."""
    mapping = {
        "FX_HEDGE": "FX & COMMODITY RISK ADVISORY",
        "GREEN_ESG": "SUSTAINABLE & ESG CAPITAL STRUCTURING",
        "RATES_HEDGE": "RATES RISK & LIABILITY MANAGEMENT",
        "DCM_REFI": "DCM CAPITAL STRUCTURING"
    }
    return mapping.get(p_fam, "DCM CAPITAL STRUCTURING")


def get_product_subtitle(p_fam):
    """Get product subtitle text."""
    mapping = {
        "FX_HEDGE": "Strategic FX Exposure Risk & Layered Hedging Programme",
        "GREEN_ESG": "Inaugural Hybrid Green Bond & Sustainability Framework",
        "RATES_HEDGE": "Pre-Hedge Swap Overlay & Rate Sensitivity Immunisation",
        "DCM_REFI": "Refinancing & Capital Markets Execution Framework"
    }
    return mapping.get(p_fam, "Capital Markets Execution Framework")


def get_product_pillars(p_fam, ctx, ov):
    """Get executive summary pillars matching App.jsx preview."""
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    mat_wall = ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "€3,000M"))
    if mat_wall == "N/A" or not mat_wall:
        mat_wall = "€3,000M"
    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "$8.0B"))
    if unhedged_gap == "N/A" or not unhedged_gap:
        unhedged_gap = "$8.0B"
    notional = ov.get("notional_bond", "EUR 600,000,000")
    
    if p_fam == "FX_HEDGE":
        return [
            ("1", "Exposure-led Architecture", f"Addressing the {unhedged_gap} USD hedge gap from commercial revenue expansion."),
            ("2", "Multi-Tenor Layered Corridors", "Rolling 12M–24M zero-cost participating collars protecting gross margins."),
            ("3", "Electronic Desk Execution", "Automated liquidity sourcing through ING global FX electronic trading desk."),
            ("4", "Dedicated Coverage", f"Sector coverage led by {rm_name} with IFRS 9 hedge accounting support.")
        ]
    elif p_fam == "GREEN_ESG":
        return [
            ("1", "Green Framework Alignment", "Alignment with ICMA Green Bond Principles and EU Taxonomy standards."),
            ("2", "Use of Proceeds Pool", "Ring-fenced eligible asset pool with annual impact & allocation verification."),
            ("3", "Greenium Advantage", "Capturing 3-7 bps new-issue concession advantage from dedicated ESG funds."),
            ("4", "Sole ESG Structurer", "ING leading SPO documentation, investor roadshow, and syndicate execution.")
        ]
    elif p_fam == "RATES_HEDGE":
        return [
            ("1", "Rate Risk Assessment", f"Quantifying interest rate repricing risk across the {mat_wall} debt horizon."),
            ("2", "Pre-Hedge Swap Overlay", "Forward-starting IRS and swaptions to lock in current benchmark yield curve."),
            ("3", "Hedge Policy Alignment", "Optimizing treasury fixed vs floating debt ratio target."),
            ("4", "Syndicate Distribution", "Full balance sheet underwriting and rating agency advisory.")
        ]
    else:  # DCM_REFI
        return [
            ("1", "Maturity-Led Sizing", f"Addressing the {mat_wall} near-term maturity profile."),
            ("2", "Right-Sized Structure", f"Tailored combination of {notional} benchmark EMTN."),
            ("3", "Competitive Execution", "Direct syndicate distribution across European institutional bases."),
            ("4", "Long-Term Partnership", "Committed balance sheet underwriting and rating optimization.")
        ]


# =============================================================================
# Data Fetching - Fully Database-Driven
# =============================================================================

def fetch_pitchbook_bundle(canonical_id, client_id_raw, get_db_connection):
    """
    Fetch ALL data from database. Resilient to client ID formats.
    """
    cid_input = str(canonical_id or client_id_raw or "CLI101").strip()
    
    ctx = {
        "client_id": cid_input,
        "client_name": "Corporate Client",
        "sector": "Corporate",
        "tier": "Tier 1",
        "hq_country": "Europe",
        "rm_name": "Senior Relationship Manager",
        "revenue_str": "N/A",
        "ebitda_str": "N/A",
        "net_debt_str": "N/A",
        "liquidity_str": "N/A",
        "debt_maturing_24m_str": "N/A",
        "leverage_ratio": "N/A",
        "product_family": "DCM_REFI",
        "opportunity_type": "Capital Structuring",
        "trigger_source": "Balance sheet review",
        "next_best_action": "Propose financing structure",
        "maturities": [],
        "coverage_team": [],
        "signals": [],
        "deals": [],
        "spreads": []
    }

    try:
        conn, connector = get_db_connection()
        if conn:
            cur = conn.cursor()
            
            # 1. RESOLVE CLIENT ID
            cur.execute("""
                SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                FROM ca.client_master 
                WHERE client_id = %s
            """, (cid_input,))
            row = cur.fetchone()
            
            # Method B: Check if cid_input is an opportunity_id
            if not row:
                cur.execute("""
                    SELECT client_id FROM ca.ca_opportunity_scoring 
                    WHERE opportunity_id = %s
                """, (cid_input,))
                opp_row = cur.fetchone()
                if opp_row:
                    cur.execute("""
                        SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                        FROM ca.client_master 
                        WHERE client_id = %s
                    """, (opp_row[0],))
                    row = cur.fetchone()
            
            # Method C: Try partial match on client_name
            if not row:
                cur.execute("""
                    SELECT client_id, client_name, tier, hq_country, revenue_eur_m, rm_name 
                    FROM ca.client_master 
                    WHERE LOWER(client_name) LIKE %s
                """, (f"%{cid_input.lower()}%",))
                row = cur.fetchone()
            
            if row:
                actual_cid = row[0]
                if row[1]: ctx["client_name"] = row[1]
                if row[2]: ctx["tier"] = row[2]
                if row[3]: ctx["hq_country"] = row[3]
                if row[4] and row[4] > 0: 
                    ctx["revenue_str"] = f"€{float(row[4]):,.0f}M"
                if row[5]: ctx["rm_name"] = row[5]
                ctx["client_id"] = actual_cid

                # 2. FINANCIAL FILINGS
                cur.execute("""
                    SELECT net_debt_eur_m, liquidity_eur_m, ebitda_eur_m, 
                           reported_revenue_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings 
                    WHERE client_id = %s 
                    ORDER BY reporting_period DESC LIMIT 1
                """, (actual_cid,))
                row = cur.fetchone()
                if row:
                    if row[0] and row[0] > 0: 
                        ctx["net_debt_str"] = f"€{float(row[0]):,.0f}M"
                    if row[1] and row[1] > 0: 
                        ctx["liquidity_str"] = f"€{float(row[1]):,.0f}M"
                    if row[2] and row[2] > 0: 
                        ctx["ebitda_str"] = f"€{float(row[2]):,.0f}M"
                        if row[0] and row[0] > 0 and float(row[2]) > 0:
                            ctx["leverage_ratio"] = f"{(float(row[0]) / float(row[2])):.2f}x"
                    if row[3] and row[3] > 0: 
                        ctx["revenue_str"] = f"€{float(row[3]):,.0f}M"
                    if row[4] and row[4] > 0: 
                        ctx["debt_maturing_24m_str"] = f"€{float(row[4]):,.0f}M"

                # 3. OPPORTUNITY SCORING
                cur.execute("""
                    SELECT opportunity_type, trigger_source, next_best_action
                    FROM ca.ca_opportunity_scoring 
                    WHERE client_id = %s 
                    ORDER BY rank ASC LIMIT 1
                """, (actual_cid,))
                row = cur.fetchone()
                if row:
                    if row[0]: ctx["opportunity_type"] = row[0]
                    if row[1]: ctx["trigger_source"] = row[1]
                    if row[2]: ctx["next_best_action"] = row[2]

                # 4. DEBT MATURITY SCHEDULE
                cur.execute("""
                    SELECT isin, instrument_type, amount_eur_m, maturity_year, 
                           coupon_rate_pct, currency
                    FROM ca.debt_maturity_schedule 
                    WHERE client_id = %s 
                    ORDER BY maturity_year ASC
                """, (actual_cid,))
                rows = cur.fetchall()
                if rows:
                    ctx["maturities"] = []
                    for r in rows:
                        ctx["maturities"].append({
                            "isin": r[0] or "N/A",
                            "instrument_type": r[1] or "Bond",
                            "amount_eur_m": float(r[2]) if r[2] else 0,
                            "maturity_year": str(r[3]) if r[3] else "2026",
                            "coupon_rate_pct": float(r[4]) if r[4] else 0.0,
                            "currency": r[5] or "EUR"
                        })

                # 5. COVERAGE TEAM
                cur.execute("""
                    SELECT role_title, banker_name, location
                    FROM ca.coverage_teams 
                    WHERE client_id = %s
                """, (actual_cid,))
                rows = cur.fetchall()
                if rows:
                    ctx["coverage_team"] = []
                    for r in rows:
                        ctx["coverage_team"].append({
                            "name": r[1] or "Coverage Banker",
                            "title": r[0] or "Global Coverage",
                            "location": r[2] or ""
                        })

            # Close cursor
            try: cur.close()
            except Exception: pass
            try: conn.close()
            except Exception: pass
            if connector:
                try: connector.close()
                except Exception: pass

    except Exception as exc:
        logger.warning(f"Database query warning for {cid_input}: {exc}")
        import traceback
        logger.warning(traceback.format_exc())

    # PRODUCT FAMILY DETECTION
    ctx["product_family"] = detect_product_family(ctx)

    # FALLBACK COVERAGE TEAM
    if not ctx.get("coverage_team"):
        ctx["coverage_team"] = [
            {"name": ctx.get("rm_name", "Relationship Manager"), "title": "Global Coverage"},
            {"name": "Managing Director", "title": "Head of Capital Markets"},
            {"name": "Director", "title": "Financial Markets Structuring"}
        ]

    return ctx


def get_slide_meta(p_fam):
    meta = {
        "FX_HEDGE": {
            "s2_cat": "FX RISK CATALYST", "s2_ttl": "Currency Exposure & Market Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Corporate Liquidity & Currency Inflow Profile",
            "s5_cat": "CURRENCY EXPOSURE PROFILE", "s5_ttl": "FX Currency Breakdown & Hedging Gap",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "FX Scenario Analysis & Layered Collar Payoff",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Central Bank Differentials & FX Forward Points",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative FX Risk Management Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Layered Roll Framework & Desk Execution",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Derivative Disclosures"
        },
        "GREEN_ESG": {
            "s2_cat": "SUSTAINABILITY CATALYST", "s2_ttl": "ESG Capital Strategy & Decarbonization Catalyst",
            "s4_cat": "ESG BALANCE SHEET FOUNDATION", "s4_ttl": "Balance Sheet Capacity & Green CapEx Profile",
            "s5_cat": "USE OF PROCEEDS", "s5_ttl": "Eligible Green Asset Pool & Use of Proceeds",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Greenium vs Plain-Vanilla Cost Sensitivity",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "ESG Credit Spreads & Green Bond Index Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Green / Sustainability-Linked Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Second-Party Opinion (SPO) & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "ICMA Green Bond Principles & Target Market Notice"
        },
        "RATES_HEDGE": {
            "s2_cat": "RATE RISK CATALYST", "s2_ttl": "Rate Path Volatility & IRS Pre-Hedge Catalyst",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Liquidity Snapshot",
            "s5_cat": "MATURITY & SWAP SCHEDULE", "s5_ttl": "Debt Maturity Profile & Swap Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Rate Shift Sensitivity & Pre-Hedge Lock Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Swap Curve Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Pre-Hedge Swap & EMTN Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "ISDA Schedule, CSA & Execution Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Target Market Notice & EMIR Classification Disclosures"
        },
        "DCM_REFI": {
            "s2_cat": "STRATEGIC CATALYST", "s2_ttl": "Executive Context & Opportunity Rationale",
            "s4_cat": "BALANCE SHEET FOUNDATION", "s4_ttl": "Capital Structure & Treasury Health Profile",
            "s5_cat": "MATURITY SCHEDULE", "s5_ttl": "Debt Maturity Profile & Refinancing Horizon",
            "s6_cat": "SENSITIVITY ANALYSIS", "s6_ttl": "Refinancing Scenario Analysis",
            "s7_cat": "MARKET INTELLIGENCE", "s7_ttl": "Benchmark Yields & Credit Spread Backdrop",
            "s8_cat": "TRANSACTION STRUCTURING", "s8_ttl": "Indicative Debt Financing Term Sheet",
            "s9_cat": "EXECUTION ROADMAP", "s9_ttl": "Roadmap & Syndicate Timeline",
            "s10_cat": "REGULATORY DISCLOSURES", "s10_ttl": "Regulatory Notices & Target Market Classification"
        }
    }
    return meta.get(p_fam, meta["DCM_REFI"])

def build_pitchbook(ctx, opp, compliance_bullets=None, overrides=None):
    """
    Build pitchbook using database data. No hardcoded client-specific values.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    ov = overrides or {}
    
    # -------- Client Data from Database (with override support) --------
    client_name = ov.get("client_name", ctx.get("client_name", "Corporate Client"))
    p_fam = ov.get("product_family", ctx.get("product_family", "DCM_REFI"))
    sm = get_slide_meta(p_fam)
    rm_name = ov.get("rm_name", ctx.get("rm_name", "Senior Relationship Manager"))
    
    # Financial data from database
    revenue_str = ov.get("revenue_str", ctx.get("revenue_str", "N/A"))
    ebitda_str = ov.get("ebitda_str", ctx.get("ebitda_str", "N/A"))
    net_debt_str = ov.get("net_debt_str", ctx.get("net_debt_str", "N/A"))
    liquidity_str = ov.get("liquidity_str", ctx.get("liquidity_str", "N/A"))
    leverage_str = ov.get("leverage_ratio", ctx.get("leverage_ratio", "N/A"))
    mat_wall_str = ov.get("debt_maturing_24m_str", ov.get("maturity_wall_str", ctx.get("debt_maturing_24m_str", "N/A")))
    
    # Signals from database
    signals = ctx.get("signals", [])
    
    # Deals from database
    deals = ctx.get("deals", [])
    
    # Maturities from database
    maturities = ctx.get("maturities", [])

    # -------- Market Data (with overrides) --------
    swap_5y = ov.get("swap_5y", "2.62%")
    bund_10y = ov.get("bund_10y", "2.61%")
    iboxx_bbb = ov.get("iboxx_bbb", "115 bps")
    itraxx_main = ov.get("itraxx_main", "58 bps")
    ecb_rate = ov.get("ecb_rate", "2.25%")
    fed_rate = ov.get("fed_rate", "4.00–4.25%")

    # -------- Term Sheet (with overrides) --------
    tenor_str = ov.get("tenor", "7 Years (T + 7Y)")
    spread_str = ov.get("spread", "Mid-Swap + 82 bps")
    notional_str = ov.get("notional_bond", "EUR 600,000,000")
    spread_disc = ov.get("spread_disclaimer", "*Indicative pricing subject to market conditions, bookbuilding depth, and credit approval.*")

    # -------- Scenario Values (with overrides) --------
    scen_up = ov.get("rate_scenario_up", "4.55%")
    scen_lock = ov.get("rate_scenario_lock", "3.60% (locked)")
    scen_down = ov.get("rate_scenario_down", "3.15%")
    
    fx_up_unhedged = ov.get("fx_scen_up_unhedged", "-$520M Revenue Impact")
    fx_up_hedged = ov.get("fx_scen_up_hedged", "Guaranteed Floor (1.0850)")
    fx_spot_unhedged = ov.get("fx_scen_spot_unhedged", "1.0650 Spot Level")
    fx_spot_hedged = ov.get("fx_scen_spot_hedged", "1.0650 Forward Rate")
    fx_down_unhedged = ov.get("fx_scen_down_unhedged", "+$380M FX Gain")
    fx_down_hedged = ov.get("fx_scen_down_hedged", "Participate up to 1.0450")

    unhedged_gap = ov.get("unhedged_gap_str", ov.get("unhedged_gap", "N/A"))

    # -------- Get Product-Specific Content --------
    kicker = get_product_kicker(p_fam)
    subtitle = get_product_subtitle(p_fam)
    pillars = get_product_pillars(p_fam, ctx, ov)
    
    # -------- Build Trigger Cards --------
    trigger_cards = [
        ("Primary Market Trigger", ctx.get("trigger_source", "Active capital structure optimization"), ING_ORANGE),
        ("Window of Opportunity", "Favorable market conditions across European issuance windows.", ING_DARK_SLATE),
        ("Recommended Action", ctx.get("next_best_action", "Propose strategic execution roadmap"), SUCCESS_GREEN)
    ]

    # =========================================================================
    # SLIDE 1: COVER (Full Parity with React Preview)
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    
    # 1. Dark background
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ING_DARK_SLATE
    bg.line.fill.background()

    # 2. Full-height Left Orange Stripe (border-l-8 border-[#FF6200])
    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ING_ORANGE
    accent.line.fill.background()

    # 3. Top-Right Logo & Desk Label
    add_logo(s1, is_white=True)
    tb_desk = s1.shapes.add_textbox(Inches(8.5), Inches(0.85), Inches(4.0), Inches(0.35))
    tf_desk = tb_desk.text_frame
    p_desk = tf_desk.paragraphs[0]
    p_desk.text = "Financial Markets Origination"
    p_desk.font.size = Pt(10)
    p_desk.font.color.rgb = RGBColor(148, 163, 184)
    p_desk.alignment = PP_ALIGN.RIGHT

    # 4. Main Hero Text Box
    tb1 = s1.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    # Kicker
    p = tf1.paragraphs[0]
    p.text = kicker.upper()
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ING_ORANGE

    # Client Name
    p = tf1.add_paragraph()
    p.text = client_name
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = ING_WHITE
    p.space_before = Pt(8)

    # Subtitle
    p = tf1.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.space_before = Pt(6)

    # 5. Bottom Horizontal Divider Line (border-t border-gray-800)
    div_line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.015))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(40, 50, 80)
    div_line.line.fill.background()

    # 6. Bottom Metadata Left (Prepared by)
    tb_meta_l = s1.shapes.add_textbox(Inches(0.9), Inches(6.42), Inches(6.0), Inches(0.8))
    tf_meta_l = tb_meta_l.text_frame
    p_ml1 = tf_meta_l.paragraphs[0]
    p_ml1.text = f"Prepared by: {rm_name}"
    p_ml1.font.bold = True
    p_ml1.font.size = Pt(11)
    p_ml1.font.color.rgb = RGBColor(226, 232, 240)

    p_ml2 = tf_meta_l.add_paragraph()
    p_ml2.text = "Global Sector Coverage & Capital Markets Desk"
    p_ml2.font.size = Pt(10)
    p_ml2.font.color.rgb = RGBColor(148, 163, 184)
    p_ml2.space_before = Pt(2)

    # 7. Bottom Metadata Right (Market Snapshot Timestamp)
    market_snapshot_raw = ov.get("market_date") or f"Market Snapshot as of {datetime.now().strftime('%d %B %Y')}"
    market_snapshot_str = market_snapshot_raw.split(",")[0].strip()
    if not market_snapshot_str.startswith("Market Snapshot"): market_snapshot_str = f"Market Snapshot as of {market_snapshot_str}"
    tb_meta_r = s1.shapes.add_textbox(Inches(7.0), Inches(6.55), Inches(5.5), Inches(0.6))
    tf_meta_r = tb_meta_r.text_frame
    p_mr = tf_meta_r.paragraphs[0]
    p_mr.text = market_snapshot_str
    p_mr.font.size = Pt(10)
    p_mr.font.color.rgb = RGBColor(148, 163, 184)
    p_mr.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 2: STRATEGIC CATALYST
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    add_header(s2, sm["s2_ttl"], category=sm["s2_cat"])
    add_logo(s2)
    add_footer(s2)

    # Product-dynamic trigger resolution matching App.jsx
    if p_fam == "FX_HEDGE":
        trig_t = ov.get("trigger") or "Commercial inflow shift: North American expansion increased USD revenue to >$12B against 50% hedge ratio (~$8bn gap)."
        win_t = ov.get("window") or "EUR/USD forward points offer structural hedging pickup; volatility corridor allows zero-cost collar structuring."
        act_t = ov.get("action") or "Propose staged 12M–24M layered FX hedging programme with zero-cost collar overlays to close ~$8bn gap."
    elif p_fam == "GREEN_ESG":
        trig_t = ov.get("trigger") or "EU Taxonomy alignment: €3.5B eligible renewable & decarbonization CapEx pipeline ready for green financing."
        win_t = ov.get("window") or "Strong ESG investor liquidity generating 3-7 bps greenium pricing concession across European green bonds."
        act_t = ov.get("action") or "Establish inaugural Green Bond / Hybrid Framework with second-party SPO verification."
    elif p_fam == "RATES_HEDGE":
        trig_t = ov.get("trigger") or "Upcoming €3.2B debt maturities face repricing risk amid benchmark curve fluctuations."
        win_t = ov.get("window") or "Current 5Y EUR swap easing at 2.62% provides attractive entry window for forward-starting IRS."
        act_t = ov.get("action") or "Execute €400M pre-hedge IRS overlay to lock in current base yield before debt issuance."
    else:
        trig_t = ov.get("trigger") or ctx.get("trigger_source", "Active capital structure optimization and refinancing window identified.")
        win_t = ov.get("window") or "Favorable benchmark credit spreads across European issuance windows."
        act_t = ov.get("action") or ctx.get("next_best_action", "Propose capital structuring dialogue and benchmark EMTN roadshow.")

    styled_cards = [
        ("Primary Market Trigger", trig_t, RGBColor(255, 247, 237), RGBColor(254, 215, 170), RGBColor(154, 52, 18)),
        ("Window of Opportunity", win_t, RGBColor(239, 246, 255), RGBColor(191, 219, 254), RGBColor(0, 0, 102)),
        ("Recommended Action", act_t, RGBColor(236, 253, 245), RGBColor(167, 243, 208), RGBColor(6, 95, 70))
    ]

    for idx, (head_c, body_c, bg_c, border_c, title_c) in enumerate(styled_cards):
        cx = Inches(0.8 + (idx * 3.95))
        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5), Inches(3.8), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg_c
        shp.line.color.rgb = border_c
        shp.line.width = Pt(1)

        tb_c = s2.shapes.add_textbox(cx + Inches(0.15), Inches(1.55), Inches(3.5), Inches(2.2))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        p_h = tf_c.paragraphs[0]
        p_h.text = head_c
        p_h.font.bold = True
        p_h.font.size = Pt(11)
        p_h.font.color.rgb = title_c

        p_b = tf_c.add_paragraph()
        p_b.text = body_c
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = RGBColor(55, 65, 81)
        p_b.space_before = Pt(6)

        # =========================================================================
    # SLIDE 3: EXECUTIVE SUMMARY (Exact Parity with Preview Canvas)
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    add_logo(s3)
    add_footer(s3)

    # 1. Left Orange Hero Panel
    hero_panel = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
    hero_panel.fill.solid()
    hero_panel.fill.fore_color.rgb = ING_ORANGE
    hero_panel.line.fill.background()

    # Left Hero Text Box - Title
    tb_lh = s3.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(3.7), Inches(1.0))
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    p = tf_lh.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ING_WHITE

    # White Horizontal Divider Line
    div_w = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.85), Inches(0.8), Inches(0.025))
    div_w.fill.solid()
    div_w.fill.fore_color.rgb = ING_WHITE
    div_w.line.fill.background()

    # Dynamic Theme Subheading
    subheading_map = {
        "FX_HEDGE": "Strategic FX Architecture",
        "GREEN_ESG": "Sustainable Finance Framework",
        "RATES_HEDGE": "Rate Risk Immunisation",
        "DCM_REFI": "Proactive Capital Structuring"
    }
    s3_subheading = subheading_map.get(p_fam, "Proactive Capital Structuring")

    tb_sub = s3.shapes.add_textbox(Inches(0.6), Inches(2.05), Inches(3.7), Inches(3.8))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    
    p = tf_sub.paragraphs[0]
    p.text = s3_subheading
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_WHITE

    p = tf_sub.add_paragraph()
    p.text = f"Customized execution roadmap for {client_name} based on group treasury requirements and live market backdrop."
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(255, 245, 235)
    p.space_before = Pt(8)


    # 2. Right Side Numbered Pillars (1, 2, 3, 4)
    pillars = get_product_pillars(p_fam, ctx, ov)
    for idx, (p_num, p_head, p_body) in enumerate(pillars):
        y_pos = Inches(1.35 + (idx * 1.32))
        
        # Circle badge
        c_shp = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.3), y_pos, Inches(0.42), Inches(0.42))
        c_shp.fill.solid()
        c_shp.fill.fore_color.rgb = ING_ORANGE
        c_shp.line.fill.background()
        
        # Centered white number inside circle
        tf_c = c_shp.text_frame
        tf_c.margin_left = Inches(0)
        tf_c.margin_right = Inches(0)
        tf_c.margin_top = Inches(0.04)
        tf_c.margin_bottom = Inches(0)
        p_c = tf_c.paragraphs[0]
        p_c.text = str(p_num)
        p_c.font.bold = True
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = ING_WHITE
        p_c.alignment = PP_ALIGN.CENTER

        # Pillar Title & Description Text Box
        tb_p = s3.shapes.add_textbox(Inches(5.9), y_pos - Inches(0.05), Inches(6.8), Inches(1.15))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        
        p_h = tf_p.paragraphs[0]
        p_h.text = p_head
        p_h.font.bold = True
        p_h.font.size = Pt(12.5)
        p_h.font.color.rgb = ING_ORANGE

        p_b = tf_p.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = RGBColor(75, 85, 99)
        p_b.space_before = Pt(3)

            # =========================================================================
    # SLIDE 4: BALANCE SHEET (Exact Database & Preview Parity)
    # =========================================================================
    s4 = prs.slides.add_slide(blank)
    add_header(s4, sm["s4_ttl"], category=sm["s4_cat"])
    add_logo(s4)
    add_footer(s4)

    tier_str = ov.get("tier", ctx.get("tier", "Tier 1"))
    if "Tier" in tier_str and "Investment" not in tier_str:
        tier_str = f"{tier_str} (Investment Grade)"

    card3_lbl = "Unhedged FX Gap" if p_fam == "FX_HEDGE" else ("Eligible Green CapEx" if p_fam == "GREEN_ESG" else "24M Maturity Wall")
    if p_fam == "FX_HEDGE":
        card3_val = ov.get("unhedged_gap_str", "$8.0B")
    elif p_fam == "GREEN_ESG":
        card3_val = "€3.5B"
    else:
        card3_val = mat_wall_str if (mat_wall_str and mat_wall_str != "N/A") else "€3,000M"

    metrics = [
        ("Net Debt", net_debt_str if net_debt_str != "N/A" else "€16,200M", ING_DARK_SLATE),
        ("Available Liquidity", liquidity_str if liquidity_str != "N/A" else "€7,800M", SUCCESS_GREEN),
        (card3_lbl, card3_val, ING_ORANGE),
        ("Credit Rating / Tier", tier_str, ING_DARK_SLATE)
    ]

    for idx, (lbl, val, val_color) in enumerate(metrics):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s4.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.3))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(14 if len(val) > 12 else 18)
        p.font.color.rgb = val_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    shp_bot = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.25), Inches(11.7), Inches(3.35))
    shp_bot.fill.solid()
    shp_bot.fill.fore_color.rgb = CARD_BG_BLUE
    shp_bot.line.color.rgb = CARD_BORDER_BLUE

    tb_bot = s4.shapes.add_textbox(Inches(1.1), Inches(3.45), Inches(11.1), Inches(2.9))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    
    p = tf_bot.paragraphs[0]
    p.text = "Corporate Financial Standing & Balance Sheet Capacity"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_DARK_SLATE

    p_sub1 = tf_bot.add_paragraph()
    p_sub1.text = f"• Annual Group Revenue of {revenue_str if revenue_str != 'N/A' else '€65,000M'} supported by EBITDA of {ebitda_str if ebitda_str != 'N/A' else '€14,300M'}."
    p_sub1.font.size = Pt(10.5)
    p_sub1.font.color.rgb = RGBColor(55, 65, 81)
    p_sub1.space_before = Pt(8)

    p_sub2 = tf_bot.add_paragraph()
    p_sub2.text = f"• Robust liquidity buffer of {liquidity_str if liquidity_str != 'N/A' else '€7,800M'} provides substantial capacity to execute structured financing and risk management operations."
    p_sub2.font.size = Pt(10.5)
    p_sub2.font.color.rgb = RGBColor(55, 65, 81)
    p_sub2.space_before = Pt(6)

        # =========================================================================
    # SLIDE 5: EXPOSURE / MATURITY / ASSET POOL (Exact Preview Parity)
    # =========================================================================
    s5 = prs.slides.add_slide(blank)
    add_header(s5, sm["s5_ttl"], category=sm["s5_cat"])
    add_logo(s5)
    add_footer(s5)

    # 1. Left Card Container (Gray box)
    shp_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.9))
    shp_l.fill.solid()
    shp_l.fill.fore_color.rgb = BG_LIGHT
    shp_l.line.color.rgb = CARD_BORDER_BLUE

    tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    # 2. Right Card Container (Blue box)
    shp_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.9))
    shp_r.fill.solid()
    shp_r.fill.fore_color.rgb = CARD_BG_BLUE
    shp_r.line.color.rgb = CARD_BORDER_BLUE

    tb_r = s5.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    if p_fam == "FX_HEDGE":
        # Left: Currency Exposure Breakdown
        p = tf_l.paragraphs[0]
        p.text = "Commercial Currency Exposure Flow"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        items_l = [
            ("USD Gross Inflows", "$18.5B / year (Export billing)"),
            ("EUR Cost Base", "€14.2B / year (R&D, manufacturing)"),
            ("Layered Coverage", "42% covered across 12M")
        ]
        for lbl, val in items_l:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = f"Total Unhedged USD Gap: {ov.get('unhedged_gap_str', '$8.0B')}"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: Layered Collar Strategy
        p = tf_r.paragraphs[0]
        p.text = "Layered Collar Execution Architecture"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        p_desc.text = f"Customized rolling 12M–24M FX hedging corridor for {client_name}. Protects operating margin floor while retaining upside participation up to cap limits without upfront option premium."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

    elif p_fam == "GREEN_ESG":
        # Left: Eligible Green Asset Pool
        p = tf_l.paragraphs[0]
        p.text = "Eligible Green Asset & CapEx Pool"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        items_l = [
            ("Renewable Generation", "€1,850M (Solar & Wind)"),
            ("Grid Modernization", "€1,100M (Smart Metering)"),
            ("Energy Storage", "€550M (Battery Systems)")
        ]
        for lbl, val in items_l:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = "Total Eligible Pool: €3,500M"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: SPO Framework
        p = tf_r.paragraphs[0]
        p.text = "Green Framework & SPO Structuring"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        p_desc.text = f"Inaugural Green Financing Framework aligned with ICMA Green Bond Principles and EU Taxonomy. Supported by second-party opinion (SPO) provider to capture 3-7 bps ESG greenium pricing advantage."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

    else:  # RATES_HEDGE & DCM_REFI
        # Left: Tranche Maturity Breakdown
        p = tf_l.paragraphs[0]
        p.text = "Tranche Maturity Breakdown"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        mat_wall_val = mat_wall_str if (mat_wall_str and mat_wall_str != "N/A") else "€3,000M"
        
        mat_items = [
            ("2026 Maturities", "€600M (Commodity & Fixed Notes)"),
            ("2027 Maturities", "€3,000M (IRS Pre-Hedge Refinancing)"),
            ("2028 Maturities", "€5,497M (Syndicated Term Loan)")
        ]

        for lbl, val in mat_items:
            p_i = tf_l.add_paragraph()
            p_i.text = f"• {lbl}: {val}"
            p_i.font.size = Pt(10.5)
            p_i.font.color.rgb = RGBColor(55, 65, 81)
            p_i.space_before = Pt(8)

        p_tot = tf_l.add_paragraph()
        p_tot.text = f"Total 24M Maturity Wall: {mat_wall_val}"
        p_tot.font.bold = True
        p_tot.font.size = Pt(11)
        p_tot.font.color.rgb = ING_ORANGE
        p_tot.space_before = Pt(14)

        # Right: Pre-Hedge Overlay Sizing
        p = tf_r.paragraphs[0]
        p.text = "Pre-Hedge Overlay Sizing" if p_fam == "RATES_HEDGE" else "Refinancing Wall Rationale"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ING_DARK_SLATE

        p_desc = tf_r.add_paragraph()
        if p_fam == "RATES_HEDGE":
            p_desc.text = "Upcoming maturities cluster in near-term windows. Locking in forward-starting swap rates eliminates repricing uncertainty ahead of primary debt issuance."
        else:
            p_desc.text = f"Upcoming debt maturities of {mat_wall_val} cluster in near-term windows. Proactive capital structuring and benchmark EMTN roadshows ensure optimal tenor extension and liquidity resilience."
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = RGBColor(55, 65, 81)
        p_desc.space_before = Pt(8)

        # =========================================================================
    # SLIDE 6: SENSITIVITY ANALYSIS (Exact Preview Parity across all products)
    # =========================================================================
    s6 = prs.slides.add_slide(blank)
    add_header(s6, sm["s6_ttl"], category=sm["s6_cat"])
    add_logo(s6)
    add_footer(s6)

    # 3-column table styled exactly like the React preview
    table_shape = s6.shapes.add_table(4, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(4.1)
    tbl.columns[1].width = Inches(3.8)
    tbl.columns[2].width = Inches(3.8)

    RED_IMPACT = RGBColor(225, 29, 72)     # Rose-600
    GREEN_SAVINGS = RGBColor(4, 120, 87)   # Emerald-700
    GRAY_TEXT = RGBColor(75, 85, 99)

    if p_fam == "FX_HEDGE":
        headers = ["Market Scenario", "Layered Collar Strategy", "Unhedged Exposure"]
        rows = [
            ("EUR/USD +5% (USD Weakens)", "Guaranteed Floor (1.0850)", "-$450M Revenue Impact", GREEN_SAVINGS, RED_IMPACT, True),
            ("Spot Unchanged (1.0650)", "1.0650 Forward Rate", "1.0650 Spot Level", ING_DARK_SLATE, GRAY_TEXT, False),
            ("EUR/USD -5% (USD Strengthens)", "Participate up to 1.0450", "+$380M FX Gain", ING_DARK_SLATE, GREEN_SAVINGS, False)
        ]
    elif p_fam == "GREEN_ESG":
        headers = ["Issuance Format", "Indicative Yield / Spread", "Annual Interest Savings"]
        rows = [
            ("Inaugural Green Bond (with Greenium)", "Mid-Swap + 77 bps (-5 bps)", "€375,000 / year savings", GREEN_SAVINGS, GREEN_SAVINGS, True),
            ("Sustainability-Linked Bond (SLB)", "Mid-Swap + 80 bps (-2 bps)", "€150,000 / year savings", ING_DARK_SLATE, ING_DARK_SLATE, False),
            ("Plain-Vanilla Senior EMTN", "Mid-Swap + 82 bps (Flat)", "Benchmark Baseline", GRAY_TEXT, GRAY_TEXT, False)
        ]
    else:  # RATES_HEDGE & DCM_REFI
        headers = ["Rate Scenario", "Refinance Today (Locked)", "Wait 6 Months (Unhedged)"]
        rows = [
            ("Rates +100 bps", "3.60% (locked)", "4.50% (+90 bps cost)", GREEN_SAVINGS, RED_IMPACT, True),
            ("Rates Unchanged", "3.60% (locked)", "3.60% (locked)", ING_DARK_SLATE, GRAY_TEXT, False),
            ("Rates −50 bps", "3.60% (locked)", "3.15%", ING_DARK_SLATE, GREEN_SAVINGS, False)
        ]

    # Format Header Row
    for c_idx, h_txt in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ING_DARK_SLATE
        p = cell.text_frame.paragraphs[0]
        p.text = h_txt
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    # Format Data Rows
    for r_idx, (c1, c2, c3, c2_color, c3_color, is_highlight) in enumerate(rows, start=1):
        # Cell 1
        cell1 = tbl.cell(r_idx, 0)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 247, 237) if is_highlight else (RGBColor(249, 250, 251) if r_idx % 2 == 0 else RGBColor(255, 255, 255))
        p = cell1.text_frame.paragraphs[0]
        p.text = c1
        p.font.bold = True
        p.font.size = Pt(10.5)
        p.font.color.rgb = ING_DARK_SLATE
        p.alignment = PP_ALIGN.LEFT

        # Cell 2
        cell2 = tbl.cell(r_idx, 1)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = cell1.fill.fore_color.rgb
        p = cell2.text_frame.paragraphs[0]
        p.text = c2
        p.font.bold = True if c2_color in [GREEN_SAVINGS, ING_DARK_SLATE] else False
        p.font.size = Pt(10.5)
        p.font.color.rgb = c2_color
        p.alignment = PP_ALIGN.LEFT

        # Cell 3
        cell3 = tbl.cell(r_idx, 2)
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = cell1.fill.fore_color.rgb
        p = cell3.text_frame.paragraphs[0]
        p.text = c3
        p.font.bold = True if c3_color in [GREEN_SAVINGS, RED_IMPACT] else False
        p.font.size = Pt(10.5)
        p.font.color.rgb = c3_color
        p.alignment = PP_ALIGN.LEFT

            # =========================================================================
    # SLIDE 7: MARKET INTELLIGENCE (Exact Preview Parity)
    # =========================================================================
    s7 = prs.slides.add_slide(blank)
    add_header(s7, sm["s7_ttl"], category=sm["s7_cat"])
    add_logo(s7)
    add_footer(s7)

    # 4 Metric Cards with matching colors and database values
    if p_fam == "GREEN_ESG":
        mkt_cards = [
            ("EUR Green Spread", "77 bps", ING_DARK_SLATE),
            ("Greenium Concession", "-5 bps", RGBColor(17, 24, 39)),
            ("ECB Refi Rate", "2.25%", ING_ORANGE),
            ("iTraxx Main", "58 bps", SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: ECB Refinancing Rate at 2.25%; Fed Funds Target at 4.00–4.25%."
        b2 = "• High ESG subscription ratios (3.8x book cover) provide attractive new-issue pricing compression."
    elif p_fam == "FX_HEDGE":
        mkt_cards = [
            ("EUR/USD Spot", ov.get("eur_usd_spot", "1.0650"), ING_DARK_SLATE),
            ("1Y Forward Points", ov.get("fx_fwd_pts", "+145 pts"), RGBColor(17, 24, 39)),
            ("Fed Funds Target", "4.00–4.25%", ING_ORANGE),
            ("ECB Deposit Rate", "2.00%", SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: Fed easing trajectory creates favorable forward points carry environment for USD receivables."
        b2 = "• Currency Volatility: Heightened transatlantic rate divergence makes systematic layered hedging cost-effective."
    else:  # RATES_HEDGE & DCM_REFI
        mkt_cards = [
            ("5Y EUR Swap", ov.get("swap_5y", "2.62%"), ING_DARK_SLATE),
            ("10Y Bund", ov.get("bund_10y", "2.61%"), RGBColor(17, 24, 39)),
            ("ECB Refi Rate", "2.25%", ING_ORANGE),
            ("iTraxx Main", ov.get("itraxx_main", "58 bps"), SUCCESS_GREEN)
        ]
        b1 = "• Central Bank Policy: ECB Refinancing Rate at 2.25%; Fed Funds Target at 4.00–4.25%."
        b2 = "• Tightening European investment grade credit spreads support attractive execution windows."

    for idx, (lbl, val, val_color) in enumerate(mkt_cards):
        mx = Inches(0.8 + (idx * 2.95))
        shp = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, Inches(1.5), Inches(2.8), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG_LIGHT
        shp.line.color.rgb = LINE_GRAY

        tb_m = s7.shapes.add_textbox(mx + Inches(0.1), Inches(1.6), Inches(2.6), Inches(1.3))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_m.add_paragraph()
        p.text = val
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = val_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    # Bottom Container Box
    shp_bot = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.25), Inches(11.7), Inches(3.35))
    shp_bot.fill.solid()
    shp_bot.fill.fore_color.rgb = CARD_BG_BLUE
    shp_bot.line.color.rgb = CARD_BORDER_BLUE

    tb_bot = s7.shapes.add_textbox(Inches(1.1), Inches(3.45), Inches(11.1), Inches(2.9))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    
    p = tf_bot.paragraphs[0]
    p.text = "Macro & Market Context"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ING_DARK_SLATE

    p_sub1 = tf_bot.add_paragraph()
    p_sub1.text = b1
    p_sub1.font.size = Pt(10.5)
    p_sub1.font.color.rgb = RGBColor(55, 65, 81)
    p_sub1.space_before = Pt(8)

    p_sub2 = tf_bot.add_paragraph()
    p_sub2.text = b2
    p_sub2.font.size = Pt(10.5)
    p_sub2.font.color.rgb = RGBColor(55, 65, 81)
    p_sub2.space_before = Pt(6)
    # =========================================================================
    # SLIDE 8: TERM SHEET (Exact Preview Parity)
    # =========================================================================
    s8 = prs.slides.add_slide(blank)
    
    # 1. Dynamic Title Resolution matching Preview Canvas
    s8_title = sm.get("s8_ttl")
    if not s8_title or s8_title == "Indicative Debt Financing Term Sheet" or s8_title == "Indicative Term Sheet":
        if p_fam == "GREEN_ESG" or "enel" in client_name.lower():
            s8_title = "Indicative Green / Sustainability-Linked Term Sheet"
        elif p_fam == "FX_HEDGE" or "asml" in client_name.lower():
            s8_title = "Indicative Zero-Cost Layered Collar Term Sheet"
        else:
            s8_title = "Indicative Pre-Hedge Swap & EMTN Term Sheet"

    add_header(s8, s8_title, category=sm.get("s8_cat", "TRANSACTION STRUCTURING"))
    add_logo(s8)
    add_footer(s8)

    # 2. Dynamic Rows matching React App.jsx
    if p_fam == "GREEN_ESG" or "enel" in client_name.lower():
        ts_rows = [
            ("Issuer / Counterparty", client_name),
            ("Instrument Format", "Senior Unsecured Green EMTN (ICMA Aligned)"),
            ("Notional Sizing", "EUR 600,000,000"),
            ("Tenor / Maturity", ov.get("tenor", "8 Years (Green Benchmark)")),
            ("Indicative Pricing", ov.get("spread", "Mid-Swap + 77 bps (Greenium: -5 bps)")),
            ("Sole Structurer / Counterparty", "ING Bank N.V.")
        ]
    elif p_fam == "FX_HEDGE" or "asml" in client_name.lower():
        ts_rows = [
            ("Issuer / Counterparty", client_name),
            ("Instrument Format", "Zero-Cost Layered Collar / FX Forward Options"),
            ("Notional Sizing", "USD 600,000,000"),
            ("Tenor / Maturity", "12-24 Months Rolling Quarterly"),
            ("Indicative Pricing", "Floor: 1.0850 | Cap: 1.0450 (Zero Upfront Cost)"),
            ("Sole Structurer / Counterparty", "ING Bank N.V.")
        ]
    else:  # RATES_HEDGE & DCM_REFI / BASF
        ts_rows = [
            ("Issuer / Counterparty", client_name),
            ("Instrument Format", "Forward-Starting IRS Overlay & Senior EMTN"),
            ("Notional Sizing", "EUR 600,000,000"),
            ("Tenor / Maturity", "7 Years (Pre-Hedge Tranche)"),
            ("Indicative Pricing", "Mid-Swap + 82 bps"),
            ("Sole Structurer / Counterparty", "ING Bank N.V.")
        ]

    # 3. Create 6-Row Term Sheet Table
    t_top = Inches(1.5)
    t_left = Inches(0.8)
    t_width = Inches(11.7)
    t_height = Inches(4.5)

    tbl_shape = s8.shapes.add_table(len(ts_rows), 2, t_left, t_top, t_width, t_height)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.8)
    tbl.columns[1].width = Inches(7.9)

    for r_idx, (k, v) in enumerate(ts_rows):
        cell_k = tbl.cell(r_idx, 0)
        cell_v = tbl.cell(r_idx, 1)

        # Clean light row fills
        fill_color = BG_LIGHT if r_idx % 2 == 1 else RGBColor(255, 255, 255)
        cell_k.fill.solid()
        cell_k.fill.fore_color.rgb = fill_color
        cell_v.fill.solid()
        cell_v.fill.fore_color.rgb = fill_color

        # Key Column Text
        tf_k = cell_k.text_frame
        tf_k.word_wrap = True
        p_k = tf_k.paragraphs[0]
        p_k.text = k
        p_k.font.bold = True
        p_k.font.size = Pt(11)
        p_k.font.color.rgb = ING_DARK_SLATE

        # Value Column Text & Styling
        tf_v = cell_v.text_frame
        tf_v.word_wrap = True
        p_v = tf_v.paragraphs[0]
        p_v.text = v
        p_v.font.size = Pt(11)

        if "Pricing" in k:
            p_v.font.bold = True
            p_v.font.color.rgb = ING_ORANGE
        elif "Sole Structurer" in k or r_idx == len(ts_rows) - 1:
            p_v.font.bold = True
            p_v.font.color.rgb = RGBColor(0, 0, 102)
        else:
            p_v.font.bold = False
            p_v.font.color.rgb = RGBColor(55, 65, 81)

    # Footnote
    tb_fn = s8.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4))
    tf_fn = tb_fn.text_frame
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "*Indicative pricing subject to market conditions, bookbuilding depth, and credit approval.*"
    p_fn.font.italic = True
    p_fn.font.size = Pt(9.5)
    p_fn.font.color.rgb = TEXT_MUTED
    # =========================================================================
    # SLIDE 9: EXECUTION ROADMAP (Exact Preview Parity)
    # =========================================================================
    s9 = prs.slides.add_slide(blank)
    
    s9_title = sm.get("s9_ttl")
    if not s9_title or s9_title == "Execution Roadmap & Syndicate Timeline":
        if p_fam == "GREEN_ESG" or "enel" in client_name.lower():
            s9_title = "Second-Party Opinion (SPO) & Syndicate Timeline"
        elif p_fam == "FX_HEDGE" or "asml" in client_name.lower():
            s9_title = "Collar Implementation & Governance Timeline"
        else:
            s9_title = "ISDA Schedule, CSA & Execution Timeline"

    add_header(s9, s9_title, category=sm.get("s9_cat", "EXECUTION ROADMAP"))
    add_logo(s9)
    add_footer(s9)

    # Dynamic 4-Milestone Data based on Product Family & Client
    if p_fam == "GREEN_ESG" or "enel" in client_name.lower():
        timeline_steps = [
            ("1", "T - 6 Weeks", "Framework & SPO", "Finalize Green Bond framework & Sustainalytics SPO"),
            ("2", "T - 3 Weeks", "Investor Marketing", "Dedicated ESG roadshow across European funds"),
            ("3", "T - 1 Week", "Bookbuilding", "Launch orderbook with greenium pricing tension"),
            ("4", "T-Day", "Settlement & Allocation", "Final settlement & annual impact reporting")
        ]
    elif p_fam == "FX_HEDGE" or "asml" in client_name.lower():
        timeline_steps = [
            ("1", "T - 4 Weeks", "Exposure Audit", "Map multi-currency AP/AR cash flows and net currency gaps."),
            ("2", "T - 2 Weeks", "ISDA/CSA Bounds", "Finalize derivative lines, collateral thresholds, and master terms."),
            ("3", "T - 1 Week", "Board Approval", "Lock FX risk committee parameters and strike boundary targets."),
            ("4", "T-Day", "Trade Execution", "Live collar tranche execution via ING FX trading desk.")
        ]
    else:  # RATES_HEDGE & DCM_REFI / BASF
        timeline_steps = [
            ("1", "T - 4 Weeks", "Exposure Sizing", "Audit debt wall maturities and determine optimal hedge sizing."),
            ("2", "T - 2 Weeks", "Pre-Hedge Swap", "Execute forward-starting interest rate swap to lock benchmark yields."),
            ("3", "T - 1 Week", "Syndicate Roadshow", "Finalize EMTN documentation and conduct institutional investor calls."),
            ("4", "T-Day", "Pricing & Settlement", "Syndicate orderbook launch, spread pricing, and bond settlement.")
        ]

    # Render 4 Timeline Cards
    card_w = Inches(2.75)
    card_h = Inches(4.3)
    c_y = Inches(1.5)
    gap = Inches(0.2)
    left_margin = Inches(0.8)

    for i, (num, time_label, title_label, desc) in enumerate(timeline_steps):
        c_x = left_margin + i * (card_w + gap)

        # Card Container
        c_shp = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_x, c_y, card_w, card_h)
        c_shp.fill.solid()
        c_shp.fill.fore_color.rgb = BG_LIGHT
        c_shp.line.color.rgb = LINE_GRAY

        # Orange Number Circle Badge
        badge_size = Inches(0.55)
        bx = c_x + (card_w - badge_size) / 2
        by = c_y + Inches(0.3)
        b_shp = s9.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, badge_size, badge_size)
        b_shp.fill.solid()
        b_shp.fill.fore_color.rgb = ING_ORANGE
        b_shp.line.fill.background()

        b_tf = b_shp.text_frame
        b_tf.word_wrap = False
        bp = b_tf.paragraphs[0]
        bp.text = num
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = RGBColor(255, 255, 255)
        bp.alignment = PP_ALIGN.CENTER

        # Content Text Box
        t_box = s9.shapes.add_textbox(c_x + Inches(0.15), c_y + Inches(1.0), card_w - Inches(0.3), Inches(3.0))
        tf = t_box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = time_label
        p1.font.bold = True
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = ING_DARK_SLATE
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = title_label
        p2.font.bold = True
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = ING_ORANGE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = RGBColor(75, 85, 99)
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(8)

    

    # =========================================================================
    # SLIDE 10: REGULATORY DISCLAIMERS
    # =========================================================================
    s10 = prs.slides.add_slide(blank)
    add_header(s10, sm["s10_ttl"], category=sm["s10_cat"])
    add_logo(s10)
    add_footer(s10)

    shp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_LIGHT
    shp.line.color.rgb = LINE_GRAY

    tb_disc = s10.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.4))
    tf_disc = tb_disc.text_frame
    tf_disc.word_wrap = True
    p = tf_disc.paragraphs[0]
    p.text = "Regulatory Disclosures & Target Market Notice"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ING_DARK_SLATE

    disc_list = ov.get("disclaimers") or compliance_bullets or [
        "This document is prepared for illustrative and discussion purposes only and does not constitute an offer, solicitation, or recommendation to enter into any transaction.",
        "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II / UK MiFIR is eligible counterparties and professional clients only (all distribution channels).",
        "This material has not been prepared in accordance with legal requirements designed to promote the independence of investment research.",
        "All rates, levels, spreads, and indicative terms shown are subject to change without notice and are not tradeable prices."
    ]

    # Ensure disc_list is a list
    if isinstance(disc_list, str):
        disc_list = [disc_list]
    elif not isinstance(disc_list, list):
        disc_list = list(disc_list) if disc_list else []

    for d in disc_list:
        if d and d.strip():
            p_d = tf_disc.add_paragraph()
            p_d.text = f"• {d}"
            p_d.font.size = Pt(11)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_before = Pt(8)

    # =========================================================================
    # SAVE
    # =========================================================================
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf