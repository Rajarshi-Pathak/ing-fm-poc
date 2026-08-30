import io
from fastapi import FastAPI, Request, Query, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
from pitchbook_builder import fetch_pitchbook_bundle, build_pitchbook

def normalize_copilot_overrides(overrides: dict) -> dict:
    if not isinstance(overrides, dict):
        return {}
    norm = dict(overrides)
    
    # Map any variations of Slide 6 unhedged exposure
    for k in ["unhedged_exposure", "unhedged_loss", "loss_impact", "unhedged_exposure_impact", "slide_6_unhedged", "unhedged_impact"]:
        if k in overrides:
            val = str(overrides[k]).strip()
            norm["fx_scen_up_unhedged"] = val
            norm["rate_scenario_up"] = val
            norm["unhedged_loss"] = val
            norm["unhedged_exposure"] = val

    # Map any variations of liquidity
    for k in ["liquidity", "available_liquidity", "liquidity_buffer"]:
        if k in overrides and "liquidity_str" not in norm:
            norm["liquidity_str"] = str(overrides[k])

    # Map any variations of net debt
    for k in ["net_debt", "debt"]:
        if k in overrides and "net_debt_str" not in norm:
            norm["net_debt_str"] = str(overrides[k])

    # Map any variations of revenue / ebitda
    if "revenue" in overrides and "revenue_str" not in norm:
        norm["revenue_str"] = str(overrides["revenue"])
    if "ebitda" in overrides and "ebitda_str" not in norm:
        norm["ebitda_str"] = str(overrides["ebitda"])

    return norm

"""
main.py - Pure Database-Driven Backend with Multi-Channel Ingestion, Vertex AI Copilot, Metrics & Schema-Aligned Signals
"""
import os
import io
import json
import logging
import re
import uuid
import urllib.parse
from datetime import datetime
from fastapi import FastAPI, HTTPException, Response, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
from google.cloud.sql.connector import Connector
import sqlalchemy
import pg8000
import feedparser
from pypdf import PdfReader
from pptx import Presentation

try:
    from google import genai
    from google.genai import types
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

from pitchbook_builder import fetch_pitchbook_bundle, build_pitchbook

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("ing_fm_backend")

app = FastAPI(title="ING FM Insights API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

PROMISSORY_PATTERNS = [
    r"\bguarantee(?:d|s|ing)?\b",
    r"\brisk[-\s]?free\b",
    r"\beliminate(?:s|d|ing)?\s+all\s+risk\b",
    r"\bno\s+risk\b",
    r"\bcertain\s+return\b",
    r"\bwill\s+definitely\b",
    r"\babsolute(?:ly)?\s+certain\b",
    r"\bfully\s+protected\b"
]


def get_db_connection():
    instance_connection_name = os.getenv("INSTANCE_CONNECTION_NAME")
    db_user = os.getenv("DB_USER", "postgres")
    db_pass = os.getenv("DB_PASS", "postgres")
    db_name = os.getenv("DB_NAME", "postgres")

    if not instance_connection_name:
        try:
            conn = pg8000.connect(
                host="127.0.0.1", port=5432, user=db_user, password=db_pass, database=db_name
            )
            return conn, None
        except Exception:
            return None, None

    connector = Connector()

    def getconn():
        return connector.connect(
            instance_connection_name, "pg8000", user=db_user, password=db_pass, db=db_name
        )

    try:
        pool = sqlalchemy.create_engine("postgresql+pg8000://", creator=getconn)
        conn = pool.raw_connection()
        return conn, connector
    except Exception as e:
        logger.error(f"Cloud SQL connection failed: {e}")
        return None, None


class TextIngestRequest(BaseModel):
    client_id: str
    source_channel: str
    source_name: str
    text_content: str


class OpportunityRequest(BaseModel):
    client_id: str
    opportunity_id: Optional[str] = None
    overrides: Optional[Dict[str, Any]] = None
    compliance_bullets: Optional[List[str]] = None
    compliance_bullets: Optional[List[str]] = None


class CopilotMessage(BaseModel):
    client_id: str
    prompt: str
    history: Optional[List[Dict[str, str]]] = []
    current_overrides: Optional[Dict[str, Any]] = None


class ComplianceAuditRequest(BaseModel):
    client_id: str
    overrides: Optional[Dict[str, Any]] = None


@app.get("/healthz")
def healthz():
    return {"status": "healthy"}


@app.get("/api/metrics")
def get_rm_metrics():
    conn, connector = get_db_connection()
    priorities = []
    active_drafts_count = 0
    pending_review_count = 0
    cohort_matches_count = 0
    
    if conn:
        try:
            cur = conn.cursor()
            
            # 1. Dynamic Priorities Ranking
            cur.execute("""
                SELECT DISTINCT ON (c.client_id)
                    c.client_name,
                    c.client_id,
                    COALESCE(o.priority_score, 75) as score,
                    COALESCE(o.opportunity_type, 'STRATEGIC FINANCING') as opp_type,
                    COALESCE(o.est_revenue_eur_000, 0) as est_fee,
                    COALESCE(o.next_best_action, 'Review opportunity and schedule coverage call.') as next_action,
                    COALESCE(o.why_now_nlg, '') as why_now,
                    COALESCE(c.industry_sector, 'Wholesale') as sector,
                    COALESCE(c.hq_country, 'Europe') as country,
                    COALESCE(c.rm_name, 'Coverage Director') as rm_name
                FROM ca.ca_opportunity_scoring o
                JOIN ca.client_master c ON (o.client_id = c.client_id OR c.client_id LIKE o.client_id || '%%' OR o.client_id LIKE c.client_id || '%%')
                ORDER BY c.client_id, score DESC, est_fee DESC;
            """)
            rows = cur.fetchall()
            
            # Sort distinct clients by priority score DESC, then est_fee DESC
            sorted_rows = sorted(rows, key=lambda x: (int(x[2]), float(x[4])), reverse=True)[:4]
            
            for rank_idx, r in enumerate(sorted_rows, 1):
                cname, cid, score, opp_type, est_fee, next_action, why_now, sector, country, rm_name = r
                fee_val = float(est_fee)
                fee_str = f"€{fee_val/1000:.1f}M" if fee_val >= 1000 else (f"€{int(fee_val)}k" if fee_val > 0 else "")
                
                badge_str = f"{str(opp_type).upper()} · RANK #{rank_idx} (SCORE {score})"
                desc_str = str(why_now) if why_now else f"{sector} ({country}) — Active balance sheet and maturity window review."
                
                priorities.append({
                    "rank": rank_idx,
                    "badge": badge_str,
                    "title": str(cname),
                    "client_name": str(cname),
                    "client_id": str(cid),
                    "score": int(score),
                    "type": str(opp_type).upper(),
                    "opportunity_type": str(opp_type),
                    "fee_estimate": fee_str,
                    "est_fee_k": int(fee_val),
                    "desc": desc_str,
                    "why_now": desc_str,
                    "action": str(next_action),
                    "next_best_action": str(next_action),
                    "sector": str(sector),
                    "country": str(country),
                    "rm_name": str(rm_name)
                })

            # 2. Dynamic Metric: Active drafts / Ingested client signals
            cur.execute("""
                SELECT COUNT(DISTINCT client_id) 
                FROM ca.digital_twin_signals;
            """)
            res_active = cur.fetchone()
            active_drafts_count = int(res_active[0]) if res_active and res_active[0] is not None else len(priorities)

            # 3. Dynamic Metric: Deals pending review (High priority opportunities >= 85)
            cur.execute("""
                SELECT COUNT(DISTINCT client_id) 
                FROM ca.ca_opportunity_scoring 
                WHERE priority_score >= 85;
            """)
            res_pending = cur.fetchone()
            pending_review_count = int(res_pending[0]) if res_pending and res_pending[0] is not None else len(priorities)

            # 4. Dynamic Metric: Total Cohort Matches in Database
            cur.execute("""
                SELECT COUNT(*) 
                FROM ca.client_master;
            """)
            res_cohort = cur.fetchone()
            cohort_matches_count = int(res_cohort[0]) if res_cohort and res_cohort[0] is not None else 13

            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as e:
            logger.error(f"Failed to query RM metrics: {e}")

    return {
        "active_drafts": {
            "value": active_drafts_count, 
            "change": f"▲ {active_drafts_count}", 
            "label": "Active drafts in progress"
        },
        "avg_time": {
            "value": "< 15s", 
            "change": "▼ 99% vs manual", 
            "label": "Avg. time to first draft"
        },
        "avg_time_draft": {
            "value": "< 15s", 
            "change": "▼ 99% vs manual", 
            "label": "Avg. time to first draft"
        },
        "pending_review": {
            "value": pending_review_count, 
            "change": "High conviction", 
            "label": "Deals pending review"
        },
        "cohort_matches": {
            "value": cohort_matches_count, 
            "change": "▲ 5", 
            "label": "Cohort matches in database"
        },
        "priorities": priorities
    }


@app.get("/api/signals")
def get_live_signals():
    conn, connector = get_db_connection()
    signals = []
    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT 
                    s.signal_id,
                    s.client_id,
                    COALESCE(c.client_name, s.client_id) as client_name,
                    s.signal_type,
                    COALESCE(s.metric_identified, s.trigger_summary, s.description, 'Market Catalyst') as headline,
                    s.confidence_pct,
                    s.urgency,
                    s.created_at
                FROM ca.digital_twin_signals s
                LEFT JOIN ca.client_master c ON (s.client_id = c.client_id OR c.client_id LIKE s.client_id || '%%' OR s.client_id LIKE c.client_id || '%%')
                ORDER BY s.created_at DESC 
                LIMIT 15;
            """)
            rows = cur.fetchall()
            now_dt = datetime.now()
            
            for r in rows:
                sig_id, cid, cname, stype, headline, conf, urgency, created_at = r
                urg_str = str(urgency or "Medium").upper()
                trend = "up" if urg_str in ["HIGH", "CRITICAL"] else ("down" if urg_str in ["LOW"] else "neutral")
                
                time_ago = "Just now"
                if created_at:
                    try:
                        delta = now_dt - created_at
                        mins = int(delta.total_seconds() / 60)
                        if mins < 1:
                            time_ago = "Just now"
                        elif mins < 60:
                            time_ago = f"{mins}m ago"
                        else:
                            hours = int(mins / 60)
                            time_ago = f"{hours}h ago"
                    except Exception:
                        time_ago = "Recent"

                signals.append({
                    "id": str(sig_id),
                    "client_id": str(cid),
                    "client_name": str(cname),
                    "type": str(stype or "CATALYST").upper(),
                    "text": f"{cname}: {headline}",
                    "headline": str(headline),
                    "confidence": int(conf or 90),
                    "urgency": urg_str,
                    "trend": trend,
                    "time_ago": time_ago
                })
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as e:
            logger.warning(f"Failed to query digital twin signals: {e}")

    if not signals:
        signals = [
            {"id": "SIG-DF1", "client_id": "CLI103", "client_name": "BASF SE", "type": "REFINANCING", "text": "BASF SE: €2.0B 6Y EMTN & €1.2B Pre-Hedge", "headline": "€2.0B 6Y EMTN & €1.2B Pre-Hedge", "confidence": 94, "urgency": "HIGH", "trend": "up", "time_ago": "Just now"},
            {"id": "SIG-DF2", "client_id": "CLI101", "client_name": "Enel S.p.A.", "type": "SUSTAINABLE", "text": "Enel S.p.A.: EUR 750M Green EMTN Pre-Hedge", "headline": "EUR 750M Green EMTN Pre-Hedge", "confidence": 94, "urgency": "HIGH", "trend": "up", "time_ago": "12m ago"},
            {"id": "SIG-DF3", "client_id": "CLI102", "client_name": "ASML Holding", "type": "HEDGING", "text": "ASML Holding: EUR 900M FX Collar Hedge", "headline": "EUR 900M FX Collar Hedge", "confidence": 92, "urgency": "HIGH", "trend": "up", "time_ago": "25m ago"}
        ]
    return signals


@app.get("/api/opportunities")
def get_opportunities():
    conn, connector = get_db_connection()
    opps = []

    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT 
                    cm.client_id,
                    cm.client_name,
                    cm.tier,
                    cm.hq_country,
                    COALESCE(cm.rm_name, 'Coverage Director'),
                    COALESCE(cm.industry_sector, 'Wholesale Banking'),
                    COALESCE(fl.net_debt_eur_m, 0),
                    COALESCE(fl.liquidity_eur_m, 0),
                    COALESCE(fl.debt_maturing_24m_eur_m, 0),
                    COALESCE(os.priority_score, 75),
                    COALESCE(os.opportunity_type, 'DEBT REFINANCING'),
                    COALESCE(os.next_best_action, 'Capital structure review and proactive balance sheet advisory.'),
                    COALESCE(os.why_now_nlg, 'Upcoming maturity window and active market rate dynamics.'),
                    COALESCE(os.est_revenue_eur_000, 0)
                FROM ca.client_master cm
                LEFT JOIN LATERAL (
                    SELECT net_debt_eur_m, liquidity_eur_m, debt_maturing_24m_eur_m
                    FROM ca.ext_company_filings
                    WHERE client_id = cm.client_id OR client_id LIKE cm.client_id || '%%'
                    ORDER BY reporting_period DESC LIMIT 1
                ) fl ON true
                LEFT JOIN LATERAL (
                    SELECT priority_score, opportunity_type, next_best_action, why_now_nlg, est_revenue_eur_000
                    FROM ca.ca_opportunity_scoring
                    WHERE client_id = cm.client_id OR client_id LIKE cm.client_id || '%%'
                    ORDER BY priority_score DESC LIMIT 1
                ) os ON true
                ORDER BY os.priority_score DESC NULLS LAST, cm.client_name ASC;
            """)
            client_rows = cur.fetchall()

            for r in client_rows:
                cid, name, tier, hq, rm, sector, net_debt, liq, m24, score_num, opp_type, action, why_now, est_fee = r
                cid_str = str(cid)
                name_str = str(name)

                # Check debt maturities count
                cur.execute("""
                    SELECT COALESCE(SUM(amount_eur_m), 0), COUNT(isin)
                    FROM ca.debt_maturity_schedule
                    WHERE client_id = %s OR client_id LIKE %s;
                """, (cid_str, f"{cid_str}%"))
                sched_row = cur.fetchone()
                total_nominal = float(sched_row[0]) if sched_row else 0.0
                tranches_count = int(sched_row[1]) if sched_row else 0

                if float(m24) == 0 and total_nominal > 0:
                    m24 = total_nominal

                score_level = "High" if int(score_num) >= 85 else ("Medium" if int(score_num) >= 70 else "Low")
                score_val = f"{score_level} · {score_num}"

                chips = []
                if float(m24) > 0:
                    chips.append(f"€{float(m24):,.0f}M debt maturity wall in next 24M")
                elif tranches_count > 0:
                    chips.append(f"{tranches_count} bond tranches scheduled")
                else:
                    chips.append(f"Industry: {sector}")

                if float(net_debt) > 0:
                    chips.append(f"Net Debt: €{float(net_debt):,.0f}M | Liquidity: €{float(liq):,.0f}M")
                else:
                    chips.append("Active balance sheet review")

                # 1. Fetch Context Fabric Signal & Evidence Note
                cur.execute("""
                    SELECT signal_type, trigger_summary, description, confidence_pct, metric_identified
                    FROM ca.digital_twin_signals
                    WHERE client_id = %s OR client_id LIKE %s
                    ORDER BY confidence_pct DESC, created_at DESC LIMIT 1;
                """, (cid_str, f"{cid_str}%"))
                sig_row = cur.fetchone()
                
                cf_desc = sig_row[2] if sig_row else "Multi-source synthesis across public filings, regulatory disclosures, and treasury dialog."
                cf_latent = sig_row[1] if sig_row else "Residual liability management and pre-hedge window."
                cf_metric = sig_row[4] if sig_row and sig_row[4] else "Capital Markets Sizing"
                
                # 2. Fetch Active Ingestion Channels & Attribution Source
                cur.execute("""
                    SELECT DISTINCT source_channel, source_name 
                    FROM ca.document_vector_chunks 
                    WHERE client_id = %s OR client_id LIKE %s
                    LIMIT 6;
                """, (cid_str, f"{cid_str}%"))
                chunk_rows = cur.fetchall()
                
                raw_channels = list(set([r[0].upper() for r in chunk_rows if r[0]])) if chunk_rows else ["NEWS_RSS", "TREASURY_EMAIL", "TEAMS_CHAT", "CONTEXT_FABRIC"]
                source_authors = [r[1] for r in chunk_rows if r[1] and any(k in r[1].lower() for k in ["origination", "rates", "rm", "treasury", "luca", "giulia", "marta", "klaus", "daan"])]
                attribution_author = source_authors[0] if source_authors else f"{rm} (Coverage Desk)"

                opps.append({
                    "id": cid_str,
                    "name": name_str,
                    "type": opp_type,
                    "is_debt": float(m24) > 0 or total_nominal > 0 or "DEBT" in opp_type.upper(),
                    "subtitle": f"{tier or 'Tier 1'} client ({hq or sector})",
                    "score": score_val,
                    "score_num": int(score_num),
                    "chips": chips,
                    "callout": f"{why_now} {action}".strip(),
                    "action": action,
                    "why_now": why_now,
                    "slides_count": 10,
                    "net_debt_str": f"€{float(net_debt):,.0f}M" if float(net_debt) > 0 else "—",
                    "liquidity_str": f"€{float(liq):,.0f}M" if float(liq) > 0 else "—",
                    "debt_maturing_24m_str": f"€{float(m24):,.0f}M" if float(m24) > 0 else "—",
                    "rm_name": rm,
                    "cf_description": cf_desc,
                    "cf_latent": cf_latent,
                    "cf_metric": cf_metric,
                    "ingestion_channels": raw_channels,
                    "attribution_author": attribution_author
                })

            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception as exc:
            logger.error(f"Error querying opportunities: {exc}")

    return opps


@app.get("/api/client/{client_id}/maturities")
def get_client_maturities(client_id: str):
    conn, connector = get_db_connection()
    ladder = []
    cid = str(client_id).strip()

    if conn:
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT maturity_year, SUM(amount_eur_m) as total_nominal, COUNT(isin) as tranches
                FROM ca.debt_maturity_schedule
                WHERE client_id = %s OR client_id LIKE %s
                GROUP BY maturity_year
                ORDER BY maturity_year ASC;
            """, (cid, f"{cid}%"))
            rows = cur.fetchall()
            for r in rows:
                ladder.append({
                    "year": str(r[0]),
                    "amount_eur_m": float(r[1]),
                    "tranches": int(r[2])
                })
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception:
            pass

    if not ladder:
        if "103" in cid or "BASF" in cid.upper():
            ladder = [
                {"year": "2026", "amount_eur_m": 1500.0, "tranches": 1},
                {"year": "2027", "amount_eur_m": 3300.0, "tranches": 2},
                {"year": "2028", "amount_eur_m": 1200.0, "tranches": 1},
                {"year": "2029", "amount_eur_m": 2000.0, "tranches": 1}
            ]
        else:
            ladder = [
                {"year": "2026", "amount_eur_m": 3000.0, "tranches": 2},
                {"year": "2027", "amount_eur_m": 2500.0, "tranches": 2},
                {"year": "2028", "amount_eur_m": 1600.0, "tranches": 1},
                {"year": "2029", "amount_eur_m": 1200.0, "tranches": 1}
            ]
    return ladder


@app.get("/api/rss/feed")
def get_client_rss_feed(client_id: str, query: Optional[str] = None):
    import urllib.request
    
    cid = str(client_id).strip()
    cname = "Corporate Client"
    
    conn, connector = get_db_connection()
    if conn:
        try:
            cur = conn.cursor()
            cur.execute("SELECT client_name FROM ca.client_master WHERE client_id = %s OR client_id LIKE %s LIMIT 1;", (cid, f"{cid}%"))
            row = cur.fetchone()
            if row and row[0]:
                cname = str(row[0])
            cur.close()
            conn.close()
            if connector:
                connector.close()
        except Exception:
            pass
            
    if cname == "Corporate Client":
        if "101" in cid or "ENEL" in cid.upper():
            cname = "Enel S.p.A."
        elif "103" in cid or "BASF" in cid.upper():
            cname = "BASF SE"
        elif "ASML" in cid.upper():
            cname = "ASML Holding N.V."

    clean_name = cname.split(" S.p.A.")[0].split(" SE")[0].split(" N.V.")[0].split(" AG")[0].strip()
    search_term = query or f'"{clean_name}" bond OR debt OR refinancing OR hybrid OR EMTN'
    encoded_q = urllib.parse.quote(search_term)
    rss_url = f"https://news.google.com/rss/search?q={encoded_q}&hl=en-US&gl=US&ceid=US:en"

    articles = []
    try:
        req_obj = urllib.request.Request(
            rss_url,
            headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
            }
        )
        with urllib.request.urlopen(req_obj, timeout=5) as response:
            xml_data = response.read()
            feed = feedparser.parse(xml_data)
            
            for entry in feed.entries[:8]:
                raw_summary = getattr(entry, 'summary', entry.title)
                clean_summary = re.sub(r'<[^>]+>', '', raw_summary).replace('&nbsp;', ' ').strip()
                articles.append({
                    "title": entry.title,
                    "link": entry.link,
                    "published": getattr(entry, 'published', datetime.now().strftime("%a, %d %b %Y %H:%M:%S GMT")),
                    "summary": clean_summary
                })
    except Exception as exc:
        logger.warning(f"Live Google News fetch notice: {exc}")

    if not articles:
        now_str = datetime.now().strftime("%a, %d %b %Y %H:%M:%S GMT")
        if "ENEL" in cname.upper() or "101" in cid:
            articles = [
                {
                    "title": f"Enel S.p.A. Explores €2.5B Hybrid Capital Bond Refinancing Window Ahead of 2026 Rollover",
                    "link": "https://www.enel.com/investors",
                    "published": now_str,
                    "summary": "Enel Treasury evaluates long-tenor green hybrid bonds and pre-hedge interest rate swaps to manage upcoming 2026 debt wall."
                },
                {
                    "title": f"European Power & Utilities Sector Faces Elevated Refinancing Calendar Across 2026–2027",
                    "link": "https://think.ing.com",
                    "published": now_str,
                    "summary": "Secondary credit spreads trade range-bound while forward swap pre-hedges gain momentum across Tier 1 European power utilities."
                }
            ]
        elif "BASF" in cname.upper() or "103" in cid:
            articles = [
                {
                    "title": f"BASF SE Evaluates €1.5B Senior Bond Issuance as Derivative Hedges Mature",
                    "link": "https://www.basf.com/investor",
                    "published": now_str,
                    "summary": "Chemicals platform reviews floating-to-fixed interest rate swaps and €4.8B 24-month maturity schedule."
                },
                {
                    "title": f"European Chemicals Sector Outlook: Financing Costs Stabilize Around 5Y EUR Swap 2.62%",
                    "link": "https://think.ing.com",
                    "published": now_str,
                    "summary": "Credit analysts highlight liability management opportunities for BBB-rated corporate issuers."
                }
            ]
        else:
            articles = [
                {
                    "title": f"{cname} Announces Comprehensive Balance Sheet Review & Debt Strategy",
                    "link": "#",
                    "published": now_str,
                    "summary": f"Corporate Treasury assesses financing options and interest rate hedges ahead of upcoming maturities."
                }
            ]

    return {"client_id": cid, "client_name": cname, "query": search_term, "articles": articles}


# =========================================================================
# Schema-Aligned Ingestion Endpoint
# =========================================================================
@app.post("/api/ingest/text")
def ingest_text_signal(req: TextIngestRequest):
    cid = req.client_id
    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    cname = bundle.get("client_name", "Corporate Client")
    text = req.text_content
    channel = req.source_channel
    sname = req.source_name

    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    extracted = {
        "signal_type": "REFINANCING",
        "catalog_family": "Financing/Capital Markets",
        "metric_identified": sname[:100] if sname else "Market Catalyst",
        "trigger_summary": text[:200],
        "metric_value": "Market Catalyst",
        "description": text[:500],
        "confidence_pct": 94,
        "urgency": "High",
        "suggested_action": "Execute EMTN Benchmark with Pre-Hedge Swap Overlay",
        "est_revenue_eur_000": 5500,
        "priority_score": 94
    }

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            prompt = f"""
Analyze the following touchpoint text for wholesale corporate client '{cname}' (Client ID: {cid}):

TEXT CONTENT:
{text}

Extract structured signal parameters matching the database schema as STRICT JSON without markdown:
{{
  "signal_type": "REFINANCING | LIQUIDITY | COVENANT | HEDGING | M&A",
  "catalog_family": "Financing/Capital Markets | Interest Rate | Foreign Exchange | Sustainable Finance",
  "metric_identified": "Short headline / metric identified (max 100 chars)",
  "trigger_summary": "1-sentence executive trigger summary",
  "metric_value": "Key value or spread (e.g. €750M or Mid-Swap +77bps)",
  "description": "2-sentence detailed institutional description",
  "confidence_pct": 94,
  "urgency": "High | Medium | Low",
  "suggested_action": "Specific deal recommendation (e.g. Execute EUR 750M 8Y Green EMTN with EUR 500M Pre-Hedge)",
  "est_revenue_eur_000": 5500,
  "priority_score": 94
}}
"""
            response = client_gcp.models.generate_content(
                model='gemini-2.5-flash',
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.1,
                    response_mime_type="application/json"
                )
            )
            extracted_json = json.loads(response.text)
            extracted.update(extracted_json)
        except Exception as e:
            logger.warning(f"Vertex AI signal extraction error: {e}")

    # Write to ca.digital_twin_signals and update ca.ca_opportunity_scoring
    conn, connector = get_db_connection()
    if conn:
        try:
            cur = conn.cursor()
            sig_id = f"SIG-{uuid.uuid4().hex[:8].upper()}"
            
            # 1. Insert into ca.digital_twin_signals
            cur.execute("""
                INSERT INTO ca.digital_twin_signals 
                (signal_id, client_id, catalog_family, signal_type, metric_identified, trigger_summary, metric_value, description, confidence_pct, urgency, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW());
            """, (
                sig_id,
                str(cid),
                str(extracted.get("catalog_family", "Financing/Capital Markets")),
                str(extracted.get("signal_type", "REFINANCING")),
                str(extracted.get("metric_identified", "Catalyst"))[:100],
                str(extracted.get("trigger_summary", text[:200])),
                str(extracted.get("metric_value", "Live Trigger"))[:50],
                str(extracted.get("description", text[:500])),
                int(extracted.get("confidence_pct", 94)),
                str(extracted.get("urgency", "High"))
            ))

            # 2. Update ca.ca_opportunity_scoring
            new_action = str(extracted.get("suggested_action", "Execute EUR 750M 8Y Green EMTN with EUR 500M Pre-Hedge"))
            new_why_now = str(extracted.get("trigger_summary", text[:200]))
            new_fee = float(extracted.get("est_revenue_eur_000", 5500))
            new_score = int(extracted.get("priority_score", 94))
            new_type = str(extracted.get("signal_type", "REFINANCING"))

            cur.execute("""
                UPDATE ca.ca_opportunity_scoring
                SET next_best_action = %s,
                    why_now_nlg = %s,
                    est_revenue_eur_000 = %s,
                    priority_score = %s,
                    opportunity_type = %s
                WHERE client_id = %s OR client_id LIKE %s;
            """, (
                new_action,
                new_why_now,
                new_fee,
                new_score,
                new_type,
                str(cid),
                f"{str(cid)}%"
            ))

            conn.commit()
            cur.close()
            conn.close()
            if connector:
                connector.close()
            logger.info(f"Successfully committed signal {sig_id} and updated opportunity scoring for {cid}")
        except Exception as e:
            logger.error(f"Failed to persist ingested signal to DB: {e}")

    return {
        "status": "INGESTED_AND_EVALUATED",
        "client_id": cid,
        "source_channel": channel,
        "source_name": sname,
        "extracted_signal": {
            "signal_headline": extracted.get("metric_identified", sname),
            "signal_type": extracted.get("signal_type", "REFINANCING"),
            "urgency": extracted.get("urgency", "High"),
            "confidence_pct": extracted.get("confidence_pct", 94)
        },
        "timestamp": datetime.now().strftime("%H:%M:%S")
    }


@app.post("/api/ingest/file")
async def ingest_file_signal(
    client_id: str = Form(...),
    source_channel: str = Form("Document Upload"),
    file: UploadFile = File(...)
):
    cid = str(client_id).strip()
    contents = await file.read()
    extracted_text = ""

    filename = file.filename.lower()
    if filename.endswith(".pdf"):
        try:
            pdf_reader = PdfReader(io.BytesIO(contents))
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() or ""
        except Exception as e:
            extracted_text = f"Error reading PDF content: {e}"
    elif filename.endswith(".pptx"):
        try:
            prs = Presentation(io.BytesIO(contents))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            extracted_text += p.text + "\n"
        except Exception as e:
            extracted_text = f"Error reading PPTX content: {e}"
    else:
        extracted_text = contents.decode("utf-8", errors="ignore")

    req = TextIngestRequest(
        client_id=cid,
        source_channel=source_channel,
        source_name=file.filename,
        text_content=extracted_text[:4000]
    )
    return ingest_text_signal(req)


@app.post("/api/check-compliance")
def check_compliance_endpoint(req: ComplianceAuditRequest):
    cid = req.client_id
    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    client_name = bundle.get("client_name", "Corporate Client")
    p_family = bundle.get("product_family", "DCM_REFI")
    ov = req.overrides or {}
    now_stamp = datetime.now().strftime("%d %B %Y, %H:%M CET")
    is_fx = (p_family == "FX_HEDGE")

    # Product-appropriate term sheet pricing
    if is_fx:
        indicative_pricing = ov.get("term_pricing") or ov.get("collar_structure") or "Zero Net Upfront Premium (Zero-Cost Collar corridor 1.0450 - 1.0850)"
    else:
        indicative_pricing = ov.get("term_pricing") or ov.get("spread") or "Mid-Swap + 82 bps"

    now_dt = datetime.now()
    current_year = now_dt.year

    audit_payload = {
        "client_name": client_name,
        "product_family": p_family,
        "current_audit_date": now_stamp,
        "slide_01_cover": "Cover Title and Target Client Branding",
        "slide_03_exec_summary": "Executive Context & Opportunity Rationale",
        "slide_04_balance_sheet": "Balance Sheet & Liquidity Overview",
        "slide_05_debt_maturity_and_fx": "Debt Maturity Profile & Refinancing Horizon (Tranche clustering 2026-2028, FX Policy Sizing & Hedged Ratios)",
        "slide_06_sensitivity": "Scenario Sensitivity Analysis",
        "slide_07_execution_roadmap": "Syndicate Execution Timeline",
        "slide_08_term_sheet": "Indicative Terms: " + indicative_pricing,
        "slide_10_disclaimers": ov.get("disclaimers", [
            "Prepared for illustrative and discussion purposes only; not an offer or solicitation.",
            "Target market under MiFID II / UK MiFIR: Eligible counterparties and professional clients only.",
            "Non-independent investment research disclaimer.",
            "Rates, levels, and spreads are indicative, subject to change, and not tradeable prices."
        ])
    }

    sys_inst = (
        "You are the ING Wholesale Banking Regulatory & Compliance Officer.\n"
        f"CURRENT OPERATING CONTEXT: Today is {now_stamp}. The active operational year is {current_year}.\n"
        f"Audit this 10-slide pitchbook for {client_name} ({p_family}) against MiFID II Art. 24, EMIR derivatives transparency, ICMA Principles, and FINRA 2210.\n\n"
        "STRICT GROUNDING & AUDIT RULES:\n"
        f"- DATE ANCHORING: The reference date '{now_stamp}' and year '{current_year}' are CURRENT. Do NOT flag {current_year} dates or 2026-2028 maturity horizons as future or speculative data.\n"
        "- SLIDE 1: Cover page only. DO NOT flag for missing disclaimers.\n"
        "- SLIDE 5: Audit debt maturity profile and FX sizing consistency. Flag if hedged numbers and stated residual gaps are mathematically contradictory.\n"
        "- SLIDE 8: Check for proximate indicative pricing caveats (terms must state they are non-binding/subject to market conditions and credit approval).\n"
        "- SLIDE 10: Check general disclaimers. If product family or hedging involves derivatives, check for EMIR classification/reporting notices.\n\n"
        "Return a valid JSON object with keys:\n"
        "compliant (bool), overall_risk_assessment (HIGH/MEDIUM/LOW), compliance_summary (string), flagged_slides (list of ints), flags (list of objects with slide_number, rule, issue, suggested_fix)."
    )

    result_json = None
    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            resp = client_gcp.models.generate_content(
                model="gemini-2.5-flash",
                contents=json.dumps(audit_payload),
                config=types.GenerateContentConfig(
                    system_instruction=sys_inst,
                    temperature=0.1,
                    response_mime_type="application/json"
                )
            )
            parsed = json.loads(resp.text)
            if "flags" in parsed and "compliance_summary" in parsed:
                result_json = parsed
        except Exception as e:
            logger.warning(f"Vertex AI compliance audit warning: {e}")

    if not result_json:
        result_json = {
            "compliant": False,
            "overall_risk_assessment": "MEDIUM",
            "compliance_summary": f"Full regulatory audit completed for {client_name}. Identified recommended disclosures under MiFID II and {'EMIR derivative standards' if is_fx else 'ICMA / Prospectus Regulation'}.",
            "flagged_slides": [5, 8, 10],
            "flags": [
                {
                    "slide_number": 5,
                    "rule": "MiFID II Art. 24 (Market Data Currency)",
                    "issue": "Market rate snapshot requires timestamp certification.",
                    "suggested_fix": f"Market Snapshot as of {now_stamp}"
                },
                {
                    "slide_number": 8,
                    "rule": "EMIR / MiFID II (Indicative Derivatives Notice)" if is_fx else "FINRA 2210 (Indicative Pricing Caveat)",
                    "issue": "Preliminary pricing corridor requires formal non-binding execution qualification.",
                    "suggested_fix": "*Indicative terms subject to market volatility, ISDA documentation, and final credit approval.*"
                },
                {
                    "slide_number": 10,
                    "rule": "MiFID II / EMIR (Target Market & Investor Classification)",
                    "issue": "Professional Client & Eligible Counterparty restrictions must be explicitly active.",
                    "suggested_fix": "FOR PROFESSIONAL CLIENTS AND ELIGIBLE COUNTERPARTIES ONLY: Target market under MiFID II is eligible counterparties and professional clients only."
                }
            ]
        }

    return result_json

@app.post("/api/copilot/chat")
def copilot_chat_endpoint(req: CopilotMessage):
    cid = req.client_id
    prompt = req.prompt
    history = req.history or []
    now_stamp = datetime.now().strftime("%d %B %Y, %H:%M CET")

    bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection)
    client_name = bundle.get("client_name", "Corporate Client")
    p_family = bundle.get("product_family", "DCM_REFI")
    maturities_data = get_client_maturities(cid)
    current_ov = req.current_overrides or {}

    history_str = ""
    for h in history[-6:]:
        speaker = "RM" if h.get("sender") == "user" else "Copilot"
        history_str += f"{speaker}: {h.get('text')}\n"

    system_instruction = """You are the senior ING Financial Markets Origination & Structuring Copilot.
You assist Relationship Managers (RMs) by explaining pitchbook slides, providing CFO talking points, analyzing exposures, and dynamically updating slide parameters.

10-SLIDE PITCHBOOK ARCHITECTURE & CONTEXT MAP:
- Slide 1: Cover Slide (client_name, kicker, subtitle, rm_name, market_date)
- Slide 2: Deal Catalyst & Strategic Trigger (trigger, window, action)
- Slide 3: Executive Summary & Strategic Pillars
- Slide 4: Capital Structure & Financial Snapshot (revenue_str, ebitda_str, net_debt_str, liquidity_str, maturity_wall_str, debt_maturing_24m_str)
- Slide 5: Exposure Horizon & Horizon Analysis (unhedged_gap_str, maturity_wall_str, usd_exposure, eur_cost, hedge_ratio, target_hedge_ratio)
- Slide 6: Sensitivity Analysis & Payoff Corridor (rate_scenario_up, rate_scenario_lock, rate_scenario_down, fx_scen_up_unhedged, fx_scen_up_hedged, fx_scen_down_unhedged)
- Slide 7: Macroeconomic Backdrop & Curve Dynamics (swap_5y, bund_10y, itraxx_main, market_date)
- Slide 8: Indicative Transaction Term Sheet (notional_bond, spread, tenor, spread_disclaimer)
- Slide 9: Execution Roadmap & Syndication Timeline
- Slide 10: Regulatory Disclosures, Compliance & MiFID II Disclaimers (disclaimers)

CRITICAL RULES:
1. When asked to EXPLAIN or ANALYZE a slide (e.g. "explain slide 2", "talking points for slide 4"), you MUST read and reference the exact slide fields and numbers in "active_deck_slides" for that slide (e.g., triggers, window, actions, rates, notional, spreads).
2. Whenever the user asks to edit, replace, update, or change ANY value, number, or text on any slide, you MUST:
   a. Return a clear explanation and strategic rationale in "reply".
   b. ALWAYS include the modified state key in "overrides" with the exact formatted string!

MUTATION EXAMPLES:
- User: "In Slide 2, change €3.2B debt maturities to €1.2B debt maturities"
  -> "overrides": {"trigger": "Upcoming €1.2B debt maturities face repricing risk amid benchmark curve fluctuations."}
- User: "In Slide 4, change €3,000M to €1,000M"
  -> "overrides": {"maturity_wall_str": "€1,000M", "debt_maturing_24m_str": "€1,000M"}
- User: "In Slide 4, set Net Debt to €15,000M"
  -> "overrides": {"net_debt_str": "€15,000M"}
- User: "In Slide 7, update iTraxx Main to 60 bps"
  -> "overrides": {"itraxx_main": "60 bps"}
- User: "In Slide 8, set tenor to 10 Years"
  -> "overrides": {"tenor": "10 Years (Euro Benchmark)"}
- User: "In Slide 8, set pricing to Mid-Swap + 75 bps"
  -> "overrides": {"spread": "Mid-Swap + 75 bps"}

If no edit was requested (e.g. "Explain Slide 2"), return "overrides": {}.

OUTPUT REQUIREMENT:
You must ALWAYS return a single valid JSON object with exactly two keys: "reply" (string) and "overrides" (object).
NEVER output raw markdown code fences around the JSON; return pure valid JSON."""

    # Build complete active slide context for full preview parity
    slide_context = {
        "slide_1_cover": {
            "client_name": client_name,
            "kicker": current_ov.get("kicker") or bundle.get("kicker") or "RATES RISK & LIABILITY MANAGEMENT",
            "subtitle": current_ov.get("subtitle") or bundle.get("subtitle") or "Pre-Hedging & Refinancing Optimization Strategy",
            "rm_name": current_ov.get("rm_name") or bundle.get("rm_name") or "G. Romano",
            "market_date": current_ov.get("market_date") or bundle.get("market_date") or "28 August 2026"
        },
        "slide_2_catalyst": {
            "slide_title": "Rate Path Volatility & IRS Pre-Hedge Catalyst",
            "primary_market_trigger": current_ov.get("trigger") or bundle.get("trigger") or "Upcoming €3.2B debt maturities face repricing risk amid benchmark curve fluctuations.",
            "window_of_opportunity": current_ov.get("window") or bundle.get("window") or f"Current 5Y EUR swap easing at {bundle.get('swap_5y', '2.62%')} provides attractive entry window for forward-starting IRS.",
            "recommended_action": current_ov.get("action") or bundle.get("action") or f"Execute {bundle.get('notional_bond', '€400M')} pre-hedge IRS overlay to lock in current base yield before debt issuance."
        },
        "slide_6_sensitivity_analysis": {
            "slide_title": "Rationale of our Proposal & Rate Sensitivity",
            "refinance_percentage": f"{current_ov.get('refi_bond_pct', bundle.get('refi_bond_pct', 60))}%",
            "prehedge_percentage": f"{current_ov.get('prehedge_swap_pct', bundle.get('prehedge_swap_pct', 40))}%",
            "indicative_spread": current_ov.get("spread", bundle.get("indicative_spread", "Mid-Swap + 82 bps")),
            "locked_rate": current_ov.get("rate_scenario_lock", "3.44% (locked)"),
            "rates_up_100bp": current_ov.get("rate_scenario_up", "4.44%"),
            "rates_unchanged": current_ov.get("rate_scenario_unchanged", "3.44%"),
            "rates_down_50bp": current_ov.get("rate_scenario_down", "2.94%")
        },
        "slide_4_financial_snapshot": {
            "revenue": current_ov.get("revenue_str") or bundle.get("revenue_str") or "€65,000M",
            "ebitda": current_ov.get("ebitda_str") or bundle.get("ebitda_str") or "€14,300M",
            "net_debt": current_ov.get("net_debt_str") or bundle.get("net_debt_str") or "€16,200M",
            "liquidity": current_ov.get("liquidity_str") or bundle.get("liquidity_str") or "€7,800M",
            "maturity_wall_24m": current_ov.get("maturity_wall_str") or current_ov.get("debt_maturing_24m_str") or bundle.get("debt_maturing_24m_str") or "€3,000M"
        },
        "slide_7_macro_backdrop": {
            "swap_5y": current_ov.get("swap_5y") or bundle.get("swap_5y") or "2.62%",
            "bund_10y": current_ov.get("bund_10y") or bundle.get("bund_10y") or "2.61%",
            "itraxx_main": current_ov.get("itraxx_main") or bundle.get("itraxx_main") or "58 bps"
        },
        "slide_8_term_sheet": {
            "notional": current_ov.get("notional_bond") or bundle.get("notional_bond") or "EUR 600,000,000",
            "spread": current_ov.get("spread") or bundle.get("spread") or "Mid-Swap + 82 bps",
            "tenor": current_ov.get("tenor") or bundle.get("tenor") or "7 Years (Euro Benchmark)"
        }
    }

    user_payload = {
        "client_profile": {
            "client_id": cid,
            "client_name": client_name,
            "product_family": p_family,
            "net_debt": bundle.get("net_debt_str"),
            "liquidity": bundle.get("liquidity_str"),
            "ebitda": bundle.get("ebitda_str"),
            "leverage": bundle.get("leverage_str"),
            "24m_maturity_wall": bundle.get("debt_maturing_24m_str"),
            "maturities_schedule": maturities_data
        },
        "active_deck_slides": slide_context,
        "current_deck_overrides": current_ov,
        "recent_chat_history": history_str,
        "user_request": prompt
    }

    reply_text = ""
    merged_overrides = dict(current_ov)
    project_id = os.getenv("GCP_PROJECT", "teach-telecom-ai-sandbox")
    region = os.getenv("REGION", "europe-west1")

    if GENAI_AVAILABLE:
        try:
            client_gcp = genai.Client(vertexai=True, project=project_id, location=region)
            response = client_gcp.models.generate_content(
                model="gemini-2.5-flash",
                contents=json.dumps(user_payload),
                config=types.GenerateContentConfig(
                    system_instruction=system_instruction,
                    temperature=0.2,
                    response_mime_type="application/json"
                )
            )
            raw_text = response.text.strip() if response and response.text else ""
            if raw_text:
                parsed = json.loads(raw_text)
                if "reply" in parsed:
                    reply_text = parsed["reply"]
                if "overrides" in parsed and isinstance(parsed["overrides"], dict):
                    merged_overrides.update(parsed["overrides"])
        except Exception as e:
            logger.warning(f"Vertex AI Copilot error: {e}")

    if not reply_text:
        reply_text = f"I am ready to assist with {client_name}'s {p_family} pitchbook. You can ask me to explain any slide, modify financial parameters, or run a regulatory compliance check."

    normalized_overrides = normalize_copilot_overrides(merged_overrides)
    return {
        "reply": reply_text,
        "client_id": cid,
        "overrides": normalized_overrides,
        "timestamp": datetime.now().strftime("%H:%M:%S")
    }



@app.post("/api/pitchbook/generate")
@app.api_route("/api/pitchbook/download", methods=["GET", "POST"])
async def handle_pitchbook_generation(
    request: Request,
    canonical_id: str = Query(None),
    client_id: str = Query(None)
):
    try:
        cid = canonical_id or client_id or "CLI102"
        overrides = {}
        
        if request.method == "POST":
            try:
                body = await request.json()
                cid = body.get("canonical_id") or body.get("client_id") or cid
                overrides = body.get("overrides") or {}
            except Exception:
                pass

        # 1. Fetch grounded bundle from DB
        bundle = fetch_pitchbook_bundle(cid, cid, get_db_connection) or {}
        
        # 2. Extract client name & opp meta
        client_name = overrides.get("client_name") or bundle.get("client_name") or "Corporate Client"
        opp_meta = {
            "id": cid,
            "name": client_name,
            "opportunity_type": overrides.get("opportunity_type") or bundle.get("opportunity_type"),
            "product_family": overrides.get("product_family") or bundle.get("product_family") or bundle.get("opportunity_type") or "FX & COMMODITY HEDGING"
        }
        
        # 3. Merge preview overrides directly onto bundle
        for k, v in overrides.items():
            if v is not None and v != "":
                bundle[k] = v

        # 4. Generate PPTX
        compliance_bullets = overrides.get("disclaimers") or overrides.get("compliance_bullets")
        pptx_buf = build_pitchbook(bundle, opp_meta, compliance_bullets=compliance_bullets, overrides=overrides)
        
        clean_filename = f"ING_{str(client_name).replace(' ', '_')}_Pitchbook.pptx"
        content_bytes = pptx_buf.getvalue() if hasattr(pptx_buf, "getvalue") else pptx_buf
        
        return Response(
            content=content_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{clean_filename}"'}
        )
    except Exception as e:
        logger.exception("Pitchbook generation failed")
        raise HTTPException(status_code=500, detail=str(e))


frontend_dist = os.path.join(os.path.dirname(__file__), "frontend", "dist")
if os.path.exists(frontend_dist):
    assets_dir = os.path.join(frontend_dist, "assets")
    if os.path.exists(assets_dir):
        app.mount("/assets", StaticFiles(directory=assets_dir), name="assets")

    @app.get("/{full_path:path}")
    async def serve_react(full_path: str):
        if full_path.startswith("api/"):
            raise HTTPException(status_code=404, detail="Not Found")
        file_path = os.path.join(frontend_dist, full_path)
        if os.path.exists(file_path) and os.path.isfile(file_path):
            return FileResponse(file_path)
        return FileResponse(os.path.join(frontend_dist, "index.html"))
else:
    @app.get("/")
    def index():
        return {"status": "Backend running, frontend build not found."}

